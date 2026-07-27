# -*- coding: utf-8 -*-
"""gate_image_coverage.py — 配图覆盖率门禁检测

职责：检查 PPTX 中配图覆盖率是否达标。
      按规则：有图片的页面数 / 总页面数 >= 阈值（默认 30%）
      覆盖率等级：<30% LOW (红色警告) | 30-50% FAIR | >=50% GOOD

图片检测逻辑：
   - 统计 slide.shapes 中 MSO_SHAPE_TYPE.PICTURE 类型的 shape
   - 图片面积需 >= 幻灯片面积 × 2%（排除微型装饰图标）
   - ImageGen 照片、PIL 生成图、嵌入的 PNG/JPG 均可检测
   - SVG 内联图片：GROUP shape 包含 ≥6 个子形状且面积 ≥5% → 计为等效配图
   - 无图片但有多形状的页面会输出警告提示

使用方式：
  python gate_image_coverage.py input.pptx [--threshold 0.3]
  python gate_image_coverage.py input.pptx --verbose    # 逐页报告
  python gate_image_coverage.py input.pptx --json       # JSON 输出

退出码:
  0 = 通过
  1 = 未通过
  2 = 错误
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR: python-pptx is required.  pip install python-pptx", file=sys.stderr)
    sys.exit(2)


# ── Slide analysis ──

def has_meaningful_image(slide, min_area_ratio=0.02):
    """Check if a slide has at least one meaningful image.
    
    "Meaningful" means the image area >= slide area × min_area_ratio.
    This filters out tiny decorative icons (< 2% of slide area).
    
    Returns (has_image: bool, image_count: int, details: list)
    """
    prs_width = slide.part.package.presentation_part.presentation.slide_width
    prs_height = slide.part.package.presentation_part.presentation.slide_height
    slide_area = prs_width * prs_height
    
    images = []
    
    for shape in slide.shapes:
        try:
            shape_type = shape.shape_type
        except ValueError:
            # Some shapes raise ValueError when accessing shape_type
            continue
        
        # Check for picture shapes
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_area = shape.width * shape.height
            area_ratio = shape_area / slide_area if slide_area > 0 else 0
            
            images.append({
                "type": "PICTURE",
                "name": getattr(shape, "name", "unknown"),
                "width_emu": shape.width,
                "height_emu": shape.height,
                "area_ratio": round(area_ratio * 100, 1),  # as percentage
                "meaningful": area_ratio >= min_area_ratio,
            })
        
        # Check for placeholder shapes that may contain images
        elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            # Placeholders with images have an image blob
            if hasattr(shape, 'image') and shape.image is not None:
                try:
                    _ = shape.image.blob  # Accessing blob confirms it has an image
                    shape_area = shape.width * shape.height
                    area_ratio = shape_area / slide_area if slide_area > 0 else 0
                    images.append({
                        "type": "PLACEHOLDER_WITH_IMAGE",
                        "name": getattr(shape, "name", "unknown"),
                        "width_emu": shape.width,
                        "height_emu": shape.height,
                        "area_ratio": round(area_ratio * 100, 1),
                        "meaningful": area_ratio >= min_area_ratio,
                    })
                except Exception:
                    pass
        
        # Check for group shapes that may contain embedded images
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            # Group shapes can contain pictures — recurse
            try:
                child_shapes = list(shape.shapes)
                found_picture_in_group = False
                for child in child_shapes:
                    try:
                        if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            child_area = child.width * child.height
                            area_ratio = child_area / slide_area if slide_area > 0 else 0
                            images.append({
                                "type": "GROUPED_PICTURE",
                                "name": getattr(child, "name", "unknown"),
                                "width_emu": child.width,
                                "height_emu": child.height,
                                "area_ratio": round(area_ratio * 100, 1),
                                "meaningful": area_ratio >= min_area_ratio,
                            })
                            found_picture_in_group = True
                    except (ValueError, AttributeError):
                        pass
                
                # K2 fix: SVG inline images are groups of native shapes (not PICTURE type).
                # Heuristic: a group with >= 6 child shapes covering >= 5% of the slide
                # is likely an SVG-derived illustration, counted as an equivalent image.
                if not found_picture_in_group and len(child_shapes) >= 6:
                    shape_area = shape.width * shape.height
                    area_ratio = shape_area / slide_area if slide_area > 0 else 0
                    if area_ratio >= 0.05:  # >= 5% of slide area
                        images.append({
                            "type": "SVG_EQUIVALENT",
                            "name": getattr(shape, "name", "unknown"),
                            "width_emu": shape.width,
                            "height_emu": shape.height,
                            "area_ratio": round(area_ratio * 100, 1),
                            "meaningful": True,  # Always meaningful if it passed the threshold
                        })
            except Exception:
                pass
    
    meaningful = [img for img in images if img["meaningful"]]
    return len(meaningful) > 0, len(images), images


def analyze_pptx(pptx_path, min_area_ratio=0.02):
    """Analyze image coverage of an entire PPTX.
    
    Returns {
        total_slides: int,
        slides_with_images: int,
        coverage: float,             # 0.0 - 1.0
        slide_details: [
            {
                slide_num: int,
                has_image: bool,
                image_count: int,
                images: [...],
            },
            ...
        ],
        warnings: list[str],
    }
    """
    prs = Presentation(pptx_path)
    total = len(prs.slides)
    details = []
    warnings = []
    
    for i, slide in enumerate(prs.slides):
        has_img, img_count, images = has_meaningful_image(slide, min_area_ratio)
        details.append({
            "slide_num": i + 1,
            "has_image": has_img,
            "image_count": img_count,
            "images": images,
        })
    
    slides_with_images = sum(1 for d in details if d["has_image"])
    coverage = slides_with_images / total if total > 0 else 0
    
    # Warnings
    svg_warning_slides = []
    for d in details:
        # Heuristic: if a slide has no detected images but has many shapes,
        # it might contain SVG-inline images that we still couldn't detect
        # (those with < 6 child shapes or < 5% area).
        # Use a layout-aware threshold: cover/divider pages have fewer shapes,
        # content pages with dense shapes are more suspicious.
        slide = prs.slides[d["slide_num"] - 1]
        shape_count = len(slide.shapes)
        # Dynamic threshold: 8 for pages with tables, 12 otherwise
        has_table = any(
            hasattr(s, 'has_table') and s.has_table
            for s in slide.shapes
        )
        threshold = 8 if has_table else 12
        if not d["has_image"] and shape_count > threshold:
            svg_warning_slides.append(d["slide_num"])
    
    if svg_warning_slides:
        warnings.append(
            f"Slides {svg_warning_slides} have many shapes but no detected images — "
            f"may contain SVG-inline images (undetectable by this script). "
            f"Consider visual inspection."
        )
    
    return {
        "total_slides": total,
        "slides_with_images": slides_with_images,
        "coverage": round(coverage, 3),
        "slide_details": details,
        "warnings": warnings,
    }


# ── Reporting ──

def print_report(result, threshold, verbose=False):
    """Print human-readable coverage report."""
    total = result["total_slides"]
    with_img = result["slides_with_images"]
    coverage = result["coverage"]
    pct = coverage * 100
    
    is_pass = coverage >= threshold
    pass_fail = "PASS" if is_pass else "FAIL"
    
    # Coverage grade
    if coverage >= 0.5:
        grade = "GOOD"
    elif coverage >= 0.3:
        grade = "FAIR"
    else:
        grade = "LOW"
    
    print(f"Image Coverage: {with_img}/{total} slides ({pct:.1f}%)  [{pass_fail}]  Grade: {grade}")
    print(f"Threshold: {threshold * 100:.0f}%")
    
    if grade == "LOW" and not is_pass:
        print("  NOTE: Coverage is critically low. Add illustrations to at least the chapter/section pages.")
    elif grade == "LOW" and is_pass:
        print("  NOTE: Coverage is low even though it passes the threshold. Consider adding more illustrations.")
    elif grade == "GOOD" and not is_pass:
        print("  NOTE: Coverage is decent but below the strict threshold. Adjust --threshold if needed.")
    
    if verbose:
        print()
        for d in result["slide_details"]:
            num = d["slide_num"]
            status = "has image" if d["has_image"] else "no image"
            img_info = ""
            if d["images"]:
                sizes = [f"{img['area_ratio']}%" for img in d["images"]]
                img_info = f" ({', '.join(sizes)})"
            print(f"  Slide {num:2d}: {status}{img_info}")
    
    if result["warnings"]:
        print()
        for w in result["warnings"]:
            print(f"  WARNING: {w}")
    
    return coverage >= threshold


def main():
    parser = argparse.ArgumentParser(
        description="Check PPTX image coverage against threshold gate."
    )
    parser.add_argument("input", help="Input PPTX file")
    parser.add_argument("--threshold", type=float, default=0.3,
                        help="Minimum coverage ratio (default: 0.3 = 30%%). "
                             "Typical values: 0.3 (conservative), 0.5 (moderate), 0.7 (strict)")
    parser.add_argument("--min-area", type=float, default=0.02,
                        help="Minimum image area ratio to count as meaningful (default: 0.02 = 2%%)")
    parser.add_argument("--verbose", "-v", action="store_true",
                        help="Show per-slide details")
    parser.add_argument("--json", action="store_true",
                        help="Output JSON instead of text")
    args = parser.parse_args()

    if not Path(args.input).exists():
        print(f"ERROR: {args.input} not found", file=sys.stderr)
        sys.exit(2)

    result = analyze_pptx(args.input, args.min_area)

    if args.json:
        # Output JSON for pipeline integration
        output = {
            "pass": result["coverage"] >= args.threshold,
            "coverage": result["coverage"],
            "threshold": args.threshold,
            "total_slides": result["total_slides"],
            "slides_with_images": result["slides_with_images"],
            "warnings": result["warnings"],
        }
        print(json.dumps(output, indent=2, ensure_ascii=False))
        passed = output["pass"]
    else:
        passed = print_report(result, args.threshold, args.verbose)

    sys.exit(0 if passed else 1)


if __name__ == "__main__":
    main()
