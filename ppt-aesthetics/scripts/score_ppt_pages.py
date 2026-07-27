#!/usr/bin/env python3
"""
PPT Aesthetic Structural Scorer — programmatic analysis of PPTX files.

Companion to the AI vision model scoring (which uses rendered PNGs).
This script reads the PPTX structurally and produces a JSON report of
detectable aesthetic issues that vision models may miss or that can be
verified deterministically:
  - colors used (count, variety, harmony)
  - fonts used (count, families, CJK/Latin mismatch)
  - text density per slide
  - alignment & grid consistency
  - margin violations (content touching slide edges)
  - shape count & layout type heuristics
  - image aspect ratio / stretch detection
  - known anti-patterns (3-equal-cards, AI-purple gradient, wall-of-text, etc.)

Usage:
    python score_ppt_pages.py <pptx_path> [--output report.json] [--slides 1,3,5]

The AI vision scoring prompt is included in the report; the agent then
combines this structural analysis with vision scores for a final grade.

Dependencies:
    pip install python-pptx Pillow
"""
from __future__ import annotations

import argparse
import json
import math
import os
import sys
from collections import Counter
from pathlib import Path

# Force UTF-8 stdout on Windows so CJK characters don't mojibake in console.
try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.dml import MSO_FILL
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx",
          file=sys.stderr)
    sys.exit(1)


EMU_PER_INCH = 914400
EMU_PER_PT = 12700

# --------------------------------------------------------------------------- #
# Scenario-based weight presets (改进项 4：场景化权重切换)
# Source: ppt-aesthetics-improvement-plan_2026-06-25.md §六
# --------------------------------------------------------------------------- #
SCENARIO_WEIGHTS = {
    "default": {
        "whitespace": 0.13, "type_scale": 0.14, "color_harmony": 0.14,
        "alignment": 0.14, "spacing": 0.09, "imagery": 0.09,
        "consistency": 0.09, "hierarchy": 0.08, "layout": 0.10,
    },
    "executive": {
        "whitespace": 0.18, "type_scale": 0.09, "color_harmony": 0.09,
        "alignment": 0.13, "spacing": 0.09, "imagery": 0.05,
        "consistency": 0.09, "hierarchy": 0.18, "layout": 0.10,
    },
    "marketing": {
        "whitespace": 0.09, "type_scale": 0.09, "color_harmony": 0.18,
        "alignment": 0.09, "spacing": 0.09, "imagery": 0.13,
        "consistency": 0.09, "hierarchy": 0.14, "layout": 0.10,
    },
    "data": {
        "whitespace": 0.05, "type_scale": 0.09, "color_harmony": 0.09,
        "alignment": 0.13, "spacing": 0.09, "imagery": 0.09,
        "consistency": 0.09, "hierarchy": 0.27, "layout": 0.10,
    },
    "gov": {
        "whitespace": 0.05, "type_scale": 0.18, "color_harmony": 0.09,
        "alignment": 0.18, "spacing": 0.09, "imagery": 0.05,
        "consistency": 0.13, "hierarchy": 0.13, "layout": 0.10,
    },
    "creative": {
        "whitespace": 0.13, "type_scale": 0.13, "color_harmony": 0.18,
        "alignment": 0.05, "spacing": 0.09, "imagery": 0.13,
        "consistency": 0.05, "hierarchy": 0.14, "layout": 0.10,
    },
    "telecom": {
        "whitespace": 0.09, "type_scale": 0.13, "color_harmony": 0.18,
        "alignment": 0.13, "spacing": 0.09, "imagery": 0.05,
        "consistency": 0.13, "hierarchy": 0.10, "layout": 0.10,
    },
}

# Anti-pattern color signatures (hex strings without #, uppercase).
AI_PURPLE_GRADIENT_KEYS = ["6366F1", "818CF8", "A855F7", "C084FC", "7C3AED", "9333EA"]
OVERSATURATED_DEFAULTS = {
    "FF0000": "pure red",
    "00FF00": "pure green",
    "0000FF": "pure blue",
    "FFFF00": "pure yellow",
    "FF00FF": "pure magenta",
    "00FFFF": "pure cyan",
}

# Fonts considered unprofessional in slides (any platform).
BAD_FONTS = {
    "comic sans ms", "comic sans", "papyrus", "lobster", "impact",
    "brush script mt", "chiller", "jokerman", "magneto", "snap itc",
    "playbill", "stencil", " curls mt",
}

# Common default/template fonts (low-effort signal).
GENERIC_FONTS = {
    "calibri", "arial", "times new roman", "carlito", "dejavu sans",
    " Liberation Sans",
}


# --------------------------------------------------------------------------- #
# Color helpers
# --------------------------------------------------------------------------- #
def _rgb_to_hex(rgb) -> str | None:
    try:
        return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}" if isinstance(rgb, (tuple, list)) else str(rgb).upper()
    except Exception:
        return None


def _hex_to_rgb(h: str):
    h = h.lstrip("#").upper()
    if len(h) == 6:
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return None


def _color_distance(h1: str, h2: str) -> float:
    """Euclidean distance in RGB space, 0..441.67."""
    r1, g1, b1 = _hex_to_rgb(h1)
    r2, g2, b2 = _hex_to_rgb(h2)
    return math.sqrt((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2)


def _color_luminance(h: str) -> float:
    """Relative luminance (0..1) per WCAG."""
    r, g, b = _hex_to_rgb(h)
    def chan(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def _contrast_ratio(h1: str, h2: str) -> float:
    """WCAG contrast ratio between two hex colors, 1..21."""
    l1 = _color_luminance(h1)
    l2 = _color_luminance(h2)
    hi, lo = max(l1, l2), min(l1, l2)
    return (hi + 0.05) / (lo + 0.05)


def _detect_slide_bg_color(slide) -> str | None:
    """Best-effort extraction of a slide's solid background fill color.

    Tries (1) python-pptx slide.background.fill, (2) direct XML parse of
    cSld/bg/bgPr/solidFill, (3) the largest shape covering >80% of slide area.
    Returns hex string 'RRGGBB' or None.
    """
    # Method 1: python-pptx background API
    try:
        from pptx.enum.dml import MSO_FILL as _MSO_FILL
        bg_fill = slide.background.fill
        if bg_fill.type == _MSO_FILL.SOLID:
            rgb = bg_fill.fore_color.rgb
            return _rgb_to_hex(rgb)
    except Exception:
        pass

    # Method 2: direct XML parse
    try:
        from pptx.oxml.ns import qn as _qn
        from lxml import etree
        cSld = slide.element.find(".//" + _qn("p:cSld"))
        if cSld is not None:
            bg = cSld.find(_qn("p:bg"))
            if bg is not None:
                solid = bg.find(".//" + _qn("a:solidFill"))
                if solid is not None:
                    srgb = solid.find(_qn("a:srgbClr"))
                    if srgb is not None and srgb.get("val"):
                        return srgb.get("val").upper()
    except Exception:
        pass

    # Method 3: largest shape covering >80% of slide area
    try:
        slide_w = slide.slide_layout.slide_master.slide_width
        slide_h = slide.slide_layout.slide_master.slide_height
        slide_area = slide_w * slide_h
        best_fill = None
        best_area = 0
        for shape in slide.shapes:
            bbox = _shape_bbox_inches(shape)
            if not bbox:
                continue
            l, t, w, h = bbox
            shape_area = w * h
            if shape_area / (slide_w / 914400 * slide_h / 914400) > 0.80:
                fill_hex = _shape_fill_hex(shape)
                if fill_hex and shape_area > best_area:
                    best_area = shape_area
                    best_fill = fill_hex
        return best_fill
    except Exception:
        pass

    return None


def _is_dark_theme(bg_hex: str | None) -> bool:
    """True if the background color is dark (luminance < 0.2)."""
    if not bg_hex:
        return False
    return _color_luminance(bg_hex) < 0.2


def _is_full_bleed_element(bbox, slide_w, slide_h) -> bool:
    """True if a shape is a thin full-width or full-height design element
    (header bars, accent lines, progress bars, side bars) that intentionally
    spans to slide edges and should be exempt from margin checks."""
    if bbox is None:
        return False
    l, t, w, h = bbox
    # Full-width thin element (header bar, accent line, progress bar)
    if w >= slide_w * 0.80 and h <= slide_h * 0.25:
        return True
    # Full-height thin element (side bar)
    if h >= slide_h * 0.80 and w <= slide_w * 0.25:
        return True
    # Full-slide background rectangle
    if w >= slide_w * 0.80 and h >= slide_h * 0.80:
        return True
    # Full-width card/container (spans most of slide width)
    if w >= slide_w * 0.75 and h <= slide_h * 0.50:
        return True
    return False


def _is_page_number_text(text: str) -> bool:
    """True if text looks like a page number, chapter number, or similar
    non-content marker (e.g. '5', '13', '01', '02', '/ 14')."""
    stripped = text.strip()
    if not stripped:
        return False
    # Pure number (possibly with leading zeros)
    if stripped.isdigit():
        return True
    # Patterns like "/ 14", "5 / 14", "5/14"
    if all(c.isdigit() or c.isspace() or c == "/" for c in stripped) and "/" in stripped:
        return True
    # Two-digit chapter numbers like "01", "02"
    if len(stripped) <= 3 and stripped.replace("0", "").isdigit():
        return True
    return False


def _union_area(rects: list[tuple]) -> float:
    """Compute the total area covered by a set of rectangles, accounting
    for overlaps. Uses coordinate compression for efficiency.

    Each rect is (left, top, width, height) in inches.
    """
    if not rects:
        return 0.0

    # Convert (l, t, w, h) -> (x1, y1, x2, y2)
    boxes = [(r[0], r[1], r[0] + r[2], r[1] + r[3]) for r in rects]

    # Collect unique x and y coordinates
    xs = sorted(set(b[0] for b in boxes) | set(b[2] for b in boxes))
    ys = sorted(set(b[1] for b in boxes) | set(b[3] for b in boxes))
    if len(xs) < 2 or len(ys) < 2:
        return 0.0

    # Build a coverage grid
    total = 0.0
    for i in range(len(xs) - 1):
        x1, x2 = xs[i], xs[i + 1]
        dx = x2 - x1
        if dx <= 0:
            continue
        for j in range(len(ys) - 1):
            y1, y2 = ys[j], ys[j + 1]
            dy = y2 - y1
            if dy <= 0:
                continue
            # Check if any rectangle covers this cell
            for bx1, by1, bx2, by2 in boxes:
                if bx1 <= x1 and bx2 >= x2 and by1 <= y1 and by2 >= y2:
                    total += dx * dy
                    break  # cell covered, no need to check more rects
    return total


def _color_family(h: str) -> str:
    """Coarse color family for counting distinct accents."""
    r, g, b = _hex_to_rgb(h)
    # Greyscale?
    mx, mn = max(r, g, b), min(r, g, b)
    if mx - mn < 20:
        return "neutral"
    if r > g + 30 and r > b + 30:
        return "red"
    if g > r + 30 and g > b + 30:
        return "green"
    if b > r + 30 and b > g + 30:
        return "blue"
    if r > b + 30 and g > b + 30:
        return "yellow"
    if r > g + 30 and b > g + 30:
        return "magenta"
    if g > r + 30 and b > r + 30:
        return "cyan"
    return "mixed"


def _is_neutral_color(h: str) -> bool:
    """True if the color is greyscale (neutral)."""
    return _color_family(h) == "neutral"


def _extract_design_tokens(slide_data: dict, slide_w_in: float, slide_h_in: float) -> dict:
    """
    改进项 2：从单页 slide 分析结果中提取标准化的设计 Token。
    
    Token 字段对齐 CSS/Tailwind 习惯，便于：
    - 多版本对比时结构化对比设计 token
    - 与前端项目对接时共同语言
    - 报告读者快速定位主色/字号/圆角等关键设计决策
    
    详见 references/aesthetic-scoring-rubric.md "Token 报告格式" 小节。
    """
    colors_used = slide_data.get("colors_used", {})
    fonts_used = slide_data.get("fonts_used", {})
    layout_signals = slide_data.get("layout_signals", {})
    shapes = slide_data.get("shapes", [])
    
    # ---- Color tokens ----
    # 主色：非中性色中出现频率最高的
    # 辅助色：非中性色中出现频率第二的
    # 中性最深：中性色中最深的（文字色候选）
    # 中性最浅：中性色中最浅的（背景色候选）
    non_neutral = [(h, c) for h, c in colors_used.items() if not _is_neutral_color(h)]
    neutrals = [(h, c) for h, c in colors_used.items() if _is_neutral_color(h)]
    
    non_neutral_sorted = sorted(non_neutral, key=lambda kv: kv[1], reverse=True)
    neutrals_sorted_by_lum = sorted(neutrals, key=lambda kv: _color_luminance(kv[0]))
    
    color_primary = non_neutral_sorted[0][0] if len(non_neutral_sorted) >= 1 and non_neutral_sorted[0][1] >= 2 else None
    color_accent = non_neutral_sorted[1][0] if len(non_neutral_sorted) >= 2 and non_neutral_sorted[1][1] >= 2 else None
    color_neutral_darkest = neutrals_sorted_by_lum[0][0] if neutrals_sorted_by_lum else None
    color_neutral_lightest = neutrals_sorted_by_lum[-1][0] if neutrals_sorted_by_lum else None
    
    # 背景：尝试从 slide background 提取，否则用中性最浅
    color_background = color_neutral_lightest
    try:
        # python-pptx 对 slide background 的支持有限，这里做尽力提取
        from pptx.enum.dml import MSO_FILL as _MSO_FILL
        bg_fill = None
        # 走 XML 直接读 cSld bg 标签
        from pptx.oxml.ns import qn as _qn
        cSld = slide_data.get("_cSld_element")  # 仅在 analyze_slide 时挂载，此处通常为 None
        # 简化：背景 token 留作 None 让下游人工判断
        if not color_background:
            color_background = None
    except Exception:
        pass
    
    # ---- Typography tokens ----
    heading_font = None
    body_font = None
    heading_size_pt = None
    body_size_pt = None
    
    font_families = list(fonts_used.keys())
    # 标题字体：字号 >= 24 的 run 出现频率最高的字体
    # 正文字体：字号 < 24 的 run 出现频率最高的字体
    # 由于 fonts_used 已经混合了所有字号，这里用 layout_signals 的 font_size_range 简化处理
    fs_range = layout_signals.get("font_size_range_pt")
    if fs_range and len(fs_range) == 2:
        heading_size_pt = fs_range[1]  # max
        body_size_pt = fs_range[0]    # min
    
    # 字体名：cjk_fonts 和 latin_fonts 中各取第一个
    cjk_fonts = layout_signals.get("cjk_fonts", [])
    latin_fonts = layout_signals.get("latin_fonts", [])
    # 优先用 deck 级 fonts_used 的 top（频率最高）
    if fonts_used:
        # 过滤掉 (theme) 前缀
        real_fonts = [(f, c) for f, c in fonts_used.items() if not f.startswith("(theme)")]
        real_fonts.sort(key=lambda kv: kv[1], reverse=True)
        if real_fonts:
            # 第一个通常是出现最多的字体（正文或标题共用）
            primary_font = real_fonts[0][0]
            if _is_cjk_font(primary_font):
                heading_font = primary_font
                body_font = primary_font
            elif _is_latin_only_font(primary_font):
                heading_font = primary_font
                body_font = primary_font
            # 第二个字体（如果有）
            if len(real_fonts) >= 2:
                second_font = real_fonts[1][0]
                if _is_cjk_font(second_font):
                    if not _is_cjk_font(heading_font or ""):
                        # latin + cjk 配对：latin 作 latin 部分，cjk 作 ea 部分
                        # 此处仍以出现频率为标题字体
                        pass
    
    # Type scale ratio = heading_size / body_size
    type_scale_ratio = None
    if heading_size_pt and body_size_pt and body_size_pt > 0:
        type_scale_ratio = round(heading_size_pt / body_size_pt, 3)
    
    # ---- Radius tokens ----
    # 从 shapes 中提取圆角矩形的 radius（in inches）
    # python-pptx 的 ROUNDED_RECTANGLE 通过 adjustments[0] 取 radius，但单位是相对值
    # 这里简化：统计 rounded_rectangle 形状的数量，radius 留待人工或后续完善
    card_count = sum(1 for s in shapes if s.get("auto_shape_type") == "ROUNDED_RECTANGLE")
    radius_card_in = None  # TODO: 后续从 shape XML 提取 adj 值
    radius_button_in = None
    
    # ---- Elevation tokens ----
    # python-pptx 不直接暴露 shadow depth，需读 XML，此处简化
    elevation_card_shadow_pt = None
    elevation_modal_shadow_pt = None
    
    # ---- Spacing tokens ----
    # 页边距：取所有 shape left/top/right/bottom 距离 slide 边缘的最小值
    margin_in = None
    grid_step_in = None
    left_edges = layout_signals.get("left_edges_in", [])
    top_edges = layout_signals.get("top_edges_in", [])
    if left_edges and top_edges:
        min_left = min(left_edges)
        # 右边距需要 shape 的 right edge，但 layout_signals 没存，简化用 left 近似
        margin_in = round(min_left, 3) if min_left < 2.0 else None
    # 网格步长：从 anti_patterns 的 irregular_spacing 信息推断，否则用默认 0.08
    for ap in slide_data.get("anti_patterns", []):
        if ap.startswith("irregular_spacing"):
            grid_step_in = 0.08  # 标称步长
            break
    
    return {
        "color": {
            "primary": color_primary,
            "accent": color_accent,
            "neutral_darkest": color_neutral_darkest,
            "neutral_lightest": color_neutral_lightest,
            "background": color_background,
        },
        "typography": {
            "heading_font": heading_font,
            "body_font": body_font,
            "heading_size_pt": heading_size_pt,
            "body_size_pt": body_size_pt,
            "type_scale_ratio": type_scale_ratio,
        },
        "radius": {
            "card_in": radius_card_in,
            "button_in": radius_button_in,
            "card_count": card_count,
        },
        "elevation": {
            "card_shadow_pt": elevation_card_shadow_pt,
            "modal_shadow_pt": elevation_modal_shadow_pt,
        },
        "spacing": {
            "margin_in": margin_in,
            "grid_step_in": grid_step_in,
        },
    }


def _is_ai_purple(h: str) -> bool:
    """Detect the AI-generated indigo/violet family (#6366F1, #818CF8, etc.).

    Distinguish from corporate navy/blue (#1B3A5C, #4A90D9): true AI purple
    has BOTH high blue AND notable red (violet tint), with green clearly
    lower than both. Navy has very low red; mid-blue has balanced r/g.
    """
    target = _hex_to_rgb(h)
    if not target:
        return False
    r, g, b = target
    # Need saturated blue dominant.
    if b < 200:
        return False
    # Indigo/violet: red is mid (60-180), green clearly lower than both.
    if not (60 <= r <= 180):
        return False
    if g >= 180:
        return False
    # Both red and blue well above green → violet/indigo tint.
    return (r - g > 30) and (b - g > 40)


# --------------------------------------------------------------------------- #
# Font helpers
# --------------------------------------------------------------------------- #
def _is_cjk_font(name: str) -> bool:
    """Heuristic: does this font name look CJK-targeted?"""
    n = (name or "").lower()
    # Strip "(theme)" prefix that we add when attributing to theme defaults.
    if n.startswith("(theme)"):
        n = n[len("(theme)"):].strip().lower()
    cjk_markers = [
        # Chinese names
        "微软雅黑", "雅黑", "黑体", "宋体", "楷体", "仿宋", "幼圆",
        "华文", "方正", "汉仪", "文泉", "思源",
        # English / pinyin aliases
        "microsoft yahei", "yahei", "pingfang", "ping fang", "pingfang sc",
        "songti", "song ti", "heiti", "hei ti", "kaiti", "kai ti",
        "fangsong", "fangsong", "simhei", "simsun", "nsimsun",
        "source han", "source han sans", "source han serif",
        "noto sans cjk", "noto serif cjk",
        "fzxiaobiaosong", "fzshusong", "fzhei",
        # Korean / Japanese (treat as CJK for font-coverage purposes)
        "malgun", "gulim", "dotum", "batang", "meiryo", "yu gothic",
        "ms gothic", "hiragino",
    ]
    return any(m in n for m in cjk_markers)


def _is_latin_only_font(name: str) -> bool:
    """Heuristic: an obviously Latin font (no CJK glyphs)."""
    n = (name or "").lower()
    if _is_cjk_font(n):
        return False
    latin_markers = [
        "arial", "helvetica", "inter", "roboto", "calibri", "segoe",
        "georgia", "times", "garamond", "din", "impact", "trebuchet",
        "verdana", "tahoma", "outfit", "geist", "satoshi", "montserrat",
        "futura", "gill", "baskerville", "franklin",
    ]
    return any(m in n for m in latin_markers)


# --------------------------------------------------------------------------- #
# Shape helpers
# --------------------------------------------------------------------------- #
def _shape_bbox_inches(shape):
    try:
        return (
            shape.left / EMU_PER_INCH,
            shape.top / EMU_PER_INCH,
            shape.width / EMU_PER_INCH,
            shape.height / EMU_PER_INCH,
        )
    except Exception:
        return None


def _is_edge_touching(bbox, slide_w, slide_h, margin_pct=0.08):
    """True if a shape's bbox is within `margin_pct` (default 8%) of any slide edge.

    Uses percentage-based margin per element-ratio-scoring.md spec:
    "四边均 8%–12% 安全边距".  Default 8% is the minimum safe margin;
    callers may pass 0.12 for a stricter check.
    """
    if bbox is None:
        return False
    margin_h = slide_w * margin_pct   # horizontal safety margin
    margin_v = slide_h * margin_pct   # vertical safety margin
    l, t, w, h = bbox
    r, b = l + w, t + h
    return l < margin_h or t < margin_v or r > slide_w - margin_h or b > slide_h - margin_v


def _shape_fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type == MSO_FILL.SOLID:
            rgb = fill.fore_color.rgb
            return _rgb_to_hex(rgb)
    except Exception:
        pass
    return None


def _shape_line_hex(shape) -> str | None:
    try:
        line = shape.line
        if line.color and line.color.type is not None:
            return _rgb_to_hex(line.color.rgb)
    except Exception:
        pass
    return None


def _shape_auto_shape_type(shape) -> str | None:
    """Return the auto-shape type name (e.g. 'ROUNDED_RECTANGLE') if applicable.

    The string form from python-pptx is e.g. 'ROUNDED_RECTANGLE (5)'; we strip
    the trailing ' (N)' so callers can compare against the bare name.
    """
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            raw = str(shape.auto_shape_type)
            # raw looks like "ROUNDED_RECTANGLE (5)" — strip the enum value.
            if " (" in raw:
                raw = raw.split(" (", 1)[0]
            return raw
    except Exception:
        pass
    return None


def _resolve_theme_fonts(prs) -> tuple[str | None, str | None]:
    """Return (major_font, minor_font) from the deck's theme, or (None, None)."""
    try:
        theme = prs.slide_masters[0].element.getparent().getparent()
        # Walk to the theme part — easier: read the theme XML directly.
        from pptx.oxml.ns import qn
        master = prs.slide_masters[0]
        theme_part = master.part.package.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        ) if hasattr(master.part, "package") else None
        if theme_part is None:
            return None, None
        import xml.etree.ElementTree as ET
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        root = ET.fromstring(theme_part.blob)
        major = root.find(".//a:fontScheme/a:majorFont/a:latin", ns)
        minor = root.find(".//a:fontScheme/a:minorFont/a:latin", ns)
        return (
            major.get("typeface") if major is not None else None,
            minor.get("typeface") if minor is not None else None,
        )
    except Exception:
        return None, None


def _resolve_run_font(run, theme_minor_font, theme_major_font) -> str | None:
    """Resolve the effective font name for a run, falling back to theme."""
    fn = run.font.name
    if fn:
        return fn
    # Try font scheme via the run's paragraph properties (rPr lvl)
    try:
        rPr = run._r.get_or_add_rPr()
        if rPr is not None and rPr.get("altLang") is not None:
            pass
    except Exception:
        pass
    # Heading-level paragraphs use major font; body uses minor.
    # We can't easily tell here, so report both as "theme:..." placeholder
    # that downstream logic can treat as "default font, possibly fine".
    return None


# --------------------------------------------------------------------------- #
# Slide analysis
# --------------------------------------------------------------------------- #
def analyze_slide(slide, slide_index, slide_w_in, slide_h_in,
                  theme_major_font=None, theme_minor_font=None):
    """Extract structural and visual metadata from one slide."""
    result = {
        "slide_number": slide_index + 1,
        "shapes": [],
        "text_blocks": [],
        "colors_used": Counter(),
        "text_colors": Counter(),  # only text run colors (for low_contrast check)
        "fonts_used": Counter(),
        "images": [],
        "anti_patterns": [],
        "layout_signals": {},
    }

    shape_bboxes = []
    left_edges, right_edges, top_edges, bottom_edges = [], [], [], []
    text_runs = []  # (text, font_name, font_size_pt)
    font_sizes = []  # all non-None run sizes in pt (for type-scale check)
    rounded_rect_bboxes = []  # for 3-equal-cards detection
    fill_color_usage = Counter()  # hex -> count, for overuse_accent detection

    # Extract layout name for default_template detection
    slide_layout_name = ""
    try:
        slide_layout_name = slide.slide_layout.name if slide.slide_layout else ""
    except Exception:
        pass

    for shape in slide.shapes:
        bbox = _shape_bbox_inches(shape)
        if bbox:
            shape_bboxes.append(bbox)
            l, t, w, h = bbox
            left_edges.append(l)
            right_edges.append(l + w)
            top_edges.append(t)
            bottom_edges.append(t + h)

        auto_type = _shape_auto_shape_type(shape)
        shape_info = {
            "name": getattr(shape, "name", str(shape.shape_id)),
            "type": str(shape.shape_type),
            "auto_shape_type": auto_type,
            "bbox_in": [round(b, 3) for b in bbox] if bbox else None,
            "has_text_frame": shape.has_text_frame,
        }
        if auto_type == "ROUNDED_RECTANGLE" and bbox:
            rounded_rect_bboxes.append(bbox)

        # Fill / line colors
        fill_hex = _shape_fill_hex(shape)
        if fill_hex:
            result["colors_used"][fill_hex] += 1
            fill_color_usage[fill_hex] += 1
        line_hex = _shape_line_hex(shape)
        if line_hex and line_hex != fill_hex:
            result["colors_used"][line_hex] += 1

        # Images — track for stretch detection
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                iw, ih = img.size  # px
                aspect_img = iw / ih if ih else 0
                aspect_box = (bbox[2] / bbox[3]) if bbox and bbox[3] else 0
                stretch = None
                if aspect_img and aspect_box:
                    ratio = aspect_box / aspect_img
                    if ratio < 0.92 or ratio > 1.08:
                        stretch = "compressed" if ratio < 0.92 else "stretched"
                result["images"].append({
                    "bbox_in": [round(b, 3) for b in bbox] if bbox else None,
                    "px_size": [iw, ih],
                    "box_aspect": round(aspect_box, 3),
                    "img_aspect": round(aspect_img, 3),
                    "stretch": stretch,
                })
                if stretch:
                    result["anti_patterns"].append(
                        f"stretched_image: image at {bbox} is {stretch} "
                        f"(box {aspect_box:.2f} vs img {aspect_img:.2f})"
                    )
            except Exception:
                pass

        # Text
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                result["text_blocks"].append({
                    "text": text[:120],
                    "length": len(text),
                    "shape_name": shape_info["name"],
                })
                for run in paragraph.runs:
                    fn = run.font.name
                    fs = run.font.size
                    fs_pt = fs.pt if fs else None
                    if fs_pt:
                        font_sizes.append(fs_pt)
                    if fn:
                        result["fonts_used"][fn] += 1
                    # If font name not set on the run, the deck falls back
                    # to the theme's minor (body) / major (heading) font.
                    # Track that as a synthetic entry so downstream logic
                    # knows which theme font is in play.
                    elif theme_minor_font or theme_major_font:
                        # Heuristic: large fonts (>24pt) → heading (major)
                        # else body (minor). Mark with theme: prefix.
                        key = (theme_major_font
                               if (fs_pt and fs_pt >= 24)
                               else theme_minor_font)
                        if key:
                            result["fonts_used"][f"(theme){key}"] += 1
                    # Run color
                    try:
                        if run.font.color and run.font.color.type is not None:
                            ch = _rgb_to_hex(run.font.color.rgb)
                            if ch:
                                result["colors_used"][ch] += 1
                                result["text_colors"][ch] += 1
                    except Exception:
                        pass
                    text_runs.append((run.text, fn, fs_pt))

        result["shapes"].append(shape_info)

    # ---- Detect slide background color (used by multiple anti-pattern checks) ----
    slide_bg_hex = _detect_slide_bg_color(slide)
    result["slide_bg_color"] = slide_bg_hex
    is_dark_bg = _is_dark_theme(slide_bg_hex)

    # ---- Anti-pattern detection ----
    # 1. Wall of text: > 60 words body or > 400 chars
    total_text = sum(t["length"] for t in result["text_blocks"])
    total_words = sum(len(t["text"].split()) for t in result["text_blocks"])
    if total_words > 60:
        result["anti_patterns"].append(
            f"wall_of_text: {total_words} words across {len(result['text_blocks'])} blocks"
        )
    elif total_text > 400:
        result["anti_patterns"].append(
            f"high_text_density: {total_text} chars (consider splitting)"
        )

    # 2. Too many fonts (> 3 families)
    font_families = list(result["fonts_used"].keys())
    if len(font_families) > 3:
        result["anti_patterns"].append(
            f"too_many_fonts: {len(font_families)} families -> {font_families}"
        )

    # 3. Bad fonts
    for fn in font_families:
        low = fn.lower()
        if any(b in low for b in BAD_FONTS):
            result["anti_patterns"].append(f"unprofessional_font: {fn}")

    # 4. CJK + Latin-only font mismatch (e.g. body text in Helvetica but
    #    contains 中文 — Helvetica has no CJK glyphs, will fallback messily)
    cjk_fonts = [f for f in font_families if _is_cjk_font(f)]
    latin_fonts = [f for f in font_families if _is_latin_only_font(f)]
    has_cjk_text = any(_contains_cjk(t["text"]) for t in result["text_blocks"])
    if has_cjk_text and not cjk_fonts:
        result["anti_patterns"].append(
            "cjk_text_no_cjk_font: 中文 text present but no CJK font assigned"
        )

    # 5. Color discipline
    #    Dark themes need more shades (bg, card, header, text variants, accent)
    #    so raise the threshold from 5 to 8 for dark backgrounds.
    non_neutral_colors = [
        h for h in result["colors_used"].keys()
        if _color_family(h) != "neutral"
    ]
    distinct_families = set(_color_family(h) for h in non_neutral_colors)
    color_threshold = 8 if is_dark_bg else 5
    if len(non_neutral_colors) > color_threshold:
        result["anti_patterns"].append(
            f"too_many_colors: {len(result['colors_used'])} distinct colors "
            f"({len(non_neutral_colors)} non-neutral, families: {distinct_families})"
        )

    # 6. AI purple gradient signature
    ai_purple_hits = [h for h in result["colors_used"].keys() if _is_ai_purple(h)]
    if ai_purple_hits:
        result["anti_patterns"].append(
            f"ai_purple_palette: detected indigo/violet colors {ai_purple_hits[:3]}"
        )

    # 7. Oversaturated pure RGB defaults
    oversaturated = [
        h for h in result["colors_used"].keys() if h in OVERSATURATED_DEFAULTS
    ]
    if oversaturated:
        result["anti_patterns"].append(
            f"oversaturated_pure_rgb: {oversaturated} (use tinted shades instead)"
        )

    # 8. Margin violations (content within safety margin of slide edges)
    #    Exempt full-bleed design elements (header bars, accent lines, etc.)
    #    Use a stricter 3% threshold (not 8%) — 8% flags normal card layouts
    #    where cards span ~80% width with 0.5in margins. 3% only flags
    #    truly edge-touching content (< 0.4in on a 13.33in slide).
    margin_violations = [
        s["name"] for s, b in zip(result["shapes"], shape_bboxes)
        if b and _is_edge_touching(b, slide_w_in, slide_h_in, margin_pct=0.03)
        and not _is_full_bleed_element(b, slide_w_in, slide_h_in)
    ]
    if margin_violations:
        result["anti_patterns"].append(
            f"no_margins: {len(margin_violations)} shape(s) within 3% safety margin of slide edge"
        )

    # 9. 3-equal-cards heuristic: 3 rounded rectangles same width/height,
    #    side by side, equal spacing
    if len(rounded_rect_bboxes) >= 3:
        # Sort by left edge, then look at consecutive trios.
        rounded_rect_bboxes.sort(key=lambda b: b[0])
        for i in range(len(rounded_rect_bboxes) - 2):
            trio = rounded_rect_bboxes[i:i + 3]
            widths = [b[2] for b in trio]
            heights = [b[3] for b in trio]
            tops = [b[1] for b in trio]
            if (max(widths) - min(widths) < 0.15 and
                max(heights) - min(heights) < 0.15 and
                max(tops) - min(tops) < 0.15):
                result["anti_patterns"].append(
                    "three_equal_cards: 3 same-sized rounded rectangles in a row "
                    "(may lack visual hierarchy; consider varying size/weight)"
                )
                break

    # ---- Web-design-engineer principle checks ----
    # (Type Scale, 60-30-10 overuse, pure black/white, long bullets,
    #  too many bullets, irregular spacing, low-res images)
    # See references/web-design-principles.md for the full rationale.

    # 10. Type Scale — ratio between max and min font size on the slide.
    if len(font_sizes) >= 2:
        fs_max = max(font_sizes)
        fs_min = min(font_sizes)
        if fs_min > 0:
            ratio = fs_max / fs_min
            # Use 1.2 (not 1.25) as the floor to tolerate float rounding
            # when fix_type_scale snaps to Major Third rungs.
            if ratio < 1.2:
                result["anti_patterns"].append(
                    f"weak_type_scale: max/min font ratio = {ratio:.2f} "
                    f"({fs_max}pt/{fs_min}pt); need >= 1.25 (Major Third) "
                    f"for visible hierarchy"
                )
            elif ratio > 3.0:
                result["anti_patterns"].append(
                    f"extreme_type_scale: max/min font ratio = {ratio:.2f} "
                    f"({fs_max}pt/{fs_min}pt); cover-sized type may have "
                    f"leaked into a content slide"
                )

    # 11. Accent color overuse — 60-30-10 rule: accent should be ~10%.
    # Identify the most-used non-neutral fill color; if it covers > 40% of
    # filled shapes, it is no longer an accent.
    non_neutral_fills = {
        h: c for h, c in fill_color_usage.items()
        if _color_family(h) != "neutral"
    }
    if non_neutral_fills:
        total_fills = sum(fill_color_usage.values())
        top_accent, top_count = max(non_neutral_fills.items(), key=lambda kv: kv[1])
        if total_fills and top_count / total_fills > 0.4 and top_count > 5:
            result["anti_patterns"].append(
                f"overuse_accent: color #{top_accent} covers "
                f"{top_count}/{total_fills} filled shapes "
                f"({top_count/total_fills*100:.0f}%); 60-30-10 rule says "
                f"accent should be ~10%"
            )

    # 12. Pure black text / pure white card — use 800-950 / 50-100 tints
    # instead of #000 / #FFF for a more refined feel.
    if "000000" in result["colors_used"]:
        result["anti_patterns"].append(
            "pure_black_text: #000000 detected; use a near-black tint "
            "(e.g. #1A1A1A) for warmer, more refined text"
        )
    if "FFFFFF" in fill_color_usage and fill_color_usage["FFFFFF"] > 2:
        result["anti_patterns"].append(
            "pure_white_card: #FFFFFF used as card fill; use a slightly "
            "off-white surface tint (e.g. #F5F5F5) to avoid harshness"
        )

    # 13. Long bullet — single text block > 35 CJK chars or > 15 EN words.
    for tb in result["text_blocks"]:
        text = tb["text"]
        cjk_count = sum(1 for ch in text if "\u4e00" <= ch <= "\u9fff")
        en_words = len(text.split())
        if cjk_count > 35:
            result["anti_patterns"].append(
                f"long_bullet: text block has {cjk_count} CJK chars "
                f"(>35); split or shorten — ideal 15-25"
            )
        elif en_words > 15 and cjk_count == 0:
            result["anti_patterns"].append(
                f"long_bullet: text block has {en_words} EN words "
                f"(>15); split or shorten — ideal 8-12"
            )

    # 14. Too many bullets — > 6 content text blocks on one slide.
    #     Exclude page numbers and very short labels (non-content markers)
    content_blocks = [
        tb for tb in result["text_blocks"]
        if not _is_page_number_text(tb["text"]) and len(tb["text"]) > 4
    ]
    if len(content_blocks) > 6:
        result["anti_patterns"].append(
            f"too_many_bullets: {len(content_blocks)} content text blocks "
            f"on one slide (max 6, ideal 3-4)"
        )

    # 15. Irregular spacing — vertical gaps between stacked shapes should
    # cluster around 0.08in multiples (4/8px grid).
    if len(shape_bboxes) >= 3:
        sorted_by_top = sorted(shape_bboxes, key=lambda b: b[1])
        v_gaps = []
        for i in range(1, len(sorted_by_top)):
            prev_bottom = sorted_by_top[i-1][1] + sorted_by_top[i-1][3]
            cur_top = sorted_by_top[i][1]
            gap = round(cur_top - prev_bottom, 3)
            if 0 < gap < 2.0:  # ignore huge gaps (different columns)
                v_gaps.append(gap)
        if v_gaps:
            on_grid = sum(
                1 for g in v_gaps
                if abs(g - round(g / 0.08) * 0.08) < 0.02
            )
            if on_grid / len(v_gaps) < 0.5:
                result["anti_patterns"].append(
                    f"irregular_spacing: only {on_grid}/{len(v_gaps)} vertical "
                    f"gaps align to 0.08in grid; use 4/8px rhythm "
                    f"(0.08/0.16/0.24/0.32in)"
                )

    # 16. Low-res image — original pixel dimension below threshold.
    for img_info in result["images"]:
        px = img_info.get("px_size", [0, 0])
        if px and (px[0] < 800 or px[1] < 600):
            result["anti_patterns"].append(
                f"low_res_image: image {px[0]}x{px[1]}px is below "
                f"800x600; will look blurry when projected"
            )

    # 17. Top-heavy layout — content centroid above 40% of slide height.
    #     Detects when all content is crammed into the upper portion,
    #     leaving a large empty gap at the bottom. This is a common
    #     artifact of AI-generated slides that over-weight "higher =
    #     heavier" visual balance rules.
    if shape_bboxes:
        total_area = 0.0
        weighted_top = 0.0
        for b in shape_bboxes:
            area = b[2] * b[3]
            centroid_y = b[1] + b[3] / 2.0
            total_area += area
            weighted_top += area * centroid_y
        if total_area > 0:
            content_centroid_y = weighted_top / total_area
            centroid_ratio = content_centroid_y / slide_h_in
            if centroid_ratio < 0.40:
                result["anti_patterns"].append(
                    f"top_heavy: content centroid at {centroid_ratio:.0%} "
                    f"of slide height (y={content_centroid_y:.2f}in / "
                    f"{slide_h_in:.2f}in); should be 40%-60% for balanced "
                    f"layout — move content down or increase card height"
                )

    # ---- Element ratio scoring (改进项: 页面填充率 + 象限空白检测) ----
    # page_fill_rate: 元素面积占页面面积的比例，最优 60-70%
    # Use union area to avoid double-counting overlapping shapes
    slide_area = slide_w_in * slide_h_in
    element_area_sum = _union_area(shape_bboxes)
    page_fill_rate = element_area_sum / slide_area if slide_area > 0 else 0.0

    # quadrant_empty: 检测四个象限中是否有没有任何元素的象限
    mid_x = slide_w_in / 2.0
    mid_y = slide_h_in / 2.0
    quadrants = [False, False, False, False]  # TL, TR, BL, BR
    for b in shape_bboxes:
        l, t, w, h = b
        center_x = l + w / 2.0
        center_y = t + h / 2.0
        # 元素重心落入的象限
        if center_x < mid_x and center_y < mid_y:
            quadrants[0] = True  # TL
        elif center_x >= mid_x and center_y < mid_y:
            quadrants[1] = True  # TR
        elif center_x < mid_x and center_y >= mid_y:
            quadrants[2] = True  # BL
        else:
            quadrants[3] = True  # BR
        # 也标记元素实际覆盖的象限（大面积元素可能跨象限）
        right = l + w
        bottom = t + h
        if l < mid_x:
            if t < mid_y:
                quadrants[0] = True
            if bottom > mid_y:
                quadrants[2] = True
        if right > mid_x:
            if t < mid_y:
                quadrants[1] = True
            if bottom > mid_y:
                quadrants[3] = True
    empty_quadrants = [i for i, q in enumerate(quadrants) if not q]
    quadrant_names = {0: "左上", 1: "右上", 2: "左下", 3: "右下"}

    # 18. Low fill rate — 元素面积占页面 < 60%（空旷）
    #     Threshold unified to 0.60 (中危 -1.5) to match SKILL.md and
    #     hollow_container HOLLOW_FILL_THRESHOLD = 0.60.
    #     Note: PICTURE shapes used as card backgrounds inflate fill rate;
    #     the hollow_container check (#54) catches container-level emptiness.
    if page_fill_rate < 0.60:
        result["anti_patterns"].append(
            f"low_fill_rate: element area covers {page_fill_rate:.0%} of slide "
            f"({element_area_sum:.1f}/{slide_area:.1f} sq.in); optimal is 60-70%"
        )

    # 19. Overfilled — 元素面积占页面 > 85%（拥挤）
    if page_fill_rate > 0.85:
        result["anti_patterns"].append(
            f"high_fill_rate: element area covers {page_fill_rate:.0%} of slide "
            f"({element_area_sum:.1f}/{slide_area:.1f} sq.in); optimal is 60-70%, "
            f"consider reducing content"
        )

    # 20. Empty quadrant — 某象限无任何元素（布局失衡）
    if empty_quadrants:
        empty_names = [quadrant_names[i] for i in empty_quadrants]
        result["anti_patterns"].append(
            f"empty_quadrant: {len(empty_quadrants)} empty quadrant(s) "
            f"({', '.join(empty_names)}); layout appears unbalanced"
        )

    # ---- New anti-pattern detections (previously defined but unimplemented) ----

    # 21. text_heavy_deck — 文字框总面积 > 页面 60%
    #     (distinct from wall_of_text which counts words; this measures area)
    #     Uses union area to avoid double-counting overlapping text frames
    text_bboxes = [
        b for s, b in zip(result["shapes"], shape_bboxes)
        if b and s.get("has_text_frame")
    ]
    text_area_sum = _union_area(text_bboxes)
    text_area_ratio = text_area_sum / slide_area if slide_area > 0 else 0.0
    if text_area_ratio > 0.60:
        result["anti_patterns"].append(
            f"text_heavy_deck: text frames cover {text_area_ratio:.0%} of slide "
            f"({text_area_sum:.1f}/{slide_area:.1f} sq.in); "
            f"recommend < 60% — add imagery or whitespace"
        )

    # 54. hollow_container — 容器/卡片内部填充率 < 60%（空方块）
    #     Root cause: add_bg() creates fixed-size card background while
    #     add_rich_textbox() with auto_size=None leaves declared height
    #     but text only occupies the top portion → large empty area.
    #     Fix: use add_table() instead (auto-sizes rows to content).
    #     Detection: identify container candidates (large shapes acting as
    #     card backgrounds), compute inner fill rate from overlapping
    #     text boxes, flag if < 60%.
    HOLLOW_FILL_THRESHOLD = 0.60  # container must be ≥60% filled
    HOLLOW_MIN_CONTAINER_AREA = 5.0  # sq.in — ignore tiny shapes (<5% of 100 sq.in slide)
    hollow_found = []
    # Build lists: container candidates vs content shapes (text + tables)
    container_bboxes = []  # (bbox, shape_info)
    textbox_bboxes = []    # (bbox, shape_info) — stand-alone text frames
    table_bboxes = []      # (bbox, shape_info) — tables (text is in cells, not text_frame)
    for s, b in zip(result["shapes"], shape_bboxes):
        if not b:
            continue
        area = b[2] * b[3]
        is_text = s.get("has_text_frame", False)
        auto_type = s.get("auto_shape_type")
        s_type = s.get("type", "")
        is_table = "TABLE" in s_type or "TABLE" in s.get("name", "").upper()
        # Table shapes: track their bbox for container fill calculation
        if is_table:
            table_bboxes.append((b, s))
            continue
        # Container candidates: large non-text shapes (bg pictures,
        # rounded rectangles, plain rectangles acting as card backgrounds)
        # Exclude TABLE shapes — tables ARE the solution, not the problem.
        if not is_text and area >= HOLLOW_MIN_CONTAINER_AREA:
            if "PICTURE" in s_type or auto_type in ("RECTANGLE", "ROUNDED_RECTANGLE") or area >= 8.0:
                container_bboxes.append((b, s))
        # Text boxes: shapes with text_frame that could be inside containers
        if is_text and area > 0.1:  # ignore tiny labels
            textbox_bboxes.append((b, s))
    # For each container, compute inner fill rate
    for cbox, cinfo in container_bboxes:
        cl, ct, cw, ch = cbox
        c_area = cw * ch
        c_right = cl + cw
        c_bottom = ct + ch
        # Skip "horizontal strip cards": wide (>70% slide width) and short
        # (<2.0 in) — these are valid card-bar layouts where low fill rate
        # is expected (e.g. term definition rows, feature strips).
        if cw > slide_w_in * 0.70 and ch < 2.0:
            continue
        # Sum areas of text boxes that overlap with this container
        contained_text_area = 0.0
        for tbox, tinfo in textbox_bboxes:
            tl, tt, tw, th = tbox
            # Compute overlap area (more robust than center-point check)
            # A text box straddling the container edge is partially counted.
            ox_left = max(cl, tl)
            ox_top = max(ct, tt)
            ox_right = min(c_right, tl + tw)
            ox_bottom = min(c_bottom, tt + th)
            if ox_right > ox_left and ox_bottom > ox_top:
                overlap_area = (ox_right - ox_left) * (ox_bottom - ox_top)
                text_area = tw * th
                # Count the overlap fraction of the text box
                fraction = overlap_area / text_area if text_area > 0 else 0
                contained_text_area += tw * th * fraction
        # Also count table areas overlapping this container.
        # Tables are the "good" content that fills containers, so their
        # full bbox counts as filled area (their cells auto-size to content).
        for tbox, tinfo in table_bboxes:
            tl, tt, tw, th = tbox
            # Compute overlap area (same robust method as text boxes)
            ox_left = max(cl, tl)
            ox_top = max(ct, tt)
            ox_right = min(c_right, tl + tw)
            ox_bottom = min(c_bottom, tt + th)
            if ox_right > ox_left and ox_bottom > ox_top:
                overlap_area = (ox_right - ox_left) * (ox_bottom - ox_top)
                table_area = tw * th
                fraction = overlap_area / table_area if table_area > 0 else 0
                contained_text_area += tw * th * fraction
        fill_rate = contained_text_area / c_area if c_area > 0 else 0.0
        if fill_rate < HOLLOW_FILL_THRESHOLD:
            hollow_found.append({
                "container": cinfo.get("name", "?"),
                "container_area": round(c_area, 1),
                "text_area": round(contained_text_area, 1),
                "fill_rate": round(fill_rate, 2),
            })
    if hollow_found:
        worst = min(hollow_found, key=lambda h: h["fill_rate"])
        result["anti_patterns"].append(
            f"hollow_container: {len(hollow_found)} container(s) with "
            f"inner fill rate < {HOLLOW_FILL_THRESHOLD:.0%}; "
            f"worst: '{worst['container']}' at {worst['fill_rate']:.0%} "
            f"(text {worst['text_area']:.1f}/{worst['container_area']:.1f} sq.in); "
            f"use add_table() instead of add_bg()+add_rich_textbox()"
        )

    # 55. inconsistent_spacing — 同类间距差异 > 50%（跨卡片间距不统一）
    #     For sibling shapes (similar height, vertically stacked), check that
    #     the vertical gaps between them are consistent.
    if len(shape_bboxes) >= 3:
        # Collect vertical gaps between same-column text shapes
        text_shapes_sorted = []
        for b in shape_bboxes:
            if b and b[3] > 0.2:  # height > 0.2in (ignore tiny)
                text_shapes_sorted.append(b)
        text_shapes_sorted.sort(key=lambda b: b[1])  # sort by top
        if len(text_shapes_sorted) >= 3:
            gaps = []
            for i in range(1, len(text_shapes_sorted)):
                prev_bottom = text_shapes_sorted[i-1][1] + text_shapes_sorted[i-1][3]
                curr_top = text_shapes_sorted[i][1]
                gap = curr_top - prev_bottom
                if gap > 0.05:  # ignore overlapping or adjacent shapes
                    gaps.append(gap)
            if len(gaps) >= 2:
                min_gap = min(gaps)
                max_gap = max(gaps)
                if min_gap > 0 and max_gap / min_gap > 2.0:
                    result["anti_patterns"].append(
                        f"inconsistent_spacing: vertical gap ratio "
                        f"{max_gap/min_gap:.1f}x (min={min_gap:.2f}in, "
                        f"max={max_gap:.2f}in); "
                        f"unify to consistent rhythm (0.16-0.32in)"
                    )

    # 56. misaligned_elements — 同行同类型元素左边界偏移 > 0.25in
    #     Check that shapes in the same row (top within 0.3in of each other)
    #     AND with similar widths (within 2x of each other) have consistent
    #     left edges. Excludes full-width elements (>80% of slide width).
    if len(shape_bboxes) >= 3:
        # Group shapes by approximate row (similar top position)
        rows = {}
        for b in shape_bboxes:
            if not b or b[3] < 0.2:
                continue
            row_key = round(b[1] / 0.3)  # bucket by 0.3in increments
            rows.setdefault(row_key, []).append(b)
        for row_key, row_shapes in rows.items():
            if len(row_shapes) < 2:
                continue
            # Filter: only compare shapes with similar widths (ratio < 2x)
            # and exclude full-width shapes (>80% slide width)
            filtered = [b for b in row_shapes
                        if b[2] < slide_w_in * 0.80 and b[2] > 0.5]
            if len(filtered) < 2:
                continue
            # Group by width similarity (within 2x)
            width_groups = {}
            for b in filtered:
                # Find existing group with similar width
                placed = False
                for w_key, group in width_groups.items():
                    if min(w_key, b[2]) / max(w_key, b[2]) > 0.5:
                        group.append(b)
                        placed = True
                        break
                if not placed:
                    width_groups[b[2]] = [b]
            for w_key, group in width_groups.items():
                if len(group) < 2:
                    continue
                left_edges = [b[0] for b in group]
                span = max(left_edges) - min(left_edges)
                if span > 0.25:
                    result["anti_patterns"].append(
                        f"misaligned_elements: {len(group)} similar shapes in row "
                        f"have left-edge span of {span:.2f}in; "
                        f"align to grid"
                    )

    # 57. chart_without_message — 图表/表格缺少标题或结论性标注
    #     Check: if a slide has a TABLE but no text block with >6 chars
    #     that is NOT inside a table, flag as chart_without_message.
    #     Use text_blocks (which have shape_name) instead of shapes dict
    #     (which lacks a "text" key).
    has_table = any("TABLE" in s.get("type", "") or "TABLE" in s.get("name", "").upper()
                    for s in result["shapes"])
    if has_table:
        table_names = {s.get("name", "") for s in result["shapes"]
                       if "TABLE" in s.get("type", "") or "TABLE" in s.get("name", "").upper()}
        has_message = any(
            len(tb["text"].strip()) > 6 and tb.get("shape_name", "") not in table_names
            for tb in result["text_blocks"]
        )
        if not has_message:
            result["anti_patterns"].append(
                "chart_without_message: table/chart present but no title "
                "or insight annotation; add So what? summary"
            )
    text_color_families = set()
    for h in result["colors_used"].keys():
        fam = _color_family(h)
        if fam != "neutral":
            text_color_families.add(fam)
    if len(text_color_families) > 3:
        result["anti_patterns"].append(
            f"rainbow_text: {len(text_color_families)} distinct non-neutral "
            f"color families in text ({text_color_families}); "
            f"limit to 2-3 per 60-30-10 rule"
        )

    # 23. tiny_text — 正文字号 < 12pt
    tiny_sizes = [p for p in font_sizes if p < 12]
    if tiny_sizes:
        result["anti_patterns"].append(
            f"tiny_text: {len(tiny_sizes)} run(s) below 12pt "
            f"(min={min(font_sizes):.1f}pt); minimum readable size is 12pt"
        )

    # 24. no_visual_hierarchy — 标题与正文字号比 < 1.2 且无粗细差异
    if len(font_sizes) >= 2:
        size_counts = {}
        for p in font_sizes:
            size_counts[round(p, 1)] = size_counts.get(round(p, 1), 0) + 1
        if len(size_counts) <= 2:  # only 1-2 distinct sizes
            fs_max = max(font_sizes)
            fs_min = min(font_sizes)
            ratio = fs_max / fs_min if fs_min > 0 else 0
            if ratio < 1.2:
                result["anti_patterns"].append(
                    f"no_visual_hierarchy: only {len(size_counts)} distinct "
                    f"font sizes, ratio={ratio:.2f}; "
                    f"need ≥ 3 sizes with ratio ≥ 1.25 for visible hierarchy"
                )

    # 25. orphan_widow — 多行段落末行字符数 < 5（CJK < 3）
    #     Skip page numbers, chapter markers, and single-line blocks
    #     (orphan/widow only applies to multi-line paragraphs)
    for tb in result["text_blocks"]:
        text = tb["text"]
        if _is_page_number_text(text):
            continue
        lines = text.split("\n")
        # Only check multi-line paragraphs (orphan/widow = last line of a wrapped paragraph)
        if len(lines) < 2:
            continue
        last_line = lines[-1].strip()
        if not last_line:
            continue
        # Compare last line to the average of other lines
        other_lines = [l.strip() for l in lines[:-1] if l.strip()]
        if not other_lines:
            continue
        avg_other_len = sum(len(l) for l in other_lines) / len(other_lines)
        # Only flag if last line is much shorter than other lines (true orphan)
        if len(last_line) < 5 and len(last_line) < avg_other_len * 0.4:
            result["anti_patterns"].append(
                f"orphan_widow: last line of text block has only "
                f"{len(last_line)} chars ('{last_line[:20]}'); "
                f"adjust line spacing or text box width"
            )
            break  # report once per slide

    # 26. bullet_soup — 要点 > 4 且无缩进/分组
    #     Exclude page numbers and short non-content markers
    soup_blocks = [
        tb for tb in result["text_blocks"]
        if not _is_page_number_text(tb["text"]) and len(tb["text"]) > 4
    ]
    if len(soup_blocks) > 4:
        # Simple heuristic: > 4 content blocks with no clear size variation = bullet_soup
        block_lengths = [tb["length"] for tb in soup_blocks]
        if block_lengths:
            length_range = max(block_lengths) - min(block_lengths)
            if length_range < 30:  # all blocks similar length = no grouping
                result["anti_patterns"].append(
                    f"bullet_soup: {len(soup_blocks)} text blocks "
                    f"with similar lengths (range={length_range}); "
                    f"group into 2-3 categories or use visual layout"
                )

    # 27. low_contrast — 文字颜色与背景对比度不足 (WCAG AA)
    #     Uses actual slide background and WCAG contrast ratio.
    #     Only checks TEXT run colors (not shape fill/line colors).
    #     slide_bg_hex and is_dark_bg already detected above.
    # Default background: white if we can't detect
    bg_hex = slide_bg_hex or "FFFFFF"
    bg_desc = f"#{bg_hex}" if slide_bg_hex else "white (assumed)"
    for h, count in result["text_colors"].items():
        ratio = _contrast_ratio(h, bg_hex)
        # Only flag if contrast is below WCAG AA (4.5:1 for body, 3.0:1 for large)
        if ratio < 3.0:
            result["anti_patterns"].append(
                f"low_contrast: text color #{h} has low contrast "
                f"against {bg_desc} background (ratio≈{ratio:.1f}:1); "
                f"WCAG AA requires ≥ 4.5:1 for body text"
            )
            break  # report once

    # 28. default_template — 检测 Office 默认模板名称
    DEFAULT_TEMPLATES = {
        "wisp", "median", "newsroom", "facet", "badge", "celestial",
        "slice", "integral", "retrospect", "parallax", "banded",
        "capital", "contour", "crop", "depth", "ductape", "frames",
        "grid", "headlines", "madison", "mainevent", "mesh",
        "metropolitan", "moilere", "office", "organic", "otilia",
        "savon", "slate", "view", "whisp", "ion", "badge",
    }
    if slide_layout_name:
        layout_lower = slide_layout_name.lower()
        if any(dt in layout_lower for dt in DEFAULT_TEMPLATES):
            result["anti_patterns"].append(
                f"default_template: layout '{slide_layout_name}' appears "
                f"to be an Office default template; create a custom design"
            )

    # 29. top_text_bottom_image — 固定上文下图布局
    #     所有文字在上半区、所有图片在下半区，导致图文割裂
    if result["images"] and result["text_blocks"]:
        text_bboxes_all = [
            b for s, b in zip(result["shapes"], shape_bboxes)
            if b and s.get("has_text_frame")
        ]
        img_bboxes_all = [img["bbox_in"] for img in result["images"] if img.get("bbox_in")]
        if text_bboxes_all and img_bboxes_all:
            text_centers_y = [(b[1] + b[3] / 2.0) for b in text_bboxes_all]
            img_centers_y = [(b[1] + b[3] / 2.0) for b in img_bboxes_all]
            text_avg_y = sum(text_centers_y) / len(text_centers_y)
            img_avg_y = sum(img_centers_y) / len(img_centers_y)
            mid_y = slide_h_in / 2.0
            # All text in top half, all images in bottom half
            text_all_top = all(cy < mid_y for cy in text_centers_y)
            img_all_bottom = all(cy > mid_y for cy in img_centers_y)
            if text_all_top and img_all_bottom and (img_avg_y - text_avg_y) > slide_h_in * 0.25:
                result["anti_patterns"].append(
                    f"top_text_bottom_image: all text in top half, all images "
                    f"in bottom half (text avg y={text_avg_y:.1f}, img avg y={img_avg_y:.1f}); "
                    f"use left-text-right-image or card grid layout instead"
                )

    # 30. no_card_containers — 内容页无卡片/容器
    #     >3 个文本块且 >20 词，但没有圆角矩形作为信息容器
    content_block_count = len([
        tb for tb in result["text_blocks"]
        if not _is_page_number_text(tb["text"]) and len(tb["text"]) > 4
    ])
    has_cards = len(rounded_rect_bboxes) > 0
    if content_block_count > 3 and total_words > 20 and not has_cards:
        result["anti_patterns"].append(
            f"no_card_containers: {content_block_count} content text blocks "
            f"with no rounded rectangle containers; text is spread directly "
            f"on background — add semi-transparent cards to group information"
        )

    # 31. code_without_monospace — 代码/路径/函数名未使用等宽字体
    import re as _re
    CODE_PATTERNS = [
        r'\w+\.\w{2,4}\b',           # file.ext (SKILL.md, config.json)
        r'def\s+\w+',                  # def function
        r'\w+\([^)]*\)',              # function(args)
        r'[A-Z]:\\',                  # Windows path
        r'/\w+/\w+',                  # Unix path
        r'\$\{?\w+\}?',               # ${VAR} / $VAR
        r'--\w+',                     # CLI flags --flag
        r'```',                       # code blocks
    ]
    code_text_found = False
    monospace_fonts = {"consolas", "courier", "mono", "menlo", "fira code",
                       "source code pro", "jetbrains mono", "cascadia",
                       "ubuntu mono", "dejavu sans mono"}
    for tb in result["text_blocks"]:
        text = tb["text"]
        for pattern in CODE_PATTERNS:
            if _re.search(pattern, text):
                code_text_found = True
                # Check if this text block uses a monospace font
                shape_name = tb.get("shape_name", "")
                block_fonts = [
                    fn for fn in result["fonts_used"].keys()
                    if not fn.startswith("(theme)")
                ]
                has_mono = any(
                    any(ms in fn.lower() for ms in monospace_fonts)
                    for fn in block_fonts
                )
                if not has_mono:
                    result["anti_patterns"].append(
                        f"code_without_monospace: code/path text '{text[:30]}' "
                        f"found but no monospace font in use (fonts: {block_fonts[:3]}); "
                        f"use Consolas/Courier New for code, paths, function names"
                    )
                    break  # report once per slide
        if code_text_found:
            break

    # 32. key_data_not_emphasized — 关键数据（数字/百分比）未放大强调
    #     检测正文文本中的数字/百分比，检查是否有更大的字号
    KEY_DATA_PATTERN = _re.compile(r'\b\d+[%+]?\b|\b\d+,\d{3}\b')
    if font_sizes and len(font_sizes) >= 2:
        max_fs = max(font_sizes)
        body_fs_candidates = [fs for fs in font_sizes if fs < max_fs]
        body_fs = min(body_fs_candidates) if body_fs_candidates else max_fs
        for tb in result["text_blocks"]:
            text = tb["text"]
            if _is_page_number_text(text):
                continue
            # Look for key data patterns: percentages, large numbers, time values
            data_matches = KEY_DATA_PATTERN.findall(text)
            if data_matches:
                # Check if any run in this text block has a larger font
                # We don't have per-block font sizes, so use heuristic:
                # if body_fs is used and no run is >= body_fs * 1.5, flag it
                has_large_run = any(fs >= body_fs * 1.5 for fs in font_sizes)
                # Only flag if there ARE large fonts available but the slide
                # doesn't seem to use them for data emphasis
                # (This is a soft check — too noisy if no large fonts exist at all)
                if not has_large_run and len(data_matches) >= 2:
                    result["anti_patterns"].append(
                        f"key_data_not_emphasized: slide contains key data "
                        f"({', '.join(data_matches[:3])}) but no font size "
                        f">= {body_fs * 1.5:.0f}pt for emphasis; "
                        f"enlarge numbers/percentages to 1.5-2x body size"
                    )
                    break  # report once

    # 33. image_at_edge — 图片紧贴幻灯片边缘（投屏易裁切）
    for img in result["images"]:
        ib = img.get("bbox_in")
        if ib and _is_edge_touching(ib, slide_w_in, slide_h_in, margin_pct=0.03):
            result["anti_patterns"].append(
                f"image_at_edge: image at ({ib[0]:.1f},{ib[1]:.1f}) "
                f"is within 3% of slide edge; add safety margin to avoid "
                f"projection cropping"
            )
            break  # report once

    # 34. bare_text_no_card — 文本框不在任何卡片容器区域内
    #     检测文本框 bbox 是否落在任何圆角矩形（卡片）bbox 内部，
    #     如果大量文本在卡片外，说明信息未分组包裹
    #     排除：含表格的页面（表格或矩形网格自带容器）
    has_table = any(
        (s.get("auto_shape_type") or "").upper() == "TABLE"
        or str(s.get("type", "")) == "TABLE (19)"
        for s in result["shapes"]
    )
    # Also detect manual tables: ≥6 same-height RECTANGLEs arranged in a grid
    if not has_table and len(result["shapes"]) >= 6:
        rect_heights = []
        for s in result["shapes"]:
            bb = s.get("bbox_in")
            at = (s.get("auto_shape_type") or "").upper()
            if bb and at == "RECTANGLE":
                rect_heights.append(round(bb[3], 2))
        if rect_heights:
            from collections import Counter as _C5
            height_counts = _C5(rect_heights)
            most_common_h, most_common_count = height_counts.most_common(1)[0]
            if most_common_count >= 6:
                has_table = True
    if rounded_rect_bboxes and len(result["text_blocks"]) > 3 and not has_table:
        def _rect_contains(outer, inner):
            """True if inner rect is mostly inside outer rect."""
            ol, ot, ow, oh = outer
            il, it, iw, ih = inner
            # inner center inside outer?
            icx, icy = il + iw / 2, it + ih / 2
            return ol <= icx <= ol + ow and ot <= icy <= ot + oh
        card_bboxes_for_check = rounded_rect_bboxes
        text_bboxes_for_check = [
            (s.get("bbox_in", [0, 0, 0, 0]) if s.get("bbox_in") else (0, 0, 0, 0))
            for s in result["shapes"]
            if s.get("has_text_frame") and s.get("bbox_in")
        ]
        bare_count = 0
        for tb in text_bboxes_for_check:
            if not any(_rect_contains(cb, tb) for cb in card_bboxes_for_check):
                bare_count += 1
        total_tb = len(text_bboxes_for_check)
        if total_tb > 0 and bare_count / total_tb > 0.6:
            result["anti_patterns"].append(
                f"bare_text_no_card: {bare_count}/{total_tb} text boxes "
                f"are outside card containers; wrap text in semi-transparent "
                f"cards to group information visually"
            )

    # 35. table_no_zebra — 表格无隔行底色
    #     检测 Table 形状，如果行数 > 3 且没有交替行填充色,
    #     判定为缺少斑马条纹
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                tbl = shape.table
                if len(tbl.rows) > 3:
                    row_fills = []
                    for ri, row in enumerate(tbl.rows):
                        # Check first cell's fill color
                        cell = row.cells[0]
                        try:
                            cf = cell.fill
                            if cf.type is not None:
                                row_fills.append(_rgb_to_hex(cf.fore_color.rgb) if cf.fore_color and cf.fore_color.rgb else "none")
                            else:
                                row_fills.append("none")
                        except Exception:
                            row_fills.append("none")
                    # If all fills are "none" or all the same → no zebra
                    unique_fills = set(row_fills)
                    if len(unique_fills) <= 1 or (len(unique_fills) == 2 and "none" in unique_fills):
                        result["anti_patterns"].append(
                            f"table_no_zebra: table with {len(tbl.rows)} rows "
                            f"has no alternating row colors; add zebra striping "
                            f"for readability"
                        )
            except Exception:
                pass

    # 36. table_weak_header — 表头行与正文视觉无区分
    #     表格第一行与后续行的字体大小/粗细/颜色无明显差异
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                tbl = shape.table
                if len(tbl.rows) >= 2:
                    hdr_bold = None
                    hdr_size = None
                    hdr_color = None
                    body_bold = None
                    body_size = None
                    body_color = None
                    # Sample first cell of header and second row
                    for pi, para in enumerate(tbl.rows[0].cells[0].text_frame.paragraphs):
                        for run in para.runs:
                            if run.font.bold is not None:
                                hdr_bold = run.font.bold
                            if run.font.size:
                                hdr_size = run.font.size.pt
                            try:
                                if run.font.color and run.font.color.rgb:
                                    hdr_color = _rgb_to_hex(run.font.color.rgb)
                            except Exception:
                                pass
                            break
                        break
                    for pi, para in enumerate(tbl.rows[1].cells[0].text_frame.paragraphs):
                        for run in para.runs:
                            if run.font.bold is not None:
                                body_bold = run.font.bold
                            if run.font.size:
                                body_size = run.font.size.pt
                            try:
                                if run.font.color and run.font.color.rgb:
                                    body_color = _rgb_to_hex(run.font.color.rgb)
                            except Exception:
                                pass
                            break
                        break
                    # Check if header is visually distinct
                    distinct = False
                    if hdr_bold and not body_bold:
                        distinct = True
                    if hdr_size and body_size and hdr_size > body_size * 1.1:
                        distinct = True
                    if hdr_color and body_color and hdr_color != body_color:
                        distinct = True
                    if not distinct:
                        result["anti_patterns"].append(
                            f"table_weak_header: table header row is not visually "
                            f"distinct from body rows; add header background color, "
                            f"bold font, or larger size to differentiate"
                        )
            except Exception:
                pass

    # 37. mixed_numbering — 同页混用不同序号格式
    #     检测纯数字(1,2,3)、带前导零(01,02)、圆形标识(①②)、括号(1) 2) 3))
    #     同时出现两种以上视为混用
    import re as _re2
    _NUM_STYLES = {
        "plain": _re2.compile(r'(?<!\w)\d{1,2}(?!\w)'),
        "leading_zero": _re2.compile(r'(?<!\w)0\d(?!\w)'),
        "circled": _re2.compile(r'[①②③④⑤⑥⑦⑧⑨⑩]'),
        "paren": _re2.compile(r'(?<!\w)\d{1,2}[)\），]'),
    }
    found_styles = set()
    for tb in result["text_blocks"]:
        text = tb["text"]
        for style_name, pattern in _NUM_STYLES.items():
            if pattern.search(text):
                found_styles.add(style_name)
    if len(found_styles) >= 3:
        result["anti_patterns"].append(
            f"mixed_numbering: {len(found_styles)} numbering styles on one "
            f"slide ({', '.join(sorted(found_styles))}); unify to a single "
            f"numbering system per slide"
        )

    # 38. unbalanced_layout — 左右内容区面积比悬殊
    #     将 slide 沿中轴分为左右两半，统计每半的元素总面积，
    #     如果比值 > 3:1 且较小侧有内容，标记为失衡
    #     排除：低文字量页面（封面/结尾），这些页面偏移布局是设计意图
    left_area = 0.0
    right_area = 0.0
    if len(shape_bboxes) >= 4 and total_words >= 10:
        mid_x = slide_w_in / 2.0
        for bbox in shape_bboxes:
            l, t, w, h = bbox
            cx = l + w / 2.0
            area = w * h
            if cx < mid_x:
                left_area += area
            else:
                right_area += area
        if left_area > 0 and right_area > 0:
            ratio = max(left_area, right_area) / min(left_area, right_area)
            if ratio > 3.0:
                heavier = "left" if left_area > right_area else "right"
                result["anti_patterns"].append(
                    f"unbalanced_layout: left/right area ratio {ratio:.1f}:1 "
                    f"({heavier} side heavier); balance content distribution "
                    f"or use asymmetric layout intentionally"
                )

    # ---- Layout signals ----
    result["layout_signals"] = {
        "shape_count": len(result["shapes"]),
        "text_block_count": len(result["text_blocks"]),
        "total_text_chars": total_text,
        "total_words": total_words,
        "image_count": len(result["images"]),
        "left_edges_in": sorted(round(l, 2) for l in left_edges),
        "top_edges_in": sorted(round(t, 2) for t in top_edges),
        "distinct_color_count": len(result["colors_used"]),
        "distinct_font_count": len(result["fonts_used"]),
        "font_families": font_families,
        "has_cjk_text": has_cjk_text,
        "cjk_fonts": cjk_fonts,
        "latin_fonts": latin_fonts,
        "font_size_range_pt": (
            [min(font_sizes), max(font_sizes)] if font_sizes else None
        ),
        "font_size_ratio": (
            round(max(font_sizes) / min(font_sizes), 3)
            if len(font_sizes) >= 2 and min(font_sizes) > 0 else None
        ),
        "page_fill_rate": round(page_fill_rate, 3),
        "element_area_sqin": round(element_area_sum, 2),
        "quadrants_filled": quadrants,
        "empty_quadrants": empty_quadrants,
        # New fields for deck-level checks
        "has_cards": len(rounded_rect_bboxes) > 0,
        "card_count": len(rounded_rect_bboxes),
        "content_block_count": content_block_count,
        "text_colors": dict(result["text_colors"]) if result["text_colors"] else {},
        "is_dark_bg": is_dark_bg,
        "slide_bg_color": slide_bg_hex,
        # Short text snippets for cross-slide repetition check
        "short_texts": [
            tb["text"][:50] for tb in result["text_blocks"]
            if not _is_page_number_text(tb["text"]) and len(tb["text"]) < 30
        ],
        # Left/right area for unbalanced_layout deck check
        "left_area_ratio": round(left_area / (left_area + right_area), 3) if (left_area + right_area) > 0 else 0.5,
        # Section labels for cross-slide position consistency check
        "section_labels": [
            {"text": tb["text"][:50], "shape_idx": i}
            for i, tb in enumerate(result["text_blocks"])
            if _is_section_label(tb["text"])
        ],
        # Per-shape bbox list with text for header nav detection
        "shape_bboxes_with_text": [
            (s.get("bbox_in"), s.get("has_text_frame", False))
            for s in result["shapes"]
        ],
    }

    # Convert counters to plain dicts for JSON
    result["colors_used"] = dict(result["colors_used"])
    result["fonts_used"] = dict(result["fonts_used"])
    
    # 改进项 2：提取设计 Token（多版本对比 / 前端 token 对接 / 报告诊断用）
    result["tokens"] = _extract_design_tokens(result, slide_w_in, slide_h_in)
    
    return result


def _contains_cjk(text: str) -> bool:
    """True if the string contains CJK Unified Ideographs."""
    return any("\u4e00" <= ch <= "\u9fff" for ch in text)


def _is_section_label(text: str) -> bool:
    """True if text looks like a section/chapter label.
    Patterns: 'SECTION 01xxx', '第1章', '01 某某', 'PART I', etc.
    """
    if not text or len(text) > 80:
        return False
    import re as _re3
    patterns = [
        r'^SECTION\s+\d',          # SECTION 01 ...
        r'^PART\s+[IVX\d]',        # PART I / PART 1
        r'^第[一二三四五六七八九十\d]+[章节部分]',  # 第1章 / 第二章
        r'^0\d\s',                  # 01 标题
        r'^CHAPTER\s+\d',           # CHAPTER 1
    ]
    for p in patterns:
        if _re3.search(p, text, _re3.IGNORECASE):
            return True
    return False


# --------------------------------------------------------------------------- #
# Deck-level aggregation
# --------------------------------------------------------------------------- #
def summarize_deck(slides_data: list[dict]) -> dict:
    """Compute deck-wide consistency signals."""
    all_fonts = Counter()
    all_colors = Counter()
    per_slide_issue_counts = []
    for s in slides_data:
        all_fonts.update(s["fonts_used"])
        all_colors.update(s["colors_used"])
        per_slide_issue_counts.append(len(s["anti_patterns"]))

    # Infer slide width from the first slide's shape bboxes (for header checks)
    slide_w_in = 13.333  # default 16:9
    try:
        for sr in slides_data:
            for sh in sr.get("shapes", []):
                bb = sh.get("bbox_in")
                if bb and bb[0] + bb[2] > slide_w_in:
                    slide_w_in = bb[0] + bb[2]
            break
    except Exception:
        pass

    deck_signals = {
        "total_slides": len(slides_data),
        "deck_font_families": list(all_fonts.keys()),
        "deck_font_count": len(all_fonts),
        "deck_color_count": len(all_colors),
        "deck_top_colors": all_colors.most_common(8),
        "deck_top_fonts": all_fonts.most_common(5),
        "avg_issues_per_slide": (
            round(sum(per_slide_issue_counts) / len(per_slide_issue_counts), 2)
            if per_slide_issue_counts else 0
        ),
        "worst_slide": (
            max(range(len(per_slide_issue_counts)),
                key=lambda i: per_slide_issue_counts[i]) + 1
            if per_slide_issue_counts else None
        ),
    }

    # Deck-level anti-patterns
    deck_aps = []
    if len(all_fonts) > 4:
        deck_aps.append(
            f"deck_too_many_fonts: {len(all_fonts)} families across the deck "
            f"(should be <= 2-3)"
        )
    if len(all_colors) > 12:
        deck_aps.append(
            f"deck_too_many_colors: {len(all_colors)} distinct colors across deck"
        )

    # Cross-slide consistency checks
    # inconsistent_fonts — different slide pairs use completely different font sets
    if len(slides_data) >= 3:
        slide_font_sets = []
        for sr in slides_data:
            ff = sr.get("layout_signals", {}).get("font_families", [])
            slide_font_sets.append(set(ff))
        # Count how many slides deviate from the most common font set
        if slide_font_sets:
            from collections import Counter as _C2
            set_counts = _C2(tuple(sorted(s)) for s in slide_font_sets)
            dominant_set = set_counts.most_common(1)[0][0]
            deviant_slides = sum(
                1 for s in slide_font_sets
                if tuple(sorted(s)) != dominant_set and len(s) > 0
            )
            if deviant_slides > len(slide_font_sets) * 0.3:
                deck_aps.append(
                    f"inconsistent_fonts: {deviant_slides}/{len(slide_font_sets)} "
                    f"slides use different font families from the deck's primary "
                    f"set ({dominant_set})"
                )

    # inconsistent_colors — high color variation across slides
    if len(slides_data) >= 3:
        slide_color_sets = []
        for sr in slides_data:
            colors = set(sr.get("colors_used", {}).keys())
            # Filter out neutrals for more meaningful comparison
            non_neutral = {h for h in colors if _color_family(h) != "neutral"}
            slide_color_sets.append(non_neutral)
        # Check if there's a dominant color palette
        all_non_neutral = set()
        for s in slide_color_sets:
            all_non_neutral.update(s)
        # If > 8 non-neutral colors across deck and most slides use different subsets
        if len(all_non_neutral) > 8:
            per_slide_count = [len(s) for s in slide_color_sets if s]
            avg_per_slide = sum(per_slide_count) / len(per_slide_count) if per_slide_count else 0
            if len(all_non_neutral) > avg_per_slide * 3:
                deck_aps.append(
                    f"inconsistent_colors: {len(all_non_neutral)} non-neutral "
                    f"colors across deck but avg {avg_per_slide:.0f} per slide; "
                    f"color palette not unified"
            )

    # 37. inconsistent_margins — 跨页内容起始边距差异大
    #     计算每页第一个主要内容区域的 top 边缘和 left 边缘，
    #     如果跨页标准差过大，说明没有统一的边距/网格体系
    if len(slides_data) >= 4:
        content_tops = []
        content_lefts = []
        for sr in slides_data:
            le = sr.get("layout_signals", {}).get("left_edges_in", [])
            te = sr.get("layout_signals", {}).get("top_edges_in", [])
            if le:
                content_lefts.append(le[0])
            if te:
                content_tops.append(te[0])
        if len(content_tops) >= 4:
            import statistics as _stats
            top_std = _stats.stdev(content_tops) if len(content_tops) >= 2 else 0
            left_std = _stats.stdev(content_lefts) if len(content_lefts) >= 2 else 0
            # If standard deviation > 0.4in, margins are wildly inconsistent
            if top_std > 0.4 or left_std > 0.4:
                deck_aps.append(
                    f"inconsistent_margins: content start position varies "
                    f"across slides (top std={top_std:.2f}in, left std={left_std:.2f}in); "
                    f"define a standard content grid with fixed margins"
                )

    # 38. inconsistent_section_pos — section 标签跨页位置不统一
    #     检测包含 SECTION/第N章 等标签的文本框，比较它们跨页的位置
    if len(slides_data) >= 3:
        section_positions = []  # list of (left, top) tuples
        for sr in slides_data:
            labels = sr.get("layout_signals", {}).get("section_labels", [])
            if labels:
                # Get bbox of the first section label shape
                idx = labels[0].get("shape_idx", 0)
                shapes = sr.get("shapes", [])
                if idx < len(shapes):
                    bbox = shapes[idx].get("bbox_in")
                    if bbox:
                        section_positions.append((bbox[0], bbox[1]))
        if len(section_positions) >= 3:
            lefts = [p[0] for p in section_positions]
            tops = [p[1] for p in section_positions]
            import statistics as _stats2
            left_range = max(lefts) - min(lefts)
            top_range = max(tops) - min(tops)
            if left_range > 0.5 or top_range > 0.3:
                deck_aps.append(
                    f"inconsistent_section_pos: section labels at inconsistent "
                    f"positions (left range={left_range:.2f}in, top range={top_range:.2f}in); "
                    f"place section labels at a fixed position in the header area"
                )

    # 39. cover_ending_mismatch — 首尾版式不闭环
    #     比较第1页和最后1页的：背景色、主要形状类型、文字对齐方式
    #     如果完全不同，说明首尾视觉风格割裂
    if len(slides_data) >= 3:
        first = slides_data[0]
        last = slides_data[-1]
        first_bg = first.get("layout_signals", {}).get("slide_bg_color")
        last_bg = last.get("layout_signals", {}).get("slide_bg_color")
        first_shapes = len(first.get("shapes", []))
        last_shapes = len(last.get("shapes", []))
        # Check background match
        bg_match = (first_bg == last_bg) if first_bg and last_bg else False
        # Check if both have decorative elements (shapes > threshold)
        first_has_deco = first_shapes > 2
        last_has_deco = last_shapes > 2
        # Check if both have centered text (cover/ending typically centered)
        first_centered = any(
            tb.get("text", "") and len(tb["text"]) > 2
            for tb in first.get("text_blocks", [])[:3]
        )
        last_centered = any(
            tb.get("text", "") and len(tb["text"]) > 2
            for tb in last.get("text_blocks", [])[:3]
        )
        # If backgrounds differ AND decoration style differs, flag it
        if not bg_match and (first_has_deco != last_has_deco):
            deck_aps.append(
                f"cover_ending_mismatch: cover (bg={'set' if first_bg else 'default'}, "
                f"shapes={first_shapes}) vs ending (bg={'set' if last_bg else 'default'}, "
                f"shapes={last_shapes}); mirror visual style between "
                f"cover and ending for cohesive framing"
            )

    # 40. density_outlier — 页面填充率偏离均值过大
    #     某些页面极疏（<均值-30%）或极密（>均值+30%），疏密失衡
    if len(slides_data) >= 4:
        fill_rates = [
            sr.get("layout_signals", {}).get("page_fill_rate", 0)
            for sr in slides_data
        ]
        avg_fill = sum(fill_rates) / len(fill_rates)
        outlier_count = sum(
            1 for fr in fill_rates
            if abs(fr - avg_fill) > 0.25  # >25 percentage points deviation
        )
        if outlier_count > len(fill_rates) * 0.3:
            sparse = sum(1 for fr in fill_rates if fr < avg_fill - 0.25)
            dense = sum(1 for fr in fill_rates if fr > avg_fill + 0.25)
            deck_aps.append(
                f"density_outlier: {outlier_count}/{len(fill_rates)} slides "
                f"deviate >25% from avg fill rate ({avg_fill:.0%}); "
                f"({sparse} too sparse, {dense} too dense); "
                f"balance content density across the deck"
            )

    # 41. no_header_navigation — 无全局页眉导航
    #     如果大部分内容页缺少固定页眉区域（logo+章节名+进度条），
    #     听众无法快速定位演讲进度
    if len(slides_data) >= 4:
        # Skip cover (first) and ending (last) — only check content slides
        content_slides = slides_data[1:-1]
        slides_with_fixed_header = 0
        for sr in content_slides:
            ls = sr.get("layout_signals", {})
            top_edges = ls.get("top_edges_in", [])
            # A fixed header has a shape starting in the top 0.5in
            # and spanning most of the slide width
            has_header_shape = False
            for shape in sr.get("shapes", []):
                bbox = shape.get("bbox_in")
                # Header: starts in top 15% of slide, spans >70% of slide width
                if bbox and bbox[1] < slide_w_in * 0.12 and bbox[2] > slide_w_in * 0.7:
                    has_header_shape = True
                    break
            if has_header_shape:
                slides_with_fixed_header += 1
        if len(content_slides) > 0 and slides_with_fixed_header < len(content_slides) * 0.5:
            deck_aps.append(
                f"no_header_navigation: only {slides_with_fixed_header}/{len(content_slides)} "
                f"content slides have a fixed header bar; add consistent header "
                f"(logo + chapter title + progress bar) for navigation"
            )

    deck_signals["deck_anti_patterns"] = deck_aps

    # 34. repeated_chapter_labels — 跨页重复的短文本（章节标注/模块名）
    #     同一段短文本出现在多页上，属于无效视觉噪音
    if len(slides_data) >= 3:
        from collections import Counter as _C3
        all_short_texts = _C3()
        for sr in slides_data:
            short_texts = sr.get("layout_signals", {}).get("short_texts", [])
            for st in short_texts:
                all_short_texts[st] += 1
        # Flag texts appearing on 3+ slides (chapter labels repeat within
        # a chapter, typically 2-4 slides per chapter)
        threshold = 3
        repeated = [(t, c) for t, c in all_short_texts.items() if c >= threshold
                    and not _is_page_number_text(t)]
        if repeated:
            examples = [f"'{t}'({c}页)" for t, c in repeated[:3]]
            deck_aps.append(
                f"repeated_chapter_labels: short text appears on many slides "
                f"({', '.join(examples)}); remove redundant per-page labels, "
                f"use progress bar or header navigation instead"
            )

    # 35. inconsistent_alignment — 跨页对齐不一致
    #     不同页面的左边缘坐标差异过大，说明没有统一网格
    if len(slides_data) >= 3:
        slide_left_edges = []
        for sr in slides_data:
            le = sr.get("layout_signals", {}).get("left_edges_in", [])
            if le:
                # Use the minimum left edge (main content start)
                slide_left_edges.append(le[0])
        if slide_left_edges:
            from collections import Counter as _C4
            edge_counts = _C4(round(e, 1) for e in slide_left_edges)
            dominant_edge = edge_counts.most_common(1)[0][0]
            deviant_count = sum(1 for e in slide_left_edges if abs(round(e, 1) - dominant_edge) > 0.5)
            if deviant_count > len(slide_left_edges) * 0.3:
                deck_aps.append(
                    f"inconsistent_alignment: {deviant_count}/{len(slide_left_edges)} "
                    f"slides have different left edges (dominant={dominant_edge}in); "
                    f"establish a unified content grid with consistent margins"
                )

    # 36. no_brand_consistency — 缺少品牌元素一致性
    #     检查是否有跨页一致的品牌栏/Logo shape
    if len(slides_data) >= 3:
        slides_with_header = 0
        for sr in slides_data:
            shapes = sr.get("shapes", [])
            # Check for a full-width shape near the top (header bar)
            ls = sr.get("layout_signals", {})
            top_edges = ls.get("top_edges_in", [])
            if top_edges and min(top_edges) < slide_w_in * 0.08:  # shape near top
                slides_with_header += 1
        if slides_with_header < len(slides_data) * 0.5:
            deck_aps.append(
                f"no_brand_consistency: only {slides_with_header}/{len(slides_data)} "
                f"slides have a top header element; add consistent brand bar "
                f"(logo + title) across all content slides for professional identity"
            )

    deck_signals["deck_anti_patterns"] = deck_aps
    return deck_signals


# --------------------------------------------------------------------------- #
# Vision scoring prompt (returned in the report for the agent to use)
# --------------------------------------------------------------------------- #
SCORING_PROMPT = """分析这张 PPT 幻灯片图片的美观度。按以下 9 个维度打分（1-10）并给出一句中文理由。

维度与权重（与结构化脚本评分同维度，便于分数合并）：
1. whitespace (13%) 留白 — 负空间是否充分？不拥挤？页边距是否充裕？元素面积占页面 60-70%？无象限完全空白？
2. type_scale (14%) 字号阶梯 — 标题/副标题/正文是否有 Major Third (1.25×) 以上阶梯？字号比是否 ≥ 1.25？
3. color_harmony (14%) 配色和谐 — 是否遵循 60-30-10？无 AI 紫蓝渐变？无纯黑 #000 / 纯白 #FFF 滥用？
4. alignment (14%) 对齐与网格 — 元素是否对齐到一致网格？页边距是否均匀？间距是否遵循 4/8px 节奏（0.08in 倍数）？
5. spacing (9%) 间距 — 元素间距是否一致？段落间距是否合理？元素间有无碰撞？
6. imagery (9%) 图像完整性 — 图片是否高清无拉伸？无占位符？图表有数据来源？
7. consistency (9%) 视觉一致 — 圆角、阴影、按钮、图标尺寸/风格是否统一？
8. hierarchy (8%) 层次 — 是否有清晰焦点？3 等大卡片？占位符文本(TODO/TBD)？字号无层级？
9. layout (10%) 布局 — 容器填充率是否 > 60%？卡片是否有空旷内腔？四象限是否平衡？图片是否触边？

等级映射：8.5-10=A, 7.0-8.4=B, 5.5-6.9=C, 4.0-5.4=D, 1.0-3.9=F

附加扣分（自动检测到以下反模式时，对应维度 -1 到 -3 分）：
- AI 紫蓝渐变 (#6366F1/#818CF8 等) → color_harmony -2
- 3 等大卡片无层级 → hierarchy -2
- 文字墙 (>60 词正文) → whitespace -2
- Comic Sans / Papyrus 等不专业字体 → type_scale -3
- 图片拉伸变形 → imagery -2
- 内容触碰边缘无页边距 → whitespace -2
- 超过 3 种主色 → color_harmony -2
- 字号比 < 1.25（无层级）→ hierarchy -1
- 强调色覆盖 > 40% 元素（违反 60-30-10）→ color_harmony -1
- 纯黑 #000 文字 / 纯白 #FFF 卡片 → color_harmony -1
- 单条要点 > 35 字（中文）→ hierarchy -1
- 间距不遵循 4/8px 节奏 → alignment -1
- 低分辨率图片 (< 800×600) → imagery -1
- 页面填充率 < 40%（空旷）或 > 85%（拥挤）→ whitespace -2
- 象限空白（1/4 页面无元素）→ whitespace -1

视觉模型额外检查（以下问题结构化脚本无法检测，需视觉判断）：
- 固定上文下图布局（文字全在上半区、图片全在下半区）→ alignment -3, hierarchy -2
- 内容无卡片/容器包裹（文字直接平铺在背景上）→ consistency -2, hierarchy -1
- 图片底色与幻灯片背景不统一（如深色页面上白底插图）→ consistency -3
- 图片荧光/高亮特效抢视觉（亮度超过正文文字）→ hierarchy -2
- 代码/路径/函数名未用等宽字体 → type_scale -2
- 关键数据（百分比/大数字/时间值）未放大强调 → hierarchy -2
- 缺少品牌栏/Logo/页眉导航 → consistency -2
- 重复的章节标注文字（每页重复同一短文本）→ consistency -1
- 章节之间无视觉过渡/分隔 → consistency -1
- 插图仅装饰无信息价值（无标注/分层/逻辑）→ imagery -2
- 配色层级单一（缺少中间过渡色区分信息等级）→ color_harmony -2
- 高亮色功能过载（同一颜色承担多种语义角色）→ color_harmony -2

严格输出 JSON（不要 markdown 代码块）：
{
  "dimensions": {
    "whitespace": {"score": N, "reason": "..."},
    "type_scale": {"score": N, "reason": "..."},
    "color_harmony": {"score": N, "reason": "..."},
    "alignment": {"score": N, "reason": "..."},
    "spacing": {"score": N, "reason": "..."},
    "imagery": {"score": N, "reason": "..."},
    "consistency": {"score": N, "reason": "..."},
    "hierarchy": {"score": N, "reason": "..."},
    "layout": {"score": N, "reason": "..."}
  },
  "weighted_total": N,
  "overall_grade": "A|B|C|D|F",
  "top_3_issues": ["...", "...", "..."],
  "improvement_suggestion": "..."
}
"""


# --------------------------------------------------------------------------- #
# Scenario-aware prompt builder (改进项 4)
# --------------------------------------------------------------------------- #
def _build_scenario_prompt(scenario: str, weights: dict) -> str:
    """根据场景权重生成定制化的视觉评分 Prompt。"""
    weight_lines = []
    dim_names = [
        ("whitespace", "留白 — 负空间是否充分？不拥挤？页边距是否充裕？元素占比是否在 60-70% 区间？"),
        ("type_scale", "字号阶梯 — 标题/副标题/正文是否有 Major Third (1.25×) 以上阶梯？"),
        ("color_harmony", "配色和谐 — 是否遵循 60-30-10？无 AI 紫蓝渐变？无纯黑 #000 / 纯白 #FFF 滥用？"),
        ("alignment", "对齐与网格 — 元素是否对齐到一致网格？页边距是否均匀？间距是否遵循 4/8px 节奏？"),
        ("spacing", "间距 — 元素间距是否一致？段落间距是否合理？"),
        ("imagery", "图像完整性 — 图片是否高清无拉伸？无占位符？图表有数据来源？"),
        ("consistency", "视觉一致 — 圆角、阴影、按钮、图标尺寸/风格是否统一？"),
        ("hierarchy", "层次 — 是否有清晰焦点？3 等大卡片？占位符文本(TODO/TBD)？字号无层级？"),
        ("layout", "布局 — 页面元素分布是否均衡？无空旷容器（填充率<60%）？无空白象限？图片不触边？"),
    ]
    for key, desc in dim_names:
        w = int(round(weights.get(key, 0) * 100))
        weight_lines.append(f"- {key} ({w}%) {desc}")
    weights_block = "\n".join(weight_lines)
    
    scenario_note = {
        "executive": "\n场景提示：本 PPT 为高管摘要/董事会场景，重点评估信息传递效率和视觉克制。",
        "marketing": "\n场景提示：本 PPT 为营销 deck 场景，重点评估视觉冲击力、配色张力和专业度。",
        "data": "\n场景提示：本 PPT 为数据密集场景（技术评审/季报），重点评估信息密度和可读性，留白可适当让步。",
        "gov": "\n场景提示：本 PPT 为党政公文场景，重点评估对齐、字体严肃性和视觉一致性。",
        "creative": "\n场景提示：本 PPT 为创意作品集场景，重点评估视觉冲击、配色大胆、留白慷慨。",
        "telecom": "\n场景提示：本 PPT 为中国电信品牌场景，重点评估配色和谐（电信红 #E60012 规范）、字号阶梯和视觉一致性。",
        "default": "",
    }.get(scenario, "")
    
    return f"""分析这张 PPT 幻灯片图片的美观度。按以下 9 个维度打分（1-10）并给出一句中文理由。
{scenario_note}
维度与权重：
{weights_block}

等级映射：8.5-10=A, 7.0-8.4=B, 5.5-6.9=C, 4.0-5.4=D, 1.0-3.9=F

附加扣分（自动检测到以下反模式时，对应维度 -1 到 -3 分）：
- AI 紫蓝渐变背景（#6366F1 / #818CF8 等）→ color_harmony -2
- 3 个等大卡片无层级 → hierarchy -2
- 文字墙（>60 词正文）→ whitespace -2
- Comic Sans / Papyrus 等不专业字体 → type_scale -3
- 图片拉伸变形 → imagery -2
- 内容触碰边缘无页边距 → whitespace -2
- 超过 3 种主色 → color_harmony -2
- 3D 图表 → imagery -2
- 默认 PowerPoint 模板未改 → consistency -2
- 空旷容器（卡片内填充率 < 60%）→ layout -2
- 四象限中有空白象限 → layout -1
- 图片触边无安全边距 → layout -1

严格输出 JSON（不要 markdown 代码块）：
{{
  "dimensions": {{
    "whitespace": {{"score": N, "reason": "..."}},
    "type_scale": {{"score": N, "reason": "..."}},
    "color_harmony": {{"score": N, "reason": "..."}},
    "alignment": {{"score": N, "reason": "..."}},
    "spacing": {{"score": N, "reason": "..."}},
    "imagery": {{"score": N, "reason": "..."}},
    "consistency": {{"score": N, "reason": "..."}},
    "hierarchy": {{"score": N, "reason": "..."}},
    "layout": {{"score": N, "reason": "..."}}
  }},
  "weighted_total": N,
  "overall_grade": "A|B|C|D|F",
  "top_3_issues": ["...", "...", "..."],
  "improvement_suggestion": "..."
}}
"""


# --------------------------------------------------------------------------- #
# Main
# --------------------------------------------------------------------------- #
def main():
    ap = argparse.ArgumentParser(description="PPT Aesthetic Structural Scorer")
    ap.add_argument("pptx_path", help="Path to the .pptx file")
    ap.add_argument("--output", "-o", default=None,
                    help="Output JSON report path (default: <pptx>.aesthetic.json)")
    ap.add_argument("--slides", "-s", default=None,
                    help="Comma-separated 1-indexed slide numbers to analyze")
    # 改进项 4：场景化权重切换
    ap.add_argument("--scenario", default="default",
                    choices=list(SCENARIO_WEIGHTS.keys()),
                    help="Scene-based weight preset (default|executive|marketing|data|gov|creative|telecom). "
                         "Affects the weighting guidance embedded in the report; "
                         "does not change anti-pattern detection.")
    args = ap.parse_args()

    pptx_path = Path(args.pptx_path)
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    slide_w_in = prs.slide_width / EMU_PER_INCH
    slide_h_in = prs.slide_height / EMU_PER_INCH

    # Resolve theme major/minor fonts so we can attribute runs that don't
    # set an explicit font name (very common — they inherit from theme).
    theme_major, theme_minor = _resolve_theme_fonts(prs)
    if theme_major or theme_minor:
        print(f"Theme fonts: major={theme_major!r}, minor={theme_minor!r}")

    if args.slides:
        idxs = [int(s.strip()) - 1 for s in args.slides.split(",") if s.strip().isdigit()]
        idxs = [i for i in idxs if 0 <= i < total]
    else:
        idxs = list(range(total))

    print(f"Analyzing {len(idxs)}/{total} slides from: {pptx_path.name}")
    print(f"Slide size: {slide_w_in:.2f} x {slide_h_in:.2f} in")

    slides_data = []
    for idx in idxs:
        slide = prs.slides[idx]
        analysis = analyze_slide(
            slide, idx, slide_w_in, slide_h_in,
            theme_major_font=theme_major, theme_minor_font=theme_minor,
        )
        slides_data.append(analysis)
        n_aps = len(analysis["anti_patterns"])
        n_fonts = analysis["layout_signals"]["distinct_font_count"]
        n_colors = analysis["layout_signals"]["distinct_color_count"]
        n_words = analysis["layout_signals"]["total_words"]
        print(f"  Slide {idx+1}: {n_fonts} fonts, {n_colors} colors, "
              f"{n_words} words, {n_aps} anti-patterns")
        for ap_ in analysis["anti_patterns"]:
            print(f"      - {ap_}")

    deck_summary = summarize_deck(slides_data)
    
    # 改进项 4：场景化权重
    scenario_weights = SCENARIO_WEIGHTS[args.scenario]
    
    report = {
        "file": pptx_path.name,
        "slide_size_in": [round(slide_w_in, 3), round(slide_h_in, 3)],
        "total_slides": total,
        "analyzed_slides": len(idxs),
        "scenario": args.scenario,
        "scenario_weights": scenario_weights,
        "deck_summary": deck_summary,
        "slides": slides_data,
        "vision_scoring_prompt": SCORING_PROMPT,
    }
    
    # 改进项 4：根据场景权重调整 vision_scoring_prompt 中的权重描述
    if args.scenario != "default":
        report["vision_scoring_prompt"] = _build_scenario_prompt(args.scenario, scenario_weights)

    out_path = args.output or str(pptx_path.with_suffix(".aesthetic.json"))
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)
    print(f"\nStructural report saved to: {out_path}")
    print("\nNext step: render slides to PNG (scripts/render_slides.py) and run")
    print("the vision_scoring_prompt on each image with an AI vision model.")
    print("Combine the structural anti_patterns above with vision scores for a")
    print("final grade per slide.")


if __name__ == "__main__":
    main()
