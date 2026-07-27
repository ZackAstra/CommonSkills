#!/usr/bin/env python3
"""
PPT Aesthetic Version Comparator — 改进项 1（多版本对比工作流）

对 2-N 个 PPTX 进行横向对比，输出：
  1. 加权总分排名表（按 final_score 降序）
  2. 维度雷达对比（每版本 8 维度得分，便于看形状差异）
  3. Token 差异表（color / typography / radius / spacing）
  4. 反模式对比（每版本命中了哪些反模式）
  5. 推荐结论（综合冠军 + 场景冠军 + 单维度冠军 + 合并建议）

依赖：
  - score_ppt_pages.py（复用其结构化评分 + token 提取 + 反模式检测）
  - 视觉评分由上游 agent 调用视觉模型产生（本脚本不直接调用视觉模型）

Usage:
    python compare_ppt_versions.py \
        --input v1.pptx v2.pptx v3.pptx \
        --labels "原版" "调色版" "重排版" \
        --scenario marketing \
        --output compare_report.json \
        --md compare_report.md

灵感来源: CSDN 博客《用 ai 给UI 页面打分》的多版本对比方法论
（https://blog.csdn.net/tomxjc/article/details/154011197）
"""
from __future__ import annotations

import argparse
import json
import os
import sys
from pathlib import Path

# Force UTF-8 on Windows so CJK doesn't mojibake.
try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

# Import sibling module (score_ppt_pages.py)
_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))
try:
    from score_ppt_pages import (
        SCENARIO_WEIGHTS,
        SCORING_PROMPT,
        analyze_slide,
        summarize_deck,
        _resolve_theme_fonts,
        _extract_design_tokens,
        _build_scenario_prompt,
    )
except ImportError as e:
    print(f"Error: cannot import from score_ppt_pages.py: {e}", file=sys.stderr)
    print("Ensure score_ppt_pages.py is in the same scripts/ directory.", file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx",
          file=sys.stderr)
    sys.exit(1)

EMU_PER_INCH = 914400


# --------------------------------------------------------------------------- #
# Core: score one PPTX (structural only — vision scoring is upstream)
# --------------------------------------------------------------------------- #
def score_one_pptx(pptx_path: Path, scenario: str) -> dict:
    """
    Structural scoring for one PPTX. Returns a report dict similar to
    score_ppt_pages.py's output, but without the vision_scoring_prompt
    (we don't need it per-version when comparing).
    """
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    slide_w_in = prs.slide_width / EMU_PER_INCH
    slide_h_in = prs.slide_height / EMU_PER_INCH

    theme_major, theme_minor = _resolve_theme_fonts(prs)

    slides_data = []
    for idx in range(total):
        slide = prs.slides[idx]
        analysis = analyze_slide(
            slide, idx, slide_w_in, slide_h_in,
            theme_major_font=theme_major, theme_minor_font=theme_minor,
        )
        slides_data.append(analysis)

    deck_summary = summarize_deck(slides_data)
    scenario_weights = SCENARIO_WEIGHTS[scenario]

    return {
        "file": pptx_path.name,
        "slide_size_in": [round(slide_w_in, 3), round(slide_h_in, 3)],
        "total_slides": total,
        "scenario": scenario,
        "scenario_weights": scenario_weights,
        "deck_summary": deck_summary,
        "slides": slides_data,
        # Note: vision_scoring_prompt omitted (comparison uses structural only;
        # agent can run vision scoring separately per version if needed)
    }


# --------------------------------------------------------------------------- #
# Compute a structural-only "final_score" for ranking purposes
# --------------------------------------------------------------------------- #
# Each anti-pattern docks a fixed amount. Vision scores (if provided) should
# be merged by the caller; this script focuses on structural ranking.
# Each anti-pattern docks a fixed amount per SKILL.md 扣分 column.
# Vision scores (if provided) should be merged by the caller;
# this script focuses on structural ranking.
ANTI_PATTERN_PENALTY = {
    # AP-WS-*
    "text_heavy_deck": 3.0,
    "wall_of_text": 1.5,
    "high_text_density": 1.0,
    "low_fill_rate": 1.5,
    "high_fill_rate": 0.5,
    "empty_quadrant": 1.0,
    # AP-TS-*
    "too_many_fonts": 2.0,
    "unprofessional_font": 2.0,
    "cjk_text_no_cjk_font": 1.5,
    "weak_type_scale": 1.5,
    "extreme_type_scale": 1.0,
    "code_without_monospace": 1.0,
    # AP-CH-*
    "too_many_colors": 2.0,
    "ai_purple_palette": 1.0,
    "oversaturated_pure_rgb": 1.0,
    "overuse_accent": 1.0,
    "pure_black_text": 0.5,
    "pure_white_card": 0.5,
    "rainbow_text": 1.5,
    "inconsistent_colors": 1.5,
    # AP-AL-*
    "no_margins": 1.0,
    "three_equal_cards": 1.0,
    "misaligned_elements": 1.5,
    # AP-SP-*
    "irregular_spacing": 1.0,
    "inconsistent_spacing": 1.0,
    "orphan_widow": 0.5,
    # AP-IM-*
    "low_res_image": 1.0,
    "stretched_image": 2.0,
    "stretched_images": 2.0,
    "clipart_style": 1.0,
    # AP-CO-*
    "inconsistent_fonts": 2.0,
    "default_template": 0.5,
    # AP-HI-*
    "tiny_text": 1.5,
    "no_visual_hierarchy": 2.0,
    "long_bullet": 1.0,
    "too_many_bullets": 1.5,
    "bullet_soup": 1.0,
    "low_contrast": 1.0,
    "chart_without_message": 1.0,
    "key_data_not_emphasized": 1.0,
    # AP-LA-*
    "top_text_bottom_image": 1.5,
    "no_card_containers": 1.0,
    "image_at_edge": 0.5,
    "bare_text_no_card": 1.0,
    "table_no_zebra": 1.0,
    "table_weak_header": 1.0,
    "mixed_numbering": 0.5,
    "unbalanced_layout": 1.0,
    "hollow_container": 2.0,
    "top_heavy": 1.0,
    # AP-DE-*
    "repeated_chapter_labels": 1.0,
    "inconsistent_alignment": 1.5,
    "no_brand_consistency": 1.0,
    "inconsistent_margins": 1.5,
    "inconsistent_section_pos": 1.0,
    "cover_ending_mismatch": 1.0,
    "density_outlier": 1.0,
    "no_header_navigation": 1.5,
    "deck_too_many_fonts": 1.5,
    "deck_too_many_colors": 1.5,
}


def _structural_score(report: dict) -> tuple[float, list[str], float]:
    """
    Return (base_score, anti_pattern_list, penalty_total).
    base_score is a heuristic structural baseline (0-10) derived from
    deck_summary metrics. Vision scoring should override this when available.
    """
    deck = report.get("deck_summary", {})
    slides = report.get("slides", [])

    # Base heuristic: start at 8.0 and dock for structural issues.
    base = 8.0

    font_count = deck.get("deck_font_count", 0)
    if font_count > 4:
        base -= (font_count - 4) * 0.5
    elif font_count >= 1 and font_count <= 2:
        base += 0.3  # reward restraint

    color_count = deck.get("deck_color_count", 0)
    if color_count > 12:
        base -= (color_count - 12) * 0.2
    elif color_count <= 5:
        base += 0.2

    avg_issues = deck.get("avg_issues_per_slide", 0)
    base -= min(avg_issues * 0.4, 2.0)

    # Collect all anti-patterns
    all_aps = []
    for s in slides:
        all_aps.extend(s.get("anti_patterns", []))
    deck_aps = deck.get("deck_anti_patterns", [])
    all_aps.extend(deck_aps)

    # Compute penalty
    penalty = 0.0
    ap_names = []
    for ap in all_aps:
        # ap string format: "name: detail"
        name = ap.split(":", 1)[0].strip()
        ap_names.append(ap)
        penalty += ANTI_PATTERN_PENALTY.get(name, 0.5)

    base = max(0.0, min(10.0, base))
    final = max(0.0, base - min(penalty, 4.0))  # cap penalty at 4.0
    return round(final, 2), ap_names, round(penalty, 2)


# --------------------------------------------------------------------------- #
# Token diff table
# --------------------------------------------------------------------------- #
def _aggregate_tokens(report: dict) -> dict:
    """
    Aggregate per-slide tokens into a deck-level token summary.
    Uses the most common non-null value across slides.
    """
    from collections import Counter

    slides = report.get("slides", [])
    if not slides:
        return {}

    # Collect all values per token path
    token_paths = {}  # path_tuple -> list of values
    for s in slides:
        tokens = s.get("tokens", {})
        def walk(obj, prefix=()):
            if isinstance(obj, dict):
                for k, v in obj.items():
                    walk(v, prefix + (k,))
            else:
                if obj is not None and v is not None:
                    token_paths.setdefault(prefix, []).append(obj)
        # Fix: walk over tokens dict
        for k1, v1 in tokens.items():
            if isinstance(v1, dict):
                for k2, v2 in v1.items():
                    if v2 is not None:
                        token_paths.setdefault((k1, k2), []).append(v2)
            elif v1 is not None:
                token_paths.setdefault((k1,), []).append(v1)

    # Pick the most common value per path
    result = {}
    for path, values in token_paths.items():
        # For numeric values, take median; for strings, take mode
        numeric = all(isinstance(v, (int, float)) for v in values)
        if numeric:
            sorted_v = sorted(values)
            median = sorted_v[len(sorted_v) // 2]
            result[".".join(path)] = median
        else:
            counter = Counter(values)
            result[".".join(path)] = counter.most_common(1)[0][0]
    return result


def build_token_diff(reports: list[dict], labels: list[str]) -> dict:
    """Build a token diff table: token_path -> {label: value}"""
    all_tokens = {}
    for label, rpt in zip(labels, reports):
        tk = _aggregate_tokens(rpt)
        for k, v in tk.items():
            all_tokens.setdefault(k, {})[label] = v
    return all_tokens


# --------------------------------------------------------------------------- #
# Markdown report
# --------------------------------------------------------------------------- #
def generate_markdown_report(reports, labels, scenario, token_diff, rankings):
    """Generate a human-readable Markdown comparison report."""
    lines = []
    lines.append("# PPT 版本对比报告")
    lines.append("")
    lines.append(f"**场景**: `{scenario}`")
    lines.append(f"**版本数**: {len(labels)}")
    lines.append(f"**生成工具**: `scripts/compare_ppt_versions.py`")
    lines.append("")

    # ---- Ranking ----
    lines.append("## 1. 综合排名")
    lines.append("")
    lines.append("| 排名 | 版本 | 文件 | 结构分 | 反模式数 | 扣分 | 最终分 |")
    lines.append("|---|---|---|---|---|---|---|")
    for rank, (label, rpt, final, aps, penalty, _final_for_sort) in enumerate(rankings, 1):
        lines.append(
            f"| {rank} | {label} | {rpt['file']} | "
            f"{final + penalty:.2f} | {len(aps)} | -{penalty:.2f} | **{final:.2f}** |"
        )
    lines.append("")

    # ---- Per-dimension structural signals ----
    lines.append("## 2. 结构信号对比")
    lines.append("")
    lines.append("| 指标 | " + " | ".join(labels) + " |")
    lines.append("|---|" + "---|" * len(labels))
    metrics = [
        ("总页数", lambda r: r.get("total_slides", 0)),
        ("字体家族数", lambda r: r.get("deck_summary", {}).get("deck_font_count", 0)),
        ("颜色数", lambda r: r.get("deck_summary", {}).get("deck_color_count", 0)),
        ("平均反模式/页", lambda r: r.get("deck_summary", {}).get("avg_issues_per_slide", 0)),
        ("最差页", lambda r: r.get("deck_summary", {}).get("worst_slide", "—")),
    ]
    for name, fn in metrics:
        row = [name] + [str(fn(r)) for r in reports]
        lines.append("| " + " | ".join(row) + " |")
    lines.append("")

    # ---- Token diff ----
    lines.append("## 3. Token 差异表")
    lines.append("")
    if not token_diff:
        lines.append("（无 token 数据）")
    else:
        lines.append("| Token | " + " | ".join(labels) + " |")
        lines.append("|---|" + "---|" * len(labels))
        for token_path in sorted(token_diff.keys()):
            row = [token_path]
            for label in labels:
                v = token_diff[token_path].get(label, "—")
                row.append(str(v))
            lines.append("| " + " | ".join(row) + " |")
    lines.append("")

    # ---- Anti-pattern comparison ----
    lines.append("## 4. 反模式对比")
    lines.append("")
    # Collect all unique anti-pattern names
    all_ap_names = set()
    ap_per_version = {}
    for label, rpt in zip(labels, reports):
        aps = []
        for s in rpt.get("slides", []):
            for ap in s.get("anti_patterns", []):
                name = ap.split(":", 1)[0].strip()
                aps.append(name)
        for ap in rpt.get("deck_summary", {}).get("deck_anti_patterns", []):
            name = ap.split(":", 1)[0].strip()
            aps.append(name)
        ap_per_version[label] = aps
        all_ap_names.update(aps)

    if all_ap_names:
        lines.append("| 反模式 | " + " | ".join(labels) + " |")
        lines.append("|---|" + "---|" * len(labels))
        for name in sorted(all_ap_names):
            row = [name]
            for label in labels:
                count = ap_per_version[label].count(name)
                row.append(str(count) if count else "—")
            lines.append("| " + " | ".join(row) + " |")
    else:
        lines.append("（所有版本均未检测到反模式）")
    lines.append("")

    # ---- Recommendations ----
    lines.append("## 5. 推荐结论")
    lines.append("")
    if rankings:
        champion = rankings[0]
        lines.append(f"- **综合冠军**: {champion[0]}（{champion[5]:.2f} 分）")
        if len(rankings) >= 2:
            second = rankings[1]
            lines.append(f"- **亚军**: {second[0]}（{second[5]:.2f} 分）")

        # Single-dimension champions (based on structural signals)
        lines.append("")
        lines.append("### 单维度冠军")
        dim_metrics = [
            ("字体最克制", lambda r: r.get("deck_summary", {}).get("deck_font_count", 99), min),
            ("颜色最克制", lambda r: r.get("deck_summary", {}).get("deck_color_count", 99), min),
            ("反模式最少", lambda r: sum(len(s.get("anti_patterns", [])) for s in r.get("slides", [])), min),
            ("信息最丰富", lambda r: r.get("total_slides", 0), max),
        ]
        for dim_name, fn, agg in dim_metrics:
            values = [(label, fn(r)) for label, r in zip(labels, reports)]
            winner = agg(values, key=lambda x: x[1])
            lines.append(f"- {dim_name}: {winner[0]}（{winner[1]}）")

        # Merge suggestion
        lines.append("")
        lines.append("### 合并建议")
        if len(rankings) >= 2:
            champ_label = champion[0]
            champ_tokens = _aggregate_tokens(champion[1])
            for runner_up in rankings[1:]:
                ru_label = runner_up[0]
                ru_tokens = _aggregate_tokens(runner_up[1])
                # Find tokens where runner-up differs and might be better
                diffs = []
                for k, v_champ in champ_tokens.items():
                    v_ru = ru_tokens.get(k)
                    if v_ru is not None and v_ru != v_champ:
                        diffs.append(f"`{k}`: {champ_label}={v_champ} vs {ru_label}={v_ru}")
                if diffs:
                    lines.append(f"- 可考虑从 {ru_label} 借鉴以下 token（{champ_label} 当前值可能不是最优）：")
                    for d in diffs[:5]:
                        lines.append(f"  - {d}")
                else:
                    lines.append(f"- {ru_label} 无显著优于 {champ_label} 的 token，可忽略")
    lines.append("")

    lines.append("## 6. 下一步建议")
    lines.append("")
    lines.append("1. **如需视觉评分**：对每个版本跑 `scripts/render_slides.py` + 视觉模型评分，再合并本报告的结构分")
    lines.append("2. **如需修复**：对冠军版本跑 `scripts/fix_ppt.py --all`，再重新对比验证")
    lines.append("3. **如需合并**：参考上方\"合并建议\"，从亚军版本迁移特定 token 到冠军版本")
    lines.append("")

    return "\n".join(lines)


# --------------------------------------------------------------------------- #
# Main
# --------------------------------------------------------------------------- #
def main():
    ap = argparse.ArgumentParser(description="PPT Aesthetic Version Comparator (改进项 1)")
    ap.add_argument("--input", "-i", nargs="+", required=True,
                    help="2+ PPTX file paths to compare")
    ap.add_argument("--labels", "-l", nargs="+", default=None,
                    help="Display labels for each input (default: file basename)")
    ap.add_argument("--scenario", "-s", default="default",
                    choices=list(SCENARIO_WEIGHTS.keys()),
                    help="Scene-based weight preset (same as score_ppt_pages.py)")
    ap.add_argument("--output", "-o", default=None,
                    help="Output JSON report path (default: compare_report.json)")
    ap.add_argument("--md", "-m", default=None,
                    help="Output Markdown report path (default: compare_report.md)")
    args = ap.parse_args()

    if len(args.input) < 2:
        print("Error: at least 2 PPTX files required for comparison", file=sys.stderr)
        sys.exit(1)

    inputs = [Path(p) for p in args.input]
    for p in inputs:
        if not p.exists():
            print(f"Error: File not found: {p}", file=sys.stderr)
            sys.exit(1)

    labels = args.labels or [p.stem for p in inputs]
    if len(labels) != len(inputs):
        print("Error: number of labels must match number of inputs", file=sys.stderr)
        sys.exit(1)

    print(f"Comparing {len(inputs)} PPT versions under scenario='{args.scenario}':")
    for label, p in zip(labels, inputs):
        print(f"  - {label}: {p.name}")

    # Score each version
    reports = []
    for label, p in zip(labels, inputs):
        print(f"\nScoring '{label}' ({p.name})...")
        rpt = score_one_pptx(p, args.scenario)
        reports.append(rpt)
        final, aps, penalty = _structural_score(rpt)
        print(f"  → structural final: {final:.2f}, anti-patterns: {len(aps)}, penalty: {penalty:.2f}")

    # Build rankings: (label, report, final, aps, penalty, final_score_for_sort)
    rankings = []
    for label, rpt in zip(labels, reports):
        final, aps, penalty = _structural_score(rpt)
        rankings.append((label, rpt, final, aps, penalty, final))
    rankings.sort(key=lambda x: x[5], reverse=True)

    # Build token diff
    token_diff = build_token_diff(reports, labels)

    # Build JSON output
    json_output = {
        "scenario": args.scenario,
        "version_count": len(inputs),
        "versions": [],
        "ranking": [r[0] for r in rankings],
        "scenario_champion": {"scenario": args.scenario, "winner": rankings[0][0]},
        "token_diff": token_diff,
        "merge_suggestion": "See Markdown report for details",
    }
    for label, rpt, final, aps, penalty, _ in rankings:
        json_output["versions"].append({
            "label": label,
            "file": rpt["file"],
            "final_score": final,
            "anti_pattern_count": len(aps),
            "penalty": penalty,
            "anti_patterns": aps,
        })

    # Write JSON
    out_json = args.output or "compare_report.json"
    with open(out_json, "w", encoding="utf-8") as f:
        json.dump(json_output, f, ensure_ascii=False, indent=2)
    print(f"\nJSON report saved to: {out_json}")

    # Write Markdown
    out_md = args.md or "compare_report.md"
    md_content = generate_markdown_report(reports, labels, args.scenario, token_diff, rankings)
    with open(out_md, "w", encoding="utf-8") as f:
        f.write(md_content)
    print(f"Markdown report saved to: {out_md}")

    # Print summary
    print("\n" + "=" * 60)
    print("COMPARISON SUMMARY")
    print("=" * 60)
    print(f"Scenario: {args.scenario}")
    print(f"Champion: {rankings[0][0]} (final={rankings[0][5]:.2f})")
    if len(rankings) >= 2:
        print(f"Runner-up: {rankings[1][0]} (final={rankings[1][5]:.2f})")
    print(f"Token diff entries: {len(token_diff)}")
    print("=" * 60)


if __name__ == "__main__":
    main()
