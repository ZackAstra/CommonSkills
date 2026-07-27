# -*- coding: utf-8 -*-
"""embed_illustrations.py — JSON 配置驱动的配图嵌入脚本

职责：读取 JSON 配置表，将指定目录中的 PNG 图片嵌入 PPTX 对应页面。
      自动计算图片位置和尺寸，保持原始宽高比。

使用方式：
  python embed_illustrations.py input.pptx -o output.pptx --config embed_config.json --images dir/

embed_config.json 格式：
  [
    {
      "slide": 2,
      "image": "02_timeline.png",        // 文件名（在 --images 目录中查找）
      "left": 7.5,                        // 英寸，左侧位置
      "top": 2.0,                         // 英寸，顶部位置
      "width": 5.3,                       // 英寸，宽度（高度自动按比例计算）
      // "height": 4.5,                   // 可选，显式指定高度（忽略宽高比）
    },
    {
      "slide": 6,
      "image": "06_staircase.png",
      "left": 8.2,
      "top": 1.2,
      "width": 4.5,
    },
    ...
  ]

位置预设（简化配置）：
  若不想精确指定 left/top/width，可用 "position" 字段：
    "position": "right"    → 自动占右 1/3: left=8.5, top=1.2, width=4.3
    "position": "center"   → 自动居中: left=3.5, top=1.5, width=6.3
    "position": "full"     → 全宽: left=0.8, top=1.0, width=11.7
    "position": "left"     → 自动占左 1/3: left=0.8, top=1.2, width=4.3
"""

import argparse
import json
import os
import sys
from pathlib import Path

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import PP_PLACEHOLDER
except ImportError:
    print("ERROR: python-pptx is required.  pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ── Position presets ──
# Based on 13.333" x 7.5" widescreen slides

POSITION_PRESETS = {
    "right":  {"left": 8.5, "top": 1.2, "width": 4.3},
    "left":   {"left": 0.8, "top": 1.2, "width": 4.3},
    "center": {"left": 3.5, "top": 1.5, "width": 6.3},
    "full":   {"left": 0.8, "top": 1.0, "width": 11.7},
}


def resolve_image_size(image_path, width=None, height=None):
    """Calculate actual placement dimensions preserving aspect ratio.
    
    Priority:
      - Both width and height specified → use as-is (may stretch)
      - Only width → height from aspect ratio
      - Only height → width from aspect ratio
      - Neither → default width=4.5, height from aspect ratio
    """
    if PILImage is not None:
        with PILImage.open(image_path) as im:
            iw, ih = im.size
        aspect = iw / ih
    else:
        # Fallback: assume 16:9 aspect
        aspect = 16 / 9

    if width and height:
        return float(width), float(height)
    elif width:
        return float(width), float(width) / aspect
    elif height:
        return float(height) / aspect, float(height)
    else:
        w = 4.5
        return w, w / aspect


def embed_image(prs, slide_index, image_path, left, top, width, height):
    """Embed an image into a specific slide of the presentation.
    
    slide_index: 0-based slide index
    Returns True on success, False on failure.
    """
    if slide_index >= len(prs.slides):
        print(f"  WARNING: slide index {slide_index} out of range "
              f"(total {len(prs.slides)} slides), skipping", file=sys.stderr)
        return False

    if not os.path.exists(image_path):
        print(f"  WARNING: image not found: {image_path}, skipping", file=sys.stderr)
        return False

    slide = prs.slides[slide_index]
    w_inch, h_inch = resolve_image_size(image_path, width, height)

    slide.shapes.add_picture(
        image_path,
        Inches(left), Inches(top),
        Inches(w_inch), Inches(h_inch),
    )
    return True


def clean_placeholders(prs, verbose=False):
    """Remove empty/unfilled placeholders from all slides.
    
    This prevents PowerPoint from rendering default prompt text like
    "单击此处添加标题" when a template has title placeholders that
    were never filled by the build script (because it used add_textbox).
    
    A placeholder is considered "empty" if:
      - It has no text (or only whitespace)
      - It has no image fill
    
    Returns the number of placeholders removed.
    """
    removed = 0
    for slide_idx, slide in enumerate(prs.slides):
        sp_tree = slide.shapes._spTree
        to_remove = []
        
        for shape in slide.shapes:
            # Only process placeholders
            if not shape.is_placeholder:
                continue
            
            # Check if placeholder has any real content
            has_text = False
            has_image = False
            
            # Check text content
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    has_text = True
            
            # Check for image fill
            try:
                if shape.shape_type is not None and hasattr(shape, 'image'):
                    has_image = True
            except Exception:
                pass
            
            # Empty placeholder — mark for removal
            if not has_text and not has_image:
                to_remove.append(shape)
        
        for shape in to_remove:
            sp = shape._element
            sp_tree.remove(sp)
            removed += 1
            if verbose:
                ph_type = "?"
                try:
                    ph_type = shape.placeholder_format.type
                except Exception:
                    pass
                print(f"    Removed empty placeholder on slide {slide_idx + 1}: "
                      f"type={ph_type}")
    
    return removed


def load_config(config_path):
    """Load and validate embed config JSON."""
    with open(config_path, "r", encoding="utf-8") as f:
        config = json.load(f)
    
    if not isinstance(config, list):
        print("ERROR: config must be a JSON array of objects", file=sys.stderr)
        sys.exit(1)
    
    for i, entry in enumerate(config):
        if "slide" not in entry:
            print(f"ERROR: config entry {i} missing 'slide' field", file=sys.stderr)
            sys.exit(1)
        if "image" not in entry:
            print(f"ERROR: config entry {i} missing 'image' field", file=sys.stderr)
            sys.exit(1)
    
    return config


def resolve_position(entry):
    """Resolve position from entry, expanding presets if needed."""
    pos = entry.get("position")
    if pos:
        if pos not in POSITION_PRESETS:
            print(f"WARNING: unknown position preset '{pos}', "
                  f"available: {list(POSITION_PRESETS.keys())}", file=sys.stderr)
            pos = "right"
        preset = POSITION_PRESETS[pos]
        # Explicit fields override preset
        left = entry.get("left", preset["left"])
        top = entry.get("top", preset["top"])
        width = entry.get("width", preset["width"])
        height = entry.get("height")
    else:
        left = entry.get("left", 8.5)
        top = entry.get("top", 1.2)
        width = entry.get("width")
        height = entry.get("height")
    
    return left, top, width, height


def main():
    parser = argparse.ArgumentParser(
        description="Embed illustrations into PPTX based on JSON config."
    )
    parser.add_argument("input", help="Input PPTX file")
    parser.add_argument("-o", "--output", help="Output PPTX file (default: input_embedded.pptx)")
    parser.add_argument("--config", required=True,
                        help="JSON config specifying which images go on which slides")
    parser.add_argument("--images", required=True,
                        help="Directory containing illustration PNG files")
    parser.add_argument("--dry-run", action="store_true",
                        help="Show what would be done without modifying PPTX")
    parser.add_argument("--clean-placeholders", action="store_true",
                        help="Remove empty/unfilled placeholders (fixes '单击此处添加标题' ghost text)")
    args = parser.parse_args()

    input_path = args.input
    output_path = args.output or (str(Path(input_path).with_suffix("")) + "_embedded.pptx")
    images_dir = Path(args.images)

    # Load PPTX
    prs = Presentation(input_path)
    print(f"Loaded: {input_path} ({len(prs.slides)} slides)")

    # Clean empty placeholders before embedding (fixes ghost prompt text)
    if args.clean_placeholders and not args.dry_run:
        removed = clean_placeholders(prs, verbose=True)
        if removed:
            print(f"Cleaned {removed} empty placeholder(s)")
        else:
            print("No empty placeholders found")

    # Load config
    config = load_config(args.config)
    print(f"Config: {len(config)} illustration(s) to embed")

    # Process each entry
    success_count = 0
    for entry in config:
        slide_num = entry["slide"]  # 1-based in config
        slide_index = slide_num - 1  # 0-based for python-pptx
        image_name = entry["image"]
        image_path = images_dir / image_name

        left, top, width, height = resolve_position(entry)

        print(f"  Slide {slide_num}: {image_name} at ({left:.1f}\", {top:.1f}\") "
              f"width={width or 'auto'}\" height={height or 'auto'}\"")

        if not args.dry_run:
            ok = embed_image(prs, slide_index, str(image_path), left, top, width, height)
            if ok:
                success_count += 1
        else:
            if image_path.exists():
                success_count += 1
            else:
                print(f"    [DRY-RUN] Image not found: {image_path}", file=sys.stderr)

    # Save
    if not args.dry_run:
        prs.save(output_path)
        size_kb = os.path.getsize(output_path) / 1024
        print(f"\nSaved: {output_path} ({size_kb:.0f} KB, "
              f"{success_count}/{len(config)} images embedded)")
    else:
        print(f"\n[DRY-RUN] Would embed {success_count}/{len(config)} images")

    # Exit code
    if success_count < len(config):
        sys.exit(1)


if __name__ == "__main__":
    main()
