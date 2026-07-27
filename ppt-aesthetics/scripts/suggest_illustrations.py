# -*- coding: utf-8 -*-
"""suggest_illustrations.py — 内容感知配图建议脚本

职责：分析 PPTX 每页的文本内容，自动推断内容类型，建议合适的配图类型、
     生成方式和嵌入位置。输出 JSON 可直接用于 gen_illustrations.py 和
     embed_illustrations.py 的配置输入。

使用方式：
  python suggest_illustrations.py input.pptx [-o suggestions.json] [--verbose]

输出 JSON 格式：
  {
    "source": "input.pptx",
    "total_slides": 18,
    "image_coverage_needed": 7,
    "suggestions": [
      {
        "slide": 1,
        "title": "封面标题",
        "body_text": "...",
        "content_type": "cover",
        "illustration_type": "scene_image",
        "generation_method": "ImageGen",
        "position_preset": "full",
        "priority": "must",
        "prompt_hint": "主题关键词 + 专业风格 + 宽幅构图 16:9",
        "description": "封面场景大图"
      },
      ...
    ],
    "gen_config": [...],    // 可直接用于 gen_illustrations.py --config
    "embed_config": [...],  // 可直接用于 embed_illustrations.py --config
    "summary": {
      "by_method": {"ImageGen": 3, "SVG": 4, "python-pptx": 1},
      "by_type": {"flowchart": 2, "timeline": 1, ...},
      "coverage_rate": 0.39
    }
  }

内容类型检测规则：
  - cover/ending     → 标题含"封"/"目录"/"结束"/"谢谢"/"Q&A"，或首末页
  - timeline         → 含"演进"/"阶段"/"发展"/"时间线"/"历程"/"趋势"
  - process/flow     → 含"流程"/"步骤"/"闭环"/"方法"/"工作流"/"pipeline"
  - comparison       → 含"对比"/"区别"/"vs"/"比较"/"差异"/"优劣"
  - hierarchy        → 含"架构"/"层次"/"结构"/"体系"/"组织"/"分层"
  - data/metrics     → 含"数据"/"指标"/"百分比"/"增长"/"统计"/"报表"
  - list/enumeration → 文本含 3+ 编号项或项目符号
  - concept          → 含"概念"/"定义"/"是什么"/"理解"/"原理"/"核心"
  - default          → 未匹配上述规则的页面
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    print("ERROR: python-pptx is required.  pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ── Content type detection rules ──

CONTENT_RULES = [
    # (content_type, keywords_regex, illustration_type, generation_method)
    # 注意：规则按检测优先级排列。同一条关键词可能出现在多种类型中，
    # 消歧靠 detect_content_type() 中的标题加权打分机制，而非规则顺序。
    (
        "timeline",
        r"演进|阶段|发展|时间线|历程|趋势|演变|迭代|升级|三阶段|四阶段|5G|历[史程]",
        "timeline",
        "PIL",
    ),
    (
        "process",
        r"步骤|闭环|pipeline|回路|循环|运转|执行",
        "flowchart",
        "SVG",
    ),
    (
        "comparison",
        r"对比|区别|vs|比较|差异|优劣|对照|权衡|取舍|转向|能力对[比比]|从.*到",
        "comparison_diagram",
        "PIL",
    ),
    (
        "hierarchy",
        r"架构|层次|结构|体系|组织|分层|模块|组件|分解|分类|目录结构|要素|组成",
        "structure_diagram",
        "PIL",
    ),
    (
        "data",
        r"数据|指标|百分比|增长|统计|报表|效果|提升|降低|占比|覆盖",
        "data_chart",
        "python-pptx",
    ),
    (
        "concept",
        r"概念|定义|是什么|理解|原理|核心|本质|理论|基础|入门|突破",
        "concept_illustration",
        "ImageGen",
    ),
    # 下方两条是"弱触发"规则，仅在其他规则未匹配时才生效
    # "流程"和"工作流"单独出现时可能是正文中的附带提及，
    # 不应仅因为正文含"流程"就将整个页面判定为 process 类型
    (
        "process_weak",
        r"^(?=.*流程)(?!.*步骤)(?!.*闭环).*$",
        "flowchart",
        "SVG",
    ),
    (
        "concept_weak",
        r"^(?=.*工作流)(?!.*步骤)(?!.*闭环).*$",
        "concept_illustration",
        "ImageGen",
    ),
]

PRIORITY_RULES = {
    "cover": "must",
    "ending": "skip",
    "timeline": "high",
    "process": "high",
    "comparison": "high",
    "hierarchy": "high",
    "data": "medium",
    "concept": "medium",
    "list": "low",
    "default": "low",
}

# Position heuristics based on text volume
# slide_width=13.333in, slide_height=7.5in (standard widescreen)
POSITION_PRESETS = {
    "cover": {"position": "full"},
    "ending": {"position": "center"},
    "timeline": {"position": "full"},
    "process": {"position": "right"},
    "comparison": {"position": "full"},
    "hierarchy": {"position": "right"},
    "data": {"position": "full"},
    "concept": {"position": "right"},
    "list": {"position": "right"},
    "default": {"position": "right"},
}

# Position preset coordinate mappings (inches)
# Based on 13.333 x 7.5 slide, margin 0.8in
POSITION_COORDS = {
    "right":  {"left": 7.5,  "top": 1.2, "width": 5.3},
    "left":   {"left": 0.8,  "top": 1.2, "width": 5.3},
    "center": {"left": 3.5,  "top": 1.5, "width": 6.3},
    "full":   {"left": 0.8,  "top": 1.0, "width": 11.7},
}

ILLUSTRATION_TYPE_NAMES = {
    "timeline": "时间线示意图",
    "flowchart": "流程图",
    "comparison_diagram": "对比图",
    "structure_diagram": "结构图",
    "data_chart": "数据可视化图表",
    "concept_illustration": "概念示意图",
    "icon_cards": "图标卡片组",
    "scene_image": "场景大图/背景图",
    "closing_image": "收尾呼应图",
}


# ── Slide text extraction ──

def extract_slide_text(slide):
    """Extract title and body text from a slide.
    
    Returns: (title, body_text) where title is str and body_text is str.
    """
    title_parts = []
    body_parts = []
    
    for shape in slide.shapes:
        # Read text from tables (card layouts use tables for content)
        if shape.has_table:
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for cell in row.cells:
                    cell_text = cell.text_frame.text.strip()
                    if not cell_text:
                        continue
                    # First row is often a header/title row
                    if row_idx == 0:
                        title_parts.append(cell_text)
                    else:
                        body_parts.append(cell_text)
            continue
        
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        
        # Heuristic: if shape is named "Title*" or is the first large text, it's a title
        shape_name = shape.name.lower()
        if "title" in shape_name or "heading" in shape_name:
            title_parts.append(text)
        else:
            body_parts.append(text)
    
    # Fallback: if no title found, use first non-empty text frame
    if not title_parts and body_parts:
        title_parts.append(body_parts.pop(0))
    
    title = " ".join(title_parts).strip()
    body = "\n".join(body_parts).strip()
    return title, body


def has_images(slide):
    """Check if a slide already has significant images (non-decorative)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    img_count = 0
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_count += 1
        except (ValueError, AttributeError):
            pass
    return img_count


def count_bullet_items(body_text):
    """Count numbered/bulleted items in body text."""
    # Match patterns like 1. 2. 3. or • or - or ● or ○
    bullet_pattern = r'(?:^\s*(?:\d+[.、)]|[•\-●○◆◇]))'
    matches = re.findall(bullet_pattern, body_text, re.MULTILINE)
    return len(matches)


# ── Content type detection ──

def detect_content_type(slide_idx, total_slides, title, body_text):
    """Detect content type based on slide position and text content.
    
    Uses a title-weighted scoring approach:
    - Title keyword matches get 3x weight (title defines the page's primary topic)
    - Body keyword matches get 1x weight
    - Weak rules (process_weak, concept_weak) get 0.5x weight on body
    
    Returns: (content_type, confidence)
    """
    combined_text = f"{title} {body_text}".lower()
    title_lower = title.lower()
    body_lower = body_text.lower()
    
    # Cover page detection (first slide or explicit cover markers)
    if slide_idx == 0:
        cover_markers = ["目录", "contents", "agenda", "提纲", "概览"]
        if any(m in combined_text for m in cover_markers):
            return "cover", 0.9
        # First slide is usually cover
        return "cover", 0.7
    
    # Ending page detection (last slide or explicit markers)
    if slide_idx == total_slides - 1:
        ending_markers = ["谢谢", "感谢", "q&a", "问答", "结束", "联?系", "邮箱"]
        if any(m in combined_text for m in ending_markers):
            return "ending", 0.9
        return "ending", 0.5
    
    # Explicit ending markers on non-last slides
    ending_markers = ["谢谢", "感谢聆听", "q&a", "问答环节"]
    if any(m in combined_text for m in ending_markers):
        return "ending", 0.8
    
    # Section divider detection
    section_markers = ["目录", "contents", "章节", "第.*章", "part\\s+\\d", "section"]
    if any(re.search(m, combined_text) for m in section_markers):
        if len(body_text) < 50:  # Short text = likely a section divider
            return "cover", 0.7  # Treat like cover for illustration purposes
    
    # ── Title disambiguation (overrides scoring) ──
    # 标题中的关键词往往比正文更能说明页面主题，
    # 以下规则在打分之前处理常见的消歧场景
    TITLE_DISAMBIG = [
        # (pattern, content_type, confidence)
        (r"对比|区别|vs|比较|差异", "comparison", 0.9),      # "XX对比" → comparison
        (r"结构|目录|要素|组成|架构", "hierarchy", 0.9),      # "XX结构" → hierarchy
        (r"标准|法则|规律|规范|经验|原则|规则", "concept", 0.85),  # "XX标准" → concept
        (r"案例|实战|实践", "concept", 0.8),                   # "案例XX" → concept
    ]
    for pattern, override_type, override_conf in TITLE_DISAMBIG:
        if re.search(pattern, title_lower):
            return override_type, override_conf
    
    # ── Title-weighted scoring ──
    best_type = "default"
    best_score = 0.0
    
    for content_type, keyword_pattern, illu_type, gen_method in CONTENT_RULES:
        # Skip weak rules in first pass (they have lower weight)
        is_weak = content_type.endswith("_weak")
        base_type = content_type.replace("_weak", "")
        
        # Count keyword matches in title (3x weight)
        title_matches = len(re.findall(keyword_pattern, title_lower))
        title_score = title_matches * 3.0
        
        # Count keyword matches in body (1x weight, weak rules get 0.5x)
        body_matches = len(re.findall(keyword_pattern, body_lower))
        body_weight = 0.5 if is_weak else 1.0
        body_score = body_matches * body_weight
        
        total_score = title_score + body_score
        
        if total_score > best_score:
            best_type = base_type  # Strip "_weak" suffix
            best_score = total_score
    
    # Calculate confidence from score
    if best_score > 0:
        confidence = min(0.4 + best_score * 0.1, 0.95)
        # Boost confidence if title matched strongly
        if best_score >= 3.0:  # At least one title keyword hit
            confidence = min(confidence + 0.15, 0.95)
    else:
        confidence = 0.0
    
    # List/enumeration detection (3+ bullet items)
    if best_type == "default":
        bullet_count = count_bullet_items(body_text)
        if bullet_count >= 3:
            return "list", 0.6
    
    # If only weak rules matched, reduce confidence
    if best_score > 0 and best_score < 2.0:
        confidence = min(confidence, 0.55)
    
    return best_type, confidence


def suggest_illustration(content_type, title, body_text):
    """Suggest illustration type, generation method, and prompt hint.
    
    Returns: (illustration_type, generation_method, prompt_hint, description)
    """
    # Look up in content rules first (match base type, strip "_weak" suffix)
    for ct, _, illu_type, gen_method in CONTENT_RULES:
        if ct.replace("_weak", "") == content_type:
            name = ILLUSTRATION_TYPE_NAMES.get(illu_type, illu_type)
            prompt_hint = _build_prompt_hint(illu_type, title, body_text)
            return illu_type, gen_method, prompt_hint, name
    
    # Special cases
    if content_type == "cover":
        return ("scene_image", "ImageGen",
                f"{title}，扁平专业风格，宽幅构图，品牌色调，16:9",
                "封面场景大图/背景图")
    
    if content_type == "ending":
        return ("closing_image", "ImageGen",
                f"收束呼应，与封面风格统一，简洁意境，大留白，16:9",
                "收尾呼应图")
    
    if content_type == "list":
        return ("icon_cards", "PIL",
                f"3-5 个图标卡片，每卡片含标题+简介，卡片网格排列",
                "图标卡片组")
    
    # Default fallback
    return ("concept_illustration", "ImageGen",
            f"{title}，专业抽象风格，简化示意图，品牌色调，16:9",
            "概念示意图")


def _build_prompt_hint(illu_type, title, body_text):
    """Build a prompt hint for illustration generation."""
    hints = {
        "timeline": f"时间线/演进图，{title}，横向时间轴，3-5个阶段节点，"
                    f"箭头连接，节点含标签，扁平风格，品牌色调",
        "flowchart": f"流程图，{title}，3-5个步骤方框，箭头连接，"
                     f"简化流程，扁平风格，品牌色调",
        "comparison_diagram": f"对比图，{title}，左右或上下分栏，"
                             f"对比要素罗列，扁平风格，品牌色调",
        "structure_diagram": f"结构图/层级图，{title}，树形或嵌套结构，"
                            f"3-4层，连线标注，扁平风格，品牌色调",
        "data_chart": f"数据图表，{title}，柱状/折线/饼图，"
                     f"数据标签清晰，扁平风格，品牌色调",
        "concept_illustration": f"概念示意图，{title}，抽象图形表达核心概念，"
                               f"简化视觉隐喻，扁平风格，品牌色调",
    }
    return hints.get(illu_type, f"{title}，专业示意图，扁平风格，品牌色调，16:9")


def suggest_position(content_type, title, body_text, slide_width=13.333):
    """Suggest embedding position for an illustration.
    
    Returns: dict with left/top/width (inches) or position preset.
    """
    preset_name = POSITION_PRESETS.get(content_type, POSITION_PRESETS["default"])
    preset = preset_name.get("position", "right")
    
    body_len = len(body_text)
    title_len = len(title)
    
    # Refine position based on text volume
    if content_type in ("cover", "ending", "timeline", "comparison"):
        # These types always need full-width or center illustrations
        coord = POSITION_COORDS[preset]
    elif body_len > 200:
        # Heavy text content: use right panel, smaller illustration
        coord = {"left": 8.2, "top": 1.5, "width": 4.5}
    elif body_len > 80:
        # Medium text: right side standard
        coord = POSITION_COORDS[preset]
    elif body_len > 0:
        # Light text: illustration can be larger
        if preset == "right":
            coord = {"left": 6.7, "top": 1.2, "width": 6.0}
        else:
            coord = POSITION_COORDS[preset]
    else:
        # Title only (section divider): center or full
        coord = POSITION_COORDS.get(preset, POSITION_COORDS["center"])
    
    return {"position": preset, **coord}


# ── Illustration data extraction from page content ──

def extract_illustration_data(content_type, title, body_text):
    """从页面文本提取配图结构化数据，供 gen_pil_illustrations.py 直接使用。

    返回: data dict，与 gen_pil_illustrations.py 的 config[].data 字段格式一致。
          提取失败返回 None，由 Agent 手动补充。
    """
    extractors = {
        "timeline": _extract_timeline_data,
        "comparison": _extract_comparison_data,
        "hierarchy": _extract_hierarchy_data,
        "process": _extract_flowchart_data,
        "concept": _extract_infographic_data,
        "data": _extract_infographic_data,
        "list": _extract_list_data,
    }
    extractor = extractors.get(content_type)
    if extractor is None:
        return None
    return extractor(title, body_text)


def _clean_text(text):
    """清理文本：去首尾空白、统一空白字符。"""
    return re.sub(r'\s+', ' ', text.strip()) if text else ""


def _extract_timeline_data(title, body_text):
    """提取时间线配图数据。

    策略1：匹配"第X阶段：名称"或编号列表
    策略2：匹配"YYYY年"日期 + 附近标签词
    策略3：匹配em-dash分隔的列表项
    """
    stages = []
    combined = f"{title}\n{body_text}"

    # 策略1：编号阶段 "第一阶段：名称" 或 "1.名称" 或 "一、名称"
    stage_patterns = [
        r'第[一二三四五\d]+[阶段章部][：:]*\s*(.+?)(?:\n|$)',
        r'(\d+)[.、]\s*(.+?)(?:\n|$)',
        r'[一二三四五]、\s*(.+?)(?:\n|$)',
    ]
    # 先试编号阶段模式
    numbered_matches = re.findall(
        r'第[一二三四五\d]+[阶段章部][：:]*\s*(.+?)(?:\n|$)',
        combined
    )
    if len(numbered_matches) >= 2:
        for m in numbered_matches:
            label = _clean_text(m.split('—')[0].split('—')[0].split('：')[0])
            if label:
                stages.append({"label": label})
        # 尝试补充日期
        date_matches = re.findall(
            r'(\d{4})\s*年\s*([上下]半)?\s*年?',
            combined
        )
        for i, dm in enumerate(date_matches):
            if i < len(stages):
                sub_parts = [dm[0] + "年"]
                if dm[1]:
                    sub_parts.append(dm[1])
                stages[i]["sub"] = "".join(sub_parts)
        # 尝试补充描述（em-dash后的内容）
        desc_matches = re.findall(r'[—–]\s*(.+?)(?:\n|$)', combined)
        for i, dm in enumerate(desc_matches):
            if i < len(stages) and dm.strip():
                stages[i]["desc"] = _clean_text(dm)

    # 策略2：日期模式排列
    if not stages:
        date_pattern = r'(\d{4})\s*年\s*([上下]半)?\s*年?'
        date_matches = list(re.finditer(date_pattern, combined))
        if len(date_matches) >= 2:
            for dm in date_matches:
                sub = dm.group(1) + "年"
                if dm.group(2):
                    sub += dm.group(2)
                # 取日期前最近的短语作为label
                before_text = combined[:dm.start()]
                # 最近一行的内容
                last_line = before_text.rsplit('\n', 1)[-1].strip()
                # 去掉已有日期词
                label = re.sub(r'\d{4}\s*年\s*[上下]?\s*半?\s*年?', '', last_line).strip()
                label = re.sub(r'^[：:—–\s]+', '', label).strip()
                if not label:
                    label = sub
                stages.append({"label": label or sub, "sub": sub})

    # 策略3：em-dash列表 "名称 — 描述"
    if not stages:
        dash_items = re.findall(r'^\s*(.+?)\s*[—–]\s*(.+?)$', combined, re.MULTILINE)
        if len(dash_items) >= 2:
            for label_text, desc_text in dash_items:
                label = _clean_text(label_text)
                stages.append({"label": label, "desc": _clean_text(desc_text)})

    if len(stages) >= 2:
        return {"title": title, "stages": stages}
    return None


def _extract_comparison_data(title, body_text):
    """提取对比图配图数据。

    策略1：em-dash分隔 "类别 — 项1/项2/项3"
    策略2：冒号分隔 "类别：项1、项2"
    策略3：多段落，每段首行为类别
    """
    columns = []
    combined = f"{title}\n{body_text}"

    # 策略1：em-dash列表 "流程编排型 — Dify/Coze/伏羲"
    dash_items = re.findall(
        r'^\s*(.+?)\s*[—–]\s*(.+?)$',
        combined, re.MULTILINE
    )
    if len(dash_items) >= 2:
        for label_text, items_text in dash_items:
            label = _clean_text(label_text)
            # 拆分items：斜杠、顿号、逗号
            items = re.split(r'[/／、，,]', items_text)
            items = [_clean_text(i) for i in items if _clean_text(i)]
            if label and items:
                columns.append({"label": label, "items": items})

    # 策略2：冒号分隔 "类别：项1、项2"
    if not columns:
        colon_items = re.findall(
            r'^\s*(.+?)\s*[：:]\s*(.+?)$',
            combined, re.MULTILINE
        )
        if len(colon_items) >= 2:
            for label_text, items_text in colon_items:
                label = _clean_text(label_text)
                items = re.split(r'[、，,;/／]', items_text)
                items = [_clean_text(i) for i in items if _clean_text(i)]
                if label and items:
                    columns.append({"label": label, "items": items})

    # 策略3：表格行提取（检测 |table| 分隔符或 tab 分隔）
    if not columns:
        lines = [l.strip() for l in combined.split('\n') if l.strip()]
        # 检测分隔表头
        tab_lines = [l for l in lines if '\t' in l or '|' in l]
        if len(tab_lines) >= 2:
            for line in tab_lines:
                parts = re.split(r'\t|\|', line)
                parts = [_clean_text(p) for p in parts if _clean_text(p)]
                if len(parts) >= 2:
                    columns.append({"label": parts[0], "items": parts[1:]})

    if len(columns) >= 2:
        return {"title": title, "columns": columns}
    return None


def _extract_hierarchy_data(title, body_text):
    """提取结构图配图数据。

    策略1：目录树 ├── └──
    策略2：编号层级 1. → 1.1 → 1.1.1
    策略3：按行提取，缩进表示层级
    """
    combined = f"{title}\n{body_text}"

    # 策略1：目录树
    tree_pattern = r'^\s*([├└┌│─ └─┘┐]+)?\s*(.+?)$'
    tree_lines = []
    for line in combined.split('\n'):
        m = re.match(tree_pattern, line)
        if m:
            prefix = m.group(1) or ""
            label = _clean_text(m.group(2))
            if label and not label.startswith('#') and label != title:
                depth = len(prefix) // 2  # 粗略按字符数估计深度
                tree_lines.append({"label": label, "depth": depth})

    if len(tree_lines) >= 2:
        # 转为layers格式：根节点 + 子项
        layers = []
        # 第一项通常是根节点（如 "my-skill/"）
        root = tree_lines[0]
        layers.append({"label": root["label"]})
        for item in tree_lines[1:]:
            if item["depth"] <= root.get("depth", 0) + 1:
                layers.append({"label": item["label"]})
        if len(layers) >= 2:
            return {"title": title, "layers": layers}

    # 策略2：编号层级
    numbered = re.findall(
        r'^\s*(\d+(?:\.\d+)*)[.、）)]\s*(.+?)$',
        combined, re.MULTILINE
    )
    if len(numbered) >= 2:
        layers = []
        for num, label in numbered:
            depth = num.count('.')
            label_clean = _clean_text(label)
            if label_clean:
                layers.append({"label": label_clean, "depth": depth})
        if len(layers) >= 2:
            return {"title": title, "layers": layers}

    # 策略3：按行提取（带缩进判断）
    lines = [l for l in combined.split('\n') if l.strip()]
    if len(lines) >= 3:
        layers = []
        for line in lines[:8]:  # 最多8层
            indent = len(line) - len(line.lstrip())
            label = _clean_text(line)
            if label and label != title:
                layers.append({"label": label, "depth": indent // 2})
        if len(layers) >= 2:
            return {"title": title, "layers": layers}

    return None


def _extract_flowchart_data(title, body_text):
    """提取流程图配图数据。

    策略1：匹配 "→" 或 "->" 连接的步骤链
    策略2：匹配编号步骤 "1.步骤 2.步骤 3.步骤"
    策略3：匹配 "步骤一/步骤二" 模式
    """
    combined = f"{title}\n{body_text}"
    steps = []

    # 策略1：箭头连接 "渲染 → 评分 → 修复"
    arrow_match = re.search(
        r'(\S+(?:\s+\S+){0,3})\s*[→➡>\-]+>\s*(\S+(?:\s+\S+){0,3})\s*[→➡>\-]+>\s*(\S+(?:\s+\S+){0,3})',
        combined
    )
    if not arrow_match:
        # 更宽松：逐个 → 分割
        arrow_chain = re.split(r'\s*[→➡>\-]+>\s*', combined)
        if len(arrow_chain) >= 3:
            for step_text in arrow_chain[:5]:
                step_text = _clean_text(step_text)
                # 去掉前后非步骤文字
                step_text = re.sub(r'^.*?\n', '', step_text) if '\n' in step_text else step_text
                if step_text:
                    steps.append({"label": step_text})

    if not steps:
        # 策略1b：匹配 "名词 → 名词 → 名词" 模式（单行中）
        arrow_line = re.search(
            r'(\w[\w\s]*?)\s*→\s*(\w[\w\s]*?)\s*→\s*(\w[\w\s]*?)(?:\s*→|\s*$)',
            combined
        )
        if arrow_line:
            for g in arrow_line.groups():
                label = _clean_text(g)
                if label:
                    steps.append({"label": label})

    # 策略2：编号步骤 "1.渲染 2.评分 3.修复"
    if not steps:
        numbered = re.findall(
            r'(?:^|\n)\s*\d+[.、）)]\s*(\S+(?:\s\S+){0,2})',
            combined
        )
        if len(numbered) >= 3:
            for label in numbered[:5]:
                steps.append({"label": _clean_text(label)})

    # 策略3：关键字 "步骤" 或 "闭环" 相关步骤
    if not steps:
        # 检测 "X → Y" 的二段式
        pairs = re.findall(r'(\S+(?:\s\S+)?)\s*[→➡]\s*(\S+(?:\s\S+)?)', combined)
        if len(pairs) >= 1:
            seen = []
            for a, b in pairs:
                a, b = _clean_text(a), _clean_text(b)
                if a and a not in seen:
                    seen.append(a)
                if b and b not in seen:
                    seen.append(b)
            steps = [{"label": s} for s in seen[:5]]

    if len(steps) >= 2:
        return {"title": title, "steps": steps}
    return None


def _extract_infographic_data(title, body_text):
    """提取数据大字报配图数据。

    策略：匹配数字+单位/标签模式
    """
    combined = f"{title}\n{body_text}"
    metrics = []

    # 匹配模式："13,000+" / "60%" / "3.5万" + 后续标签
    number_pattern = r'(\d[\d,.]*\+?%?(?:万|亿)?)\s*([^\d,.\n]{2,15})'
    matches = re.findall(number_pattern, combined)
    for value, label in matches:
        label = _clean_text(label)
        # 过滤非指标标签（如"年"、"月"等太短的标签）
        if label and len(label) >= 2 and not re.match(r'^[年月份日号]$', label):
            metrics.append({"value": value, "label": label})

    if len(metrics) >= 1:
        return {"title": title, "metrics": metrics[:3]}
    return None


def _extract_list_data(title, body_text):
    """提取图标卡片组配图数据（3+ 列表项）。

    策略：提取编号/项目符号列表项
    """
    items = []
    combined = f"{title}\n{body_text}"

    # 编号列表
    numbered = re.findall(
        r'(?:^|\n)\s*(?:\d+[.、）)]|[•\-●○◆◇])\s*(.+?)$',
        combined, re.MULTILINE
    )
    for item in numbered:
        item = _clean_text(item)
        if item and len(item) >= 2:
            items.append(item)

    if len(items) >= 3:
        # 转为 columns 格式（3列均分）
        cols = [{"label": f"卡片{i+1}", "items": []} for i in range(min(3, len(items)))]
        for i, item in enumerate(items):
            col_idx = i % len(cols)
            cols[col_idx]["items"].append(item)
        return {"title": title, "columns": cols}

    return None


# ── Style suggestion ──

ACCENT_COLORS = {
    "telecom": "#E60012",
    "default": "#E60012",
    "tech": "#2563EB",
    "nature": "#059669",
}

SECONDARY_COLORS = {
    "telecom": "#2563EB",
    "default": "#2563EB",
    "tech": "#7C3AED",
    "nature": "#D97706",
}


def suggest_style(scenario="default"):
    """Suggest style parameters based on scenario/brand.

    Returns: style dict compatible with gen_pil_illustrations.py config.
    """
    accent = ACCENT_COLORS.get(scenario, ACCENT_COLORS["default"])
    secondary = SECONDARY_COLORS.get(scenario, SECONDARY_COLORS["default"])
    return {
        "accent_color": accent,
        "secondary_color": secondary,
    }


# ── Main analysis ──

def analyze_pptx(pptx_path, verbose=False, scenario="default"):
    """Analyze a PPTX file and generate illustration suggestions.
    
    Args:
        pptx_path: Path to input PPTX.
        verbose: Print detailed analysis to stderr.
        scenario: Brand/scenario name for style suggestions (e.g. "telecom").
    
    Returns: dict with suggestions, gen_config, embed_config, and summary.
    """
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    
    suggestions = []
    gen_config = []
    embed_config = []
    
    method_counts = {}
    type_counts = {}
    illustrated_count = 0
    auto_data_count = 0
    manual_data_count = 0
    
    style = suggest_style(scenario)
    
    for idx, slide in enumerate(prs.slides):
        title, body_text = extract_slide_text(slide)
        existing_images = has_images(slide)
        
        # Detect content type
        content_type, confidence = detect_content_type(idx, total_slides, title, body_text)
        
        # If slide already has images, lower priority
        if existing_images >= 2:
            priority = "skip"
        else:
            priority = PRIORITY_RULES.get(content_type, "low")
            if existing_images == 1 and priority != "must":
                priority = "low"  # Already has some visual content
        
        # Suggest illustration
        illu_type, gen_method, prompt_hint, description = suggest_illustration(
            content_type, title, body_text
        )
        
        # Suggest position
        position_info = suggest_position(content_type, title, body_text)
        
        # Extract illustration data from page content (NEW)
        illu_data = None
        data_status = "none"
        if gen_method == "PIL" and priority not in ("skip",):
            illu_data = extract_illustration_data(content_type, title, body_text)
            if illu_data is not None:
                data_status = "auto"
                auto_data_count += 1
            else:
                data_status = "manual"
                manual_data_count += 1
        
        # Build suggestion entry
        suggestion = {
            "slide": idx + 1,  # 1-based for user-facing
            "title": title[:80] if title else "(无标题)",
            "body_text": body_text[:200] if body_text else "",
            "content_type": content_type,
            "confidence": round(confidence, 2),
            "illustration_type": illu_type,
            "generation_method": gen_method,
            "position_preset": position_info["position"],
            "left": position_info.get("left"),
            "top": position_info.get("top"),
            "width": position_info.get("width"),
            "priority": priority,
            "prompt_hint": prompt_hint,
            "description": description,
            "existing_images": existing_images,
            "data_status": data_status,
        }
        suggestions.append(suggestion)
        
        # Only include in pipeline configs if priority is not skip
        if priority != "skip":
            illustrated_count += 1
            
            # gen_illustrations.py config (ENHANCED with data + style)
            # NOTE: source 字段已移除。ImageGen 产出的文件名是随机的，
            # 无法预测。Agent 需在调用 gen_illustrations.py 前：
            #   方案 A: 手动编写 config.json，将 source 字段匹配实际文件名
            #   方案 B: 不传 config，gen_illustrations.py 自动按修改时间顺序匹配
            gen_entry = {
                "slide": idx + 1,
                "type": illu_type,
                "output": f"{idx+1:02d}_{illu_type}.png",
            }
            if illu_data is not None:
                gen_entry["data"] = illu_data
                gen_entry["style"] = style
                gen_entry["data_status"] = "auto"
            elif gen_method == "PIL":
                gen_entry["data_status"] = "manual"
            gen_config.append(gen_entry)
            
            # embed_illustrations.py config
            embed_entry = {
                "slide": idx + 1,
                "image": f"{idx+1:02d}_{illu_type}.png",
            }
            if position_info.get("left") is not None:
                embed_entry["left"] = position_info["left"]
            if position_info.get("top") is not None:
                embed_entry["top"] = position_info["top"]
            if position_info.get("width") is not None:
                embed_entry["width"] = position_info["width"]
            embed_config.append(embed_entry)
            
            # Count for summary
            method_counts[gen_method] = method_counts.get(gen_method, 0) + 1
            type_counts[illu_type] = type_counts.get(illu_type, 0) + 1
    
    coverage_rate = illustrated_count / total_slides if total_slides > 0 else 0
    
    result = {
        "source": str(pptx_path),
        "total_slides": total_slides,
        "image_coverage_needed": illustrated_count,
        "suggestions": suggestions,
        "gen_config": gen_config,
        "embed_config": embed_config,
        "summary": {
            "by_method": method_counts,
            "by_type": type_counts,
            "coverage_rate": round(coverage_rate, 2),
            "coverage_grade": (
                "GOOD" if coverage_rate >= 0.5
                else "FAIR" if coverage_rate >= 0.3
                else "LOW"
            ),
            "data_extraction": {
                "auto": auto_data_count,
                "manual": manual_data_count,
                "auto_rate": round(
                    auto_data_count / max(auto_data_count + manual_data_count, 1), 2
                ),
            },
        },
    }
    
    return result


# ── CLI ──

def main():
    parser = argparse.ArgumentParser(
        description="分析 PPTX 内容，自动建议配图类型、生成方式和嵌入位置，并提取配图结构化数据"
    )
    parser.add_argument("input", help="输入 PPTX 文件路径")
    parser.add_argument("-o", "--output", help="输出 JSON 文件路径（默认 stdout）")
    parser.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="输出详细分析过程",
    )
    parser.add_argument(
        "--min-confidence",
        type=float,
        default=0.3,
        help="最低置信度阈值（默认 0.3，低于此值的建议标记为 low 优先级）",
    )
    parser.add_argument(
        "--scenario",
        default="default",
        help="品牌/场景名称，用于配色建议（如 telecom, tech, nature；默认 default）",
    )
    args = parser.parse_args()
    
    if not Path(args.input).exists():
        print(f"ERROR: 文件不存在: {args.input}", file=sys.stderr)
        sys.exit(1)
    
    if args.verbose:
        print(f"分析: {args.input}", file=sys.stderr)
        print(f"场景: {args.scenario}", file=sys.stderr)
    
    result = analyze_pptx(args.input, verbose=args.verbose, scenario=args.scenario)
    
    # Apply min-confidence filter (just flag, don't remove)
    for s in result["suggestions"]:
        if s["confidence"] < args.min_confidence and s["priority"] not in ("must", "skip"):
            s["priority"] = "low"
            s["low_confidence_note"] = f"置信度 {s['confidence']} < {args.min_confidence}，建议人工确认"
    
    json_str = json.dumps(result, ensure_ascii=False, indent=2)
    
    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(json_str)
        if args.verbose:
            print(f"已写入: {args.output}", file=sys.stderr)
    else:
        print(json_str)
    
    # Print summary to stderr
    summary = result["summary"]
    print(f"\n{'='*60}", file=sys.stderr)
    print(f"配图建议摘要", file=sys.stderr)
    print(f"{'='*60}", file=sys.stderr)
    print(f"总页数: {result['total_slides']}", file=sys.stderr)
    print(f"需配图页数: {result['image_coverage_needed']}", file=sys.stderr)
    print(f"覆盖率: {summary['coverage_rate']*100:.0f}% ({summary['coverage_grade']})", file=sys.stderr)
    de = summary.get("data_extraction", {})
    if de:
        print(f"数据提取: {de.get('auto',0)} 页自动 / {de.get('manual',0)} 页手动 "
              f"(自动率={de.get('auto_rate',0)*100:.0f}%)", file=sys.stderr)
    print(f"\n按生成方式:", file=sys.stderr)
    for method, count in sorted(summary["by_method"].items()):
        print(f"  {method}: {count} 张", file=sys.stderr)
    print(f"\n按配图类型:", file=sys.stderr)
    for itype, count in sorted(summary["by_type"].items()):
        desc = ILLUSTRATION_TYPE_NAMES.get(itype, itype)
        print(f"  {desc} ({itype}): {count} 张", file=sys.stderr)
    print(f"\n逐页建议:", file=sys.stderr)
    for s in result["suggestions"]:
        priority_icon = {"must": "★", "high": "▲", "medium": "●", "low": "○", "skip": "✗"}.get(s["priority"], "?")
        desc = s["description"]
        method = s["generation_method"]
        conf = s["confidence"]
        existing = f" [已有{s['existing_images']}图]" if s["existing_images"] > 0 else ""
        ds = s.get("data_status", "none")
        data_tag = f" [{ds}]" if ds != "none" else ""
        print(f"  P{s['slide']:2d} {priority_icon} {s['content_type']:12s} → {desc} ({method}) "
              f"置信度={conf}{existing}{data_tag}", file=sys.stderr)


if __name__ == "__main__":
    main()
