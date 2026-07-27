#!/usr/bin/env python3
"""
PPT Auto-Fixer — apply common aesthetic repairs to a PPTX file in place.

This complements the structural scorer (score_ppt_pages.py): the scorer finds
problems; this script fixes the ones that can be fixed mechanically.

Repair categories:
  --fonts        : replace unprofessional / inconsistent fonts with a curated
                   pairing (CJK + Latin). See FONT_PAIRINGS below.
  --colors       : unify all shape fills to a curated palette. See PALETTES.
  --margins      : nudge shapes that touch slide edges inward to a safe margin.
  --images       : fix stretched/compressed images by locking aspect ratio.
  --bullets      : convert over-long paragraphs to bullet points (>40 words).
  --theme-bg     : apply a single solid background color to all slides.
  --all          : run most repairs above (fonts, colors, margins, images,
                   bullets, theme_bg, hollow) with sensible defaults.
                   NOTE: --type-scale and --spacing-grid are NOT included
                   because they visibly reflow layout; add them explicitly.

Usage:
    python fix_ppt.py input.pptx --output fixed.pptx --all
    python fix_ppt.py input.pptx --fonts --pairing source-han-sans+inter
    python fix_ppt.py input.pptx --colors --palette corp-blue
    python fix_ppt.py input.pptx --margins --margin-in 0.5
    python fix_ppt.py input.pptx --images

The script never modifies the input file in place — always writes to --output
(default: <input>_fixed.pptx). Inspect the diff by rendering both versions.

Dependencies:
    pip install python-pptx
"""
from __future__ import annotations

import argparse
import copy
import os
import sys
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.dml import MSO_FILL
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx",
          file=sys.stderr)
    sys.exit(1)


EMU_PER_INCH = 914400


# --------------------------------------------------------------------------- #
# Curated font pairings — keyed by ID. Each pairing defines a CJK heading
# font, CJK body font, Latin heading font, Latin body font, and an optional
# numeric/data font. See references/font-pairings.md for rationale.
# --------------------------------------------------------------------------- #
FONT_PAIRINGS: dict[str, dict] = {
    "source-han-sans+inter": {
        "label": "思源黑体 + Inter (现代科技，跨平台)",
        "cjk_heading": "Source Han Sans CN",
        "cjk_body": "Source Han Sans CN",
        "latin_heading": "Inter",
        "latin_body": "Inter",
        "numeric": "DIN Alternate",
    },
    "microsoft-yahei+arial": {
        "label": "微软雅黑 + Arial (Windows 默认，零安装)",
        "cjk_heading": "Microsoft YaHei",
        "cjk_body": "Microsoft YaHei",
        "latin_heading": "Arial",
        "latin_body": "Arial",
        "numeric": "Arial",
    },
    "pingfang+helvetica": {
        "label": "苹方 + Helvetica (Mac 默认，跨端注意)",
        "cjk_heading": "PingFang SC",
        "cjk_body": "PingFang SC",
        "latin_heading": "Helvetica",
        "latin_body": "Helvetica",
        "numeric": "Helvetica Neue",
    },
    "songti+georgia": {
        "label": "宋体 + Georgia (学术/政府/传统品牌)",
        "cjk_heading": "SimSun",
        "cjk_body": "SimSun",
        "latin_heading": "Georgia",
        "latin_body": "Georgia",
        "numeric": "Georgia",
    },
    "fz-xiaobiaosong+fangsong": {
        "label": "方正小标宋 + 仿宋 (党政公文，正式严肃)",
        "cjk_heading": "FZXiaoBiaoSong-B05S",
        "cjk_body": "FangSong",
        "latin_heading": "Times New Roman",
        "latin_body": "Times New Roman",
        "numeric": "Times New Roman",
    },
}


# --------------------------------------------------------------------------- #
# Curated palettes — keyed by ID. Each defines primary, secondary, accent,
# and a set of neutral backgrounds/text colors. See references/color-palettes.md.
# --------------------------------------------------------------------------- #
PALETTES: dict[str, dict] = {
    "corp-blue": {
        "label": "企业蓝 (咨询/金融)",
        "primary": "1B3A5C",      # navy
        "secondary": "4A90D9",    # mid blue
        "accent": "E8792B",       # warm orange
        "bg": "FFFFFF",
        "surface": "F4F6F9",
        "text": "1A1A1A",
        "text_muted": "5A6675",
        "border": "D6DCE5",
    },
    "forest-exec": {
        "label": "森林绿 (ESG/医疗)",
        "primary": "2D4A3E",
        "secondary": "6B9E8A",
        "accent": "D4A843",
        "bg": "FFFFFF",
        "surface": "F2F5F2",
        "text": "1A2620",
        "text_muted": "5C6B62",
        "border": "D5DCD7",
    },
    "charcoal-modern": {
        "label": "炭灰现代 (科技/SaaS)",
        "primary": "333333",
        "secondary": "737373",
        "accent": "0078D4",
        "bg": "FFFFFF",
        "surface": "F5F5F5",
        "text": "1A1A1A",
        "text_muted": "5A5A5A",
        "border": "D9D9D9",
    },
    "slate-pro": {
        "label": "石板专业 (法律/咨询)",
        "primary": "3C3C50",
        "secondary": "7A7A8E",
        "accent": "C84B31",
        "bg": "FFFFFF",
        "surface": "F3F3F6",
        "text": "1F1F2E",
        "text_muted": "5C5C70",
        "border": "D6D6DD",
    },
    "mono-clean": {
        "label": "单色极简 (数据密集/投资)",
        "primary": "1A1A1A",
        "secondary": "B0B0B0",
        "accent": "2196F3",
        "bg": "FFFFFF",
        "surface": "F7F7F7",
        "text": "1A1A1A",
        "text_muted": "707070",
        "border": "DDDDDD",
    },
    "deep-stage": {
        "label": "深色舞台 (高管主题演讲)",
        "primary": "FFFFFF",
        "secondary": "B8B8C8",
        "accent": "FF6B35",
        "bg": "0F1419",
        "surface": "1A2230",
        "text": "F1F5F9",
        "text_muted": "94A3B8",
        "border": "2A3445",
    },
    "telecom-red": {
        "label": "电信红 (中国电信品牌)",
        "primary": "E60012",
        "secondary": "1B3A5C",
        "accent": "FFB800",
        "bg": "FFFFFF",
        "surface": "F9FAFB",
        "text": "1A1A1A",
        "text_muted": "5A6675",
        "border": "D6DCE5",
    },
}


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
def _hex_to_rgbcolor(h: str) -> RGBColor:
    h = h.lstrip("#").upper()
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _is_cjk_text(text: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in text)


def _resolve_pairing(pairing_id: str) -> dict:
    if pairing_id not in FONT_PAIRINGS:
        raise SystemExit(
            f"Unknown font pairing: {pairing_id}. "
            f"Available: {', '.join(FONT_PAIRINGS.keys())}"
        )
    return FONT_PAIRINGS[pairing_id]


def _resolve_palette(palette_id: str) -> dict:
    if palette_id not in PALETTES:
        raise SystemExit(
            f"Unknown palette: {palette_id}. "
            f"Available: {', '.join(PALETTES.keys())}"
        )
    return PALETTES[palette_id]


def _iter_runs(shape):
    """Yield (paragraph, run) tuples for every text run in a shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            yield para, run


# --------------------------------------------------------------------------- #
# Repair: fonts
# --------------------------------------------------------------------------- #
def fix_fonts(prs, pairing: dict, dry_run: bool = False) -> dict:
    """Replace every run's font name with the pairing's CJK or Latin font.

    For mixed CJK+Latin runs, we set the run's font name to the Latin font
    AND set the East Asian font (ea) attribute to the CJK font via XML. This
    is how PowerPoint natively handles CJK+Latin within one run.
    """
    stats = {"runs_seen": 0, "runs_changed": 0, "slides_touched": 0}
    for slide in prs.slides:
        slide_touched = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para, run in _iter_runs(shape):
                stats["runs_seen"] += 1
                text = run.text or ""
                use_cjk = _is_cjk_text(text)
                # Pick the font for this run
                fs = run.font.size
                fs_pt = fs.pt if fs else None
                if fs_pt and fs_pt >= 24:
                    latin = pairing["latin_heading"]
                    cjk = pairing["cjk_heading"]
                else:
                    latin = pairing["latin_body"]
                    cjk = pairing["cjk_body"]

                target_latin = cjk if use_cjk else latin
                old = run.font.name
                if not dry_run:
                    run.font.name = target_latin
                    # Also set the East Asian font via XML so mixed runs work.
                    rPr = run._r.get_or_add_rPr()
                    # Remove existing ea element, then add fresh.
                    for ea in rPr.findall(qn("a:ea")):
                        rPr.remove(ea)
                    ea = rPr.makeelement(qn("a:ea"), {"typeface": cjk})
                    rPr.append(ea)
                if old != target_latin:
                    stats["runs_changed"] += 1
                slide_touched = True
        if slide_touched:
            stats["slides_touched"] += 1
    return stats


# --------------------------------------------------------------------------- #
# Repair: colors
# --------------------------------------------------------------------------- #
def _looks_like_text_color(shape) -> bool:
    """Heuristic: is this shape primarily a text holder (vs. a fill shape)?"""
    if shape.has_text_frame and shape.text_frame.text.strip():
        # If the shape has no fill or transparent fill, treat as text.
        try:
            if shape.fill.type in (None, MSO_FILL.BACKGROUND):
                return True
        except Exception:
            pass
    return False


def fix_colors(prs, palette: dict, dry_run: bool = False) -> dict:
    """Unify all shape fill/line colors to the palette.

    Rules:
      - Text-bearing shapes with no fill: leave fill alone, set text color
        to palette.text or palette.primary for headings.
      - Solid-filled rectangles/cards: remap to palette.primary, .secondary,
        .accent, or .surface based on original saturation.
      - Lines: remap to palette.border.
      - Brand colors (primary, accent, secondary) are protected from remapping.
    """
    stats = {"fills_changed": 0, "lines_changed": 0, "text_colors_changed": 0}
    primary = _hex_to_rgbcolor(palette["primary"])
    secondary = _hex_to_rgbcolor(palette["secondary"])
    accent = _hex_to_rgbcolor(palette["accent"])
    surface = _hex_to_rgbcolor(palette["surface"])
    border = _hex_to_rgbcolor(palette["border"])
    text_color = _hex_to_rgbcolor(palette["text"])
    text_muted = _hex_to_rgbcolor(palette["text_muted"])

    # Build brand color whitelist — colors that should never be remapped.
    brand_whitelist = set()
    for key in ("primary", "accent", "secondary"):
        rgb = _hex_to_rgbcolor(palette[key])
        brand_whitelist.add((int(rgb[0]), int(rgb[1]), int(rgb[2])))

    for slide in prs.slides:
        for shape in slide.shapes:
            # Fill
            try:
                if shape.fill.type == MSO_FILL.SOLID:
                    orig = shape.fill.fore_color.rgb
                    # Map: saturated/strong → primary; mid → secondary;
                    # warm/standout → accent; near-white → surface.
                    new = _remap_fill_color(orig, primary, secondary, accent, surface,
                                            brand_whitelist=brand_whitelist)
                    # If remap kept the original (brand-protected), don't count as changed
                    if new != RGBColor(int(orig[0]), int(orig[1]), int(orig[2])):
                        if not dry_run:
                            shape.fill.fore_color.rgb = new
                        stats["fills_changed"] += 1
            except Exception:
                pass

            # Line color
            try:
                ln = shape.line
                if ln.color and ln.color.type is not None:
                    if not dry_run:
                        ln.color.rgb = border
                    stats["lines_changed"] += 1
            except Exception:
                pass

            # Text run colors — replace pure RGB defaults / oversaturated
            # colors / pure black with palette text/muted/primary.
            # Pure black #000 → palette.text (warmer near-black tint).
            # Brand colors are protected from remapping.
            if shape.has_text_frame:
                for para, run in _iter_runs(shape):
                    try:
                        rc = run.font.color
                        if rc and rc.type is not None:
                            orig = rc.rgb
                            orig_tuple = (int(orig[0]), int(orig[1]), int(orig[2]))
                            orig_hex = f"{orig[0]:02X}{orig[1]:02X}{orig[2]:02X}"
                            # Skip if this is a brand color (e.g. telecom red heading)
                            is_brand = any(
                                abs(orig_tuple[0] - bw[0]) + abs(orig_tuple[1] - bw[1])
                                + abs(orig_tuple[2] - bw[2]) < 60
                                for bw in brand_whitelist
                            )
                            if is_brand:
                                continue
                            new = None
                            if orig_hex == "000000":
                                new = text_color  # warmer near-black
                            elif _is_oversaturated(orig):
                                fs_pt = run.font.size.pt if run.font.size else None
                                new = primary if (fs_pt and fs_pt >= 28) else text_color
                            if new is not None:
                                if not dry_run:
                                    rc.rgb = new
                                stats["text_colors_changed"] += 1
                    except Exception:
                        pass
    return stats


def _remap_fill_color(orig, primary, secondary, accent, surface,
                      brand_whitelist: set | None = None) -> RGBColor:
    """Decide which palette color to substitute for an original fill.

    brand_whitelist: set of (R, G, B) tuples for brand colors that
    should never be remapped.  Typically derived from the palette's
    primary/accent hex values.
    """
    r, g, b = orig[0], orig[1], orig[2]

    # Brand color protection: never remap colors that are close to
    # a brand/primary/accent color (ΔE < 30 approximate threshold).
    if brand_whitelist:
        for br, bg, bb in brand_whitelist:
            if abs(int(r) - int(br)) + abs(int(g) - int(bg)) + abs(int(b) - int(bb)) < 60:
                return RGBColor(int(r), int(g), int(b))  # keep original

    # Near-white → surface
    if r > 235 and g > 235 and b > 235:
        return surface
    # Near-black or very dark → primary (keep dark headers)
    if r < 60 and g < 60 and b < 60:
        return primary
    # Warm (red/orange dominant) → accent
    if r > 180 and r > g + 30 and r > b + 30:
        return accent
    # Default → secondary
    return secondary


def _is_oversaturated(rgb) -> bool:
    r, g, b = rgb[0], rgb[1], rgb[2]
    oversaturated = {
        (255, 0, 0), (0, 255, 0), (0, 0, 255),
        (255, 255, 0), (255, 0, 255), (0, 255, 255),
    }
    return (r, g, b) in oversaturated


# --------------------------------------------------------------------------- #
# Repair: margins — nudge edge-touching shapes inward
# --------------------------------------------------------------------------- #
def fix_margins(prs, margin_pct: float = 0.08, dry_run: bool = False) -> dict:
    """Move shapes that are within margin_pct (default 8%) of a slide edge inward.

    Uses percentage-based margin per element-ratio-scoring.md spec:
    "四边均 8%–12% 安全边距".  Default 8% is the minimum safe margin.

    After nudging, checks for bounding-box overlaps with sibling shapes.
    If an overlap is detected, the nudge is skipped for that shape and
    a warning is logged in stats["overlaps_avoided"].
    
    Also respects template decoration zones: shapes are not nudged into
    positions that overlap non-placeholder decorative elements on the layout.
    """
    stats = {"shapes_moved": 0, "overlaps_avoided": 0, "decoration_avoided": 0}
    sw = prs.slide_width
    sh = prs.slide_height
    margin_h_emu = Emu(int(sw * margin_pct))
    margin_v_emu = Emu(int(sh * margin_pct))

    # Detect template decoration zones per layout
    layout_zones = _detect_layout_exclusion_zones(prs)

    def _bbox_overlap(a, b) -> bool:
        """Check if two bounding boxes (l, t, w, h) overlap."""
        al, at, aw, ah = a
        bl, bt, bw, bh = b
        return not (al + aw <= bl or bl + bw <= al or at + ah <= bt or bt + bh <= at)

    def _overlaps_decoration(l, t, w, h, zones):
        """Check if a bbox overlaps any template decoration zone."""
        pad = int(Inches(0.04))
        for zl, zt, zw, zh in zones:
            if (l < zl + zw + pad and l + w > zl - pad and
                t < zt + zh + pad and t + h > zt - pad):
                return True
        return False

    for slide in prs.slides:
        # Get decoration zones for this slide's layout
        li = _get_layout_index(slide, prs)
        zones = layout_zones.get(li, [])
        
        # Collect all shape bboxes for collision detection.
        shape_info = []
        for shape in slide.shapes:
            try:
                shape_info.append((shape, shape.left, shape.top,
                                   shape.width, shape.height))
            except Exception:
                pass

        for shape, orig_l, orig_t, w, h in shape_info:
            try:
                # Skip placeholder shapes (template-managed)
                if shape.is_placeholder:
                    continue
                    
                l, t = orig_l, orig_t
                moved = False
                if l < margin_h_emu:
                    l = margin_h_emu
                    moved = True
                if t < margin_v_emu:
                    t = margin_v_emu
                    moved = True
                if l + w > sw - margin_h_emu:
                    l = max(margin_h_emu, sw - margin_h_emu - w)
                    moved = True
                if t + h > sh - margin_v_emu:
                    t = max(margin_v_emu, sh - margin_v_emu - h)
                    moved = True
                if not moved:
                    continue
                    
                # Decoration zone awareness:
                # If original position is safe but pushed position overlaps a decoration,
                # revert the push — don't make things worse.
                if zones:
                    original_safe = not _overlaps_decoration(orig_l, orig_t, w, h, zones)
                    pushed_overlaps = _overlaps_decoration(l, t, w, h, zones)
                    
                    if original_safe and pushed_overlaps:
                        # Push made things worse — don't move
                        stats["overlaps_avoided"] += 1
                        continue
                    
                    if not original_safe and pushed_overlaps:
                        # Already overlapping and still overlapping — try nudging below
                        for zl, zt, zw, zh in sorted(zones, key=lambda z: z[1]):
                            if t < zt + zh and t + h > zt:
                                candidate_t = zt + zh + int(Inches(0.04))
                                if candidate_t + h <= sh - margin_v_emu:
                                    t = candidate_t
                                    stats["decoration_avoided"] += 1
                                    break
                        else:
                            # Couldn't fix overlap — skip this nudge
                            stats["overlaps_avoided"] += 1
                            continue
                
                # Check for overlap with other shapes at the new position
                new_bbox = (l, t, w, h)
                has_overlap = False
                for other_shape, ol, ot, ow, oh in shape_info:
                    if other_shape is shape:
                        continue
                    if _bbox_overlap(new_bbox, (ol, ot, ow, oh)):
                        has_overlap = True
                        break
                if has_overlap:
                    stats["overlaps_avoided"] += 1
                    continue  # Skip move to avoid collision
                if not dry_run:
                    shape.left = l
                    shape.top = t
                stats["shapes_moved"] += 1
            except Exception:
                pass
    return stats


# --------------------------------------------------------------------------- #
# Repair: images — fix stretched/compressed aspect ratios
# --------------------------------------------------------------------------- #
def fix_images(prs, dry_run: bool = False) -> dict:
    """Restore image aspect ratios by adjusting the container height to match."""
    stats = {"images_fixed": 0, "images_skipped": 0}
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                img = shape.image
                iw, ih = img.size
                if not iw or not ih:
                    stats["images_skipped"] += 1
                    continue
                img_aspect = iw / ih
                cur_w = shape.width
                cur_h = shape.height
                box_aspect = cur_w / cur_h if cur_h else 0
                if not box_aspect:
                    continue
                ratio = box_aspect / img_aspect
                if 0.96 <= ratio <= 1.04:
                    continue  # already correct
                # Adjust height to match width (keep width fixed)
                new_h = int(cur_w / img_aspect)
                if not dry_run:
                    shape.height = new_h
                stats["images_fixed"] += 1
            except Exception:
                stats["images_skipped"] += 1
    return stats


# --------------------------------------------------------------------------- #
# Repair: image borders — add subtle border to images for visual definition
# --------------------------------------------------------------------------- #
def fix_image_borders(prs, border_pt: float = 1.5, dry_run: bool = False) -> dict:
    """Add a subtle border to all images for visual definition.

    Uses a thin line (default 1.5pt) in a neutral color to give images
    a clean edge, especially useful on dark-themed slides where images
    may blend into the background.
    """
    from pptx.util import Pt
    stats = {"images_bordered": 0, "images_skipped": 0}
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                if not dry_run:
                    line = shape.line
                    line.width = Pt(border_pt)
                    # Use a neutral border color that works on both light/dark
                    line.color.rgb = RGBColor(0xB8, 0xBE, 0xC5)
                stats["images_bordered"] += 1
            except Exception:
                stats["images_skipped"] += 1
    return stats


# --------------------------------------------------------------------------- #
# Repair: theme background — apply a single solid background to all slides
# --------------------------------------------------------------------------- #
def fix_theme_background(prs, bg_hex: str, dry_run: bool = False) -> dict:
    """Apply a solid background fill to every slide."""
    stats = {"slides_updated": 0}
    bg_color = _hex_to_rgbcolor(bg_hex)
    for slide in prs.slides:
        try:
            # Use the slide's background fill (python-pptx >= 0.6.18)
            bg = slide.background
            fill = bg.fill
            if not dry_run:
                fill.solid()
                fill.fore_color.rgb = bg_color
            stats["slides_updated"] += 1
        except Exception:
            pass
    return stats


# --------------------------------------------------------------------------- #
# Repair: bullets — split long paragraphs into bullet points
# --------------------------------------------------------------------------- #
def fix_bullets(prs, max_words: int = 40, max_cjk_chars: int = 35,
                dry_run: bool = False) -> dict:
    """Convert long paragraphs into bullet-pointed sub-runs.

    Triggers when EITHER:
      - English word count > max_words (default 40), OR
      - CJK character count > max_cjk_chars (default 35, per web-design
        principle: Chinese single-bullet ideal 15-25, hard cap 35).

    Splits on sentence-ending delimiters (。；;！!？?．) for CJK text
    and additionally on commas (，,) for English text only.
    CJK commas (，) are NOT treated as split points because they are
    intra-sentence pauses (e.g. "包括，但不限于" must not be split).
    """
    import re
    stats = {"paragraphs_split": 0}
    # Sentence-ending delimiters (safe for both CJK and English)
    sentence_end_re = re.compile(r"[。；;！!？?．]")
    # Full splitter including commas (English text only)
    full_split_re = re.compile(r"[。；;！!？?．，,]")
    bullet_char = "• "

    def _cjk_count(s):
        return sum(1 for ch in s if "\u4e00" <= ch <= "\u9fff")

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            # Iterate over a snapshot since we'll mutate paragraphs.
            paras = list(tf.paragraphs)
            for para in paras:
                text = para.text.strip()
                if not text:
                    continue
                word_count = len(text.split())
                cjk_count = _cjk_count(text)
                # Trigger if either threshold exceeded.
                if word_count < max_words and cjk_count < max_cjk_chars:
                    continue
                # Use CJK-safe splitter (no comma) for CJK-dominant text,
                # full splitter (with comma) for English-dominant text.
                is_cjk = _is_cjk_text(text)
                split_re = sentence_end_re if is_cjk else full_split_re
                # Only split if there are actual sentence delimiters
                parts = [p.strip() for p in split_re.split(text) if p.strip()]
                if len(parts) < 2:
                    continue
                # Replace original paragraph text with the first part,
                # then add the rest as new paragraphs.
                if not dry_run:
                    # Clear runs in the original paragraph, then set first part.
                    for r in list(para.runs):
                        r.text = ""
                    if para.runs:
                        para.runs[0].text = bullet_char + parts[0]
                    else:
                        run = para.add_run()
                        run.text = bullet_char + parts[0]
                    # Append the remaining parts as new paragraphs.
                    for p in parts[1:]:
                        new_p = tf.add_paragraph()
                        # Copy basic formatting from original
                        try:
                            new_p.alignment = para.alignment
                        except Exception:
                            pass
                        run = new_p.add_run()
                        run.text = bullet_char + p
                        # Copy font from original run if present
                        if para.runs:
                            src_font = para.runs[0].font
                            try:
                                run.font.name = src_font.name
                                run.font.size = src_font.size
                                run.font.bold = src_font.bold
                            except Exception:
                                pass
                stats["paragraphs_split"] += 1
    return stats


# --------------------------------------------------------------------------- #
# Repair: type scale — re-bucket run font sizes to a Major Third / Perfect
# Fourth / golden-ratio ladder. See references/web-design-principles.md §1.
# --------------------------------------------------------------------------- #
def fix_type_scale(prs, ratio: float = 1.25, scenario: str | None = None,
                   dry_run: bool = False) -> dict:
    """Snap every run's font size to the nearest rung of a Type Scale ladder.

    The ladder is built from the deck's modal (most common) body size as the
    base, then multiplied by ``ratio`` for each successive rung:
        base, base*ratio, base*ratio**2, base*ratio**3, base*ratio**4
    e.g. ratio=1.25, base=20pt -> [20, 25, 31.25, 39, 49]
    Each run size is rounded to the nearest rung (clamped to [8, 96]pt).

    scenario: optional scenario hint that overrides the default ratio:
      - "executive":  ratio=1.5  (strong hierarchy for high-level decks)
      - "marketing":  ratio=1.333 (moderate hierarchy, visual punch)
      - "data":       ratio=1.25  (compact, space-efficient)
      - "creative":   ratio=1.333 (balanced for expressive layouts)
      - "gov":        ratio=1.25  (formal, understated hierarchy)
      - "telecom":    ratio=1.333 (moderate, brand-forward)
    If scenario is provided and no explicit ratio override is given via CLI,
    the scenario's recommended ratio is used.
    """
    from collections import Counter as _Counter

    # Scenario → recommended ratio mapping
    SCENARIO_RATIOS = {
        "executive": 1.5,
        "marketing": 1.333,
        "data": 1.25,
        "creative": 1.333,
        "gov": 1.25,
        "telecom": 1.333,
    }
    if scenario and scenario in SCENARIO_RATIOS:
        ratio = SCENARIO_RATIOS[scenario]

    stats = {"runs_seen": 0, "runs_changed": 0, "base_pt": None,
             "ratio_used": ratio, "scenario": scenario}

    # Pass 1: collect all run sizes to find the modal "body" size.
    sizes = _Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    fs = run.font.size
                    if fs:
                        sizes[round(fs.pt)] += 1
    if not sizes:
        return stats

    # The modal size is the "body" base. If multiple modes, pick the smallest
    # among the top 3 (body text is usually the most common AND smallest).
    top3 = [s for s, _ in sizes.most_common(3)]
    base = min(top3) if top3 else 18
    stats["base_pt"] = base

    # Build the ladder (5 rungs covers body -> cover title).
    ladder = [round(base * (ratio ** i), 1) for i in range(5)]
    # Clamp to a sane PPT range.
    ladder = [max(8, min(96, r)) for r in ladder]

    def nearest_rung(pt):
        """Snap to the nearest ladder rung, but always push sizes larger
        than ``base`` UP to the next rung (never back down to base).

        This guarantees that an existing title (even if only slightly
        larger than body) ends up at least one rung above body, so the
        Type Scale hierarchy becomes visible instead of flattened.
        """
        if pt <= base:
            return base
        # Find the smallest rung strictly greater than base that is >= pt.
        for r in ladder[1:]:
            if r >= pt:
                return r
        # pt exceeds the top rung — clamp to the top.
        return ladder[-1]

    # Pass 2: snap each run to the nearest rung.
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    fs = run.font.size
                    if not fs:
                        continue
                    stats["runs_seen"] += 1
                    cur = round(fs.pt, 1)
                    new = nearest_rung(cur)
                    if new != cur:
                        if not dry_run:
                            run.font.size = Pt(new)
                        stats["runs_changed"] += 1
    return stats


# --------------------------------------------------------------------------- #
# Repair: spacing grid — snap shape left/top/width/height to 0.08in multiples
# (the PPT equivalent of the web 4/8px rhythm). See web-design-principles §3.
# --------------------------------------------------------------------------- #
def _detect_layout_exclusion_zones(prs):
    """Detect decorative elements in slide layouts that content should avoid.
    
    Returns a dict mapping layout_index → list of (left, top, width, height) in EMU
    for non-placeholder shapes on each layout. These represent template decorations
    (divider lines, logos, ribbons, etc.) that content shapes should not overlap.
    """
    zones = {}
    for li, layout in enumerate(prs.slide_layouts):
        layout_zones = []
        for s in layout.shapes:
            # Only consider non-placeholder decorative shapes
            if s.is_placeholder:
                continue
            try:
                l, t, w, h = s.left, s.top, s.width, s.height
                # Skip very small shapes (< 0.05 sq.in) — likely artifacts
                if w * h < Inches(0.05) * Inches(0.05):
                    continue
                layout_zones.append((l, t, w, h))
            except Exception:
                continue
        zones[li] = layout_zones
    return zones


def _get_layout_index(slide, prs):
    """Get the 0-based layout index for a slide."""
    for li, layout in enumerate(prs.slide_layouts):
        if slide.slide_layout is layout:
            return li
    return -1


def fix_spacing_grid(prs, grid_in: float = 0.08, dry_run: bool = False) -> dict:
    """Snap every shape's EMU position/size to the nearest grid multiple.

    Preserves a safe page margin: shapes that would be pushed below
    8% of slide width/height from an edge by snapping are clamped back.
    
    Also respects template decoration zones: if a layout has decorative
    elements (divider lines, logos), content shapes are snapped to positions
    that avoid overlapping them.
    """
    stats = {"shapes_snapped": 0, "shapes_skipped": 0}
    grid_emu = Emu(Inches(grid_in))
    sw = prs.slide_width
    sh = prs.slide_height
    # Percentage-based minimum safety margin (8% per spec)
    min_margin_h = Emu(int(sw * 0.08))
    min_margin_v = Emu(int(sh * 0.08))

    # Detect template decoration zones per layout
    layout_zones = _detect_layout_exclusion_zones(prs)

    def snap(v):
        return int(round(v / grid_emu) * grid_emu)

    def shape_overlaps_zone(new_l, new_t, new_w, new_h, zones):
        """Check if a shape bbox overlaps any decoration zone."""
        # Add a small padding (0.04in) around zones for visual clearance
        pad = int(Inches(0.04))
        for zl, zt, zw, zh in zones:
            if (new_l < zl + zw + pad and new_l + new_w > zl - pad and
                new_t < zt + zh + pad and new_t + new_h > zt - pad):
                return True
        return False

    for slide in prs.slides:
        # Get layout index to find decoration zones
        li = _get_layout_index(slide, prs)
        zones = layout_zones.get(li, [])
        
        for shape in slide.shapes:
            try:
                l, t, w, h = shape.left, shape.top, shape.width, shape.height
                
                # Skip shapes that are part of the layout (decorations themselves)
                if shape.is_placeholder:
                    continue  # Don't move placeholders — they're template-managed
                
                new_l = snap(l)
                new_t = snap(t)
                new_w = snap(w)
                new_h = snap(h)
                # Clamp to safe margins (don't push content into the edge).
                new_l = max(new_l, min_margin_h)
                new_t = max(new_t, min_margin_v)
                if new_l + new_w > sw - min_margin_h:
                    new_l = max(min_margin_h, sw - min_margin_h - new_w)
                if new_t + new_h > sh - min_margin_v:
                    new_t = max(min_margin_v, sh - min_margin_v - new_h)
                
                # Decoration zone awareness:
                # If the ORIGINAL position does NOT overlap a decoration zone,
                # but the PUSHED position DOES, then the push is wrong —
                # revert to the original position (not the pushed one).
                # This prevents margin/grid logic from pushing a title
                # (originally ABOVE a divider line) down INTO the line.
                original_safe = True
                pushed_overlaps = False
                if zones:
                    original_safe = not shape_overlaps_zone(l, t, w, h, zones)
                    pushed_overlaps = shape_overlaps_zone(new_l, new_t, new_w, new_h, zones)
                
                if original_safe and pushed_overlaps:
                    # The push made things worse — revert position
                    new_l = l
                    new_t = t
                
                # If neither original nor pushed is safe, try nudging below the zone
                if not original_safe and pushed_overlaps and zones:
                    for zl, zt, zw, zh in sorted(zones, key=lambda z: z[1]):
                        if new_t < zt + zh and new_t + new_h > zt:
                            candidate_t = snap(zt + zh + int(Inches(0.04)))
                            if candidate_t + new_h <= sh - min_margin_v:
                                new_t = candidate_t
                                break
                
                if (new_l, new_t, new_w, new_h) != (l, t, w, h):
                    if not dry_run:
                        shape.left = new_l
                        shape.top = new_t
                        shape.width = new_w
                        shape.height = new_h
                    stats["shapes_snapped"] += 1
                else:
                    stats["shapes_skipped"] += 1
            except Exception:
                stats["shapes_skipped"] += 1
    return stats


# --------------------------------------------------------------------------- #
# Repair: hollow containers — shrink oversized card backgrounds to match
# actual content extent, or add a table as the content container.
# --------------------------------------------------------------------------- #
def fix_hollow(prs, dry_run: bool = False) -> dict:
    """Fix hollow containers (anti-pattern AP-LA-09).

    Detection logic mirrors score_ppt_pages.py:
      - Find large non-text shapes (RECTANGLE, ROUNDED_RECTANGLE, PICTURE
        with area >= 5 sq.in) that serve as card backgrounds.
      - Exclude TABLE shapes and horizontal strip cards (wide > 70% slide,
        height < 2.0 in).
      - Compute fill rate = (text + table area inside container) / container area.
      - If fill rate < 60%, shrink the container background to tightly wrap
        the content (content bounding box + 0.15" padding on each side).

    Returns a dict with stats.
    """
    FILL_THRESHOLD = 0.60
    MIN_CONTAINER_AREA_SQIN = 5.0
    EMU_PER_INCH = 914400

    stats = {"containers_examined": 0, "containers_fixed": 0,
             "containers_skipped": 0}

    for slide in prs.slides:
        sw = prs.slide_width
        sh = prs.slide_height

        # Collect all shapes with bounding boxes
        shapes_info = []
        for sp in slide.shapes:
            try:
                l, t, w, h = sp.left, sp.top, sp.width, sp.height
                shapes_info.append((sp, l, t, w, h))
            except Exception:
                continue

        # Identify container candidates
        for sp, l, t, w, h in shapes_info:
            area_sqin = (w / EMU_PER_INCH) * (h / EMU_PER_INCH)
            if area_sqin < MIN_CONTAINER_AREA_SQIN:
                continue

            # Skip TABLE (tables ARE the fix)
            if sp.shape_type == MSO_SHAPE_TYPE.TABLE:
                continue

            # Only check RECTANGLE, ROUNDED_RECTANGLE, or PICTURE >= 8sqin
            try:
                auto_type = sp.auto_shape_type
                # auto_shape_type returns MSO_SHAPE IntEnum, not string.
                # Use .name to get "RECTANGLE", "ROUNDED_RECTANGLE", etc.
                type_name = auto_type.name if hasattr(auto_type, "name") else str(auto_type)
                is_rect = type_name in ("RECTANGLE", "ROUNDED_RECTANGLE")
            except (ValueError, AttributeError):
                is_rect = False
            is_picture = sp.shape_type == MSO_SHAPE_TYPE.PICTURE
            if not (is_rect or is_picture or area_sqin >= 8.0):
                continue

            # Skip horizontal strip cards (wide+short)
            width_in = w / EMU_PER_INCH
            slide_width_in = sw / EMU_PER_INCH
            height_in = h / EMU_PER_INCH
            if width_in > slide_width_in * 0.70 and height_in < 2.0:
                continue

            stats["containers_examined"] += 1

            # Find content bounding box inside this container
            cl, ct, cr, cb = l, t, l + w, t + h
            content_min_x, content_min_y = cr, cb
            content_max_x, content_max_y = cl, ct

            for other_sp, ol, ot, ow, oh in shapes_info:
                if other_sp is sp:
                    continue
                # Skip other containers
                if other_sp.shape_type in (MSO_SHAPE_TYPE.TABLE,
                                           MSO_SHAPE_TYPE.PICTURE):
                    continue
                if not other_sp.has_text_frame:
                    continue

                # Check if this text box is inside the container
                ocx = ol + ow / 2
                ocy = ot + oh / 2
                if cl <= ocx <= cr and ct <= ocy <= cb:
                    content_min_x = min(content_min_x, ol)
                    content_min_y = min(content_min_y, ot)
                    content_max_x = max(content_max_x, ol + ow)
                    content_max_y = max(content_max_y, ot + oh)

            # If no content found, skip
            if content_min_x >= content_max_x or content_min_y >= content_max_y:
                stats["containers_skipped"] += 1
                continue

            # Compute content area and fill rate
            content_area = ((content_max_x - content_min_x) / EMU_PER_INCH *
                            (content_max_y - content_min_y) / EMU_PER_INCH)
            fill_rate = content_area / area_sqin if area_sqin > 0 else 1.0

            if fill_rate >= FILL_THRESHOLD:
                stats["containers_skipped"] += 1
                continue

            # Fix: shrink container to content bbox + 0.15" padding
            PAD = int(0.15 * EMU_PER_INCH)
            new_l = max(0, content_min_x - PAD)
            new_t = max(0, content_min_y - PAD)
            new_w = content_max_x - content_min_x + 2 * PAD
            new_h = content_max_y - content_min_y + 2 * PAD

            if not dry_run:
                try:
                    sp.left = new_l
                    sp.top = new_t
                    sp.width = new_w
                    sp.height = new_h
                except Exception:
                    stats["containers_skipped"] += 1
                    continue

            stats["containers_fixed"] += 1

    return stats


# --------------------------------------------------------------------------- #
# Main
# --------------------------------------------------------------------------- #
def main():
    ap = argparse.ArgumentParser(
        description="Apply common aesthetic repairs to a PPTX file.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    ap.add_argument("pptx_path", help="Input PPTX path")
    ap.add_argument("--output", "-o", default=None,
                    help="Output path (default: <input>_fixed.pptx)")
    ap.add_argument("--dry-run", action="store_true",
                    help="Report what would change without writing a file")

    repair_grp = ap.add_argument_group("Repairs")
    repair_grp.add_argument("--all", action="store_true",
                            help="Run all repairs below with defaults")
    repair_grp.add_argument("--fonts", action="store_true",
                            help="Replace fonts with a curated CJK+Latin pairing")
    repair_grp.add_argument("--colors", action="store_true",
                            help="Unify shape fills/lines to a curated palette")
    repair_grp.add_argument("--margins", action="store_true",
                            help="Nudge edge-touching shapes inward")
    repair_grp.add_argument("--images", action="store_true",
                            help="Fix stretched/compressed image aspect ratios")
    repair_grp.add_argument("--bullets", action="store_true",
                            help="Split long paragraphs into bullet points")
    repair_grp.add_argument("--theme-bg", action="store_true",
                            help="Apply a single solid background to all slides")
    repair_grp.add_argument("--type-scale", action="store_true",
                            help="Snap font sizes to a Type Scale ladder "
                                 "(web-design-engineer principle)")
    repair_grp.add_argument("--spacing-grid", action="store_true",
                            help="Snap shape positions to a 0.08in grid "
                                 "(web 4/8px rhythm)")
    repair_grp.add_argument("--hollow", action="store_true",
                            help="Fix hollow containers: shrink card backgrounds to "
                                 "match actual content size, or replace add_bg()+"
                                 "add_rich_textbox() layout with table-based layout")

    opt_grp = ap.add_argument_group("Options")
    opt_grp.add_argument("--pairing", default="microsoft-yahei+arial",
                         choices=list(FONT_PAIRINGS.keys()),
                         help="Font pairing ID (default: microsoft-yahei+arial)")
    opt_grp.add_argument("--palette", default="corp-blue",
                         choices=list(PALETTES.keys()),
                         help="Color palette ID (default: corp-blue)")
    opt_grp.add_argument("--margin-pct", type=float, default=0.08,
                         help="Margin as percentage of slide dimension (default: 0.08 = 8%%)")
    opt_grp.add_argument("--margin-in", type=float, default=None,
                         help="Margin in inches (deprecated; use --margin-pct instead)")
    opt_grp.add_argument("--bg-hex", default=None,
                         help="Background hex (default: palette.bg)")
    opt_grp.add_argument("--scenario", default=None,
                         choices=["executive", "marketing", "data", "creative",
                                  "gov", "telecom"],
                         help="Scenario hint for auto-configuring type scale ratio "
                              "and other defaults (overrides --ratio if set)")
    opt_grp.add_argument("--ratio", type=float, default=None,
                         help="Type Scale ratio (1.25 / 1.333 / 1.5 / 1.618). "
                              "Default depends on --scenario, or 1.25 if no scenario.")
    opt_grp.add_argument("--grid-in", type=float, default=0.08,
                         help="Spacing grid unit in inches (default: 0.08)")
    opt_grp.add_argument("--dark-theme", action="store_true",
                         help="Protect dark theme: skip --colors and --theme-bg "
                              "when the PPT has a dark background (auto-detected)")
    opt_grp.add_argument("--image-borders", action="store_true",
                         help="Add subtle border to images (1.5pt, matching "
                              "palette border color) for visual definition")
    opt_grp.add_argument("--issues", default=None,
                         help="Path to critique_engine JSON report. Auto-enables "
                              "repairs for detected issues.")

    args = ap.parse_args()

    in_path = Path(args.pptx_path)
    if not in_path.exists():
        print(f"Error: File not found: {in_path}", file=sys.stderr)
        sys.exit(1)

    out_path = args.output or str(in_path.parent / (in_path.stem + "_fixed.pptx"))

    # If --all, enable every repair
    if args.all:
        args.fonts = True
        args.colors = True
        args.margins = True
        args.images = True
        args.bullets = True
        args.theme_bg = True
        args.hollow = True
        # Note: type-scale and spacing-grid are NOT in --all because they
        # visibly reflow layout; opt in explicitly when desired.

    # If --issues, load the JSON report and auto-enable relevant repairs
    if args.issues:
        import json
        issues_path = Path(args.issues)
        if not issues_path.exists():
            print(f"Error: Issues file not found: {issues_path}", file=sys.stderr)
            sys.exit(1)
        with open(issues_path, encoding="utf-8") as f:
            report = json.load(f)
        ap_keys = set()
        for slide_data in report.get("slides", []):
            for ap in slide_data.get("anti_patterns", []):
                ap_keys.add(ap.split(":", 1)[0].strip())
        # Map anti-patterns to repairs
        ap_to_repair = {
            # -- hollow / card containers --
            "hollow_container": "hollow",
            "low_fill_rate": "hollow",
            "no_card_containers": "hollow",
            "bare_text_no_card": "hollow",
            "empty_quadrant": "hollow",
            # -- fonts --
            "unprofessional_font": "fonts",
            "cjk_text_no_cjk_font": "fonts",
            "too_many_fonts": "fonts",
            "inconsistent_fonts": "fonts",
            "deck_too_many_fonts": "fonts",
            "code_without_monospace": "fonts",
            # -- colors --
            "too_many_colors": "colors",
            "inconsistent_colors": "colors",
            "ai_purple_palette": "colors",
            "oversaturated_pure_rgb": "colors",
            "pure_black_text": "colors",
            "pure_white_card": "colors",
            "rainbow_text": "colors",
            "overuse_accent": "colors",
            "deck_too_many_colors": "colors",
            "no_brand_consistency": "colors",
            "low_contrast": "colors",
            # -- margins --
            "no_margins": "margins",
            "high_fill_rate": "margins",
            "image_at_edge": "margins",
            "inconsistent_margins": "margins",
            # -- images --
            "stretched_image": "images",
            # NOTE: stretched_images (plural) was a typo; score_ppt_pages emits singular
            # -- bullets / text density --
            "wall_of_text": "bullets",
            "high_text_density": "bullets",
            "text_heavy_deck": "bullets",
            "long_bullet": "bullets",
            "too_many_bullets": "bullets",
            "bullet_soup": "bullets",
            "orphan_widow": "bullets",
            # -- type_scale --
            "tiny_text": "type_scale",
            "weak_type_scale": "type_scale",
            "extreme_type_scale": "type_scale",
            "no_visual_hierarchy": "type_scale",
            "key_data_not_emphasized": "type_scale",
            # -- spacing_grid --
            "irregular_spacing": "spacing_grid",
            "inconsistent_spacing": "spacing_grid",
            "misaligned_elements": "spacing_grid",
            "top_heavy": "spacing_grid",
            "three_equal_cards": "spacing_grid",
            "unbalanced_layout": "spacing_grid",
            "inconsistent_alignment": "spacing_grid",
            "inconsistent_section_pos": "spacing_grid",
            "top_text_bottom_image": "spacing_grid",
            # -- no auto-fix available (logged but not repairable) --
            # low_res_image: source image resolution cannot be improved mechanically
            # chart_without_message: requires content authoring
            # default_template: requires design replacement
            # table_no_zebra: table styling not yet implemented
            # table_weak_header: table styling not yet implemented
            # mixed_numbering: numbering logic not yet implemented
            # cover_ending_mismatch: structural, requires manual design
            # density_outlier: deck-level, requires manual review
            # no_header_navigation: structural, requires manual design
            # repeated_chapter_labels: content-level, requires manual editing
        }
        enabled_repairs = set()
        for ap_key in ap_keys:
            repair = ap_to_repair.get(ap_key)
            if repair and repair not in enabled_repairs:
                enabled_repairs.add(repair)
                setattr(args, repair, True)
                print(f"  [issues] auto-enabling --{repair} (from {ap_key})")

    if not any([args.fonts, args.colors, args.margins,
                args.images, args.bullets, args.theme_bg,
                args.type_scale, args.spacing_grid, args.hollow]):
        print("No repair selected. Use --all or pick specific repairs. "
              "See --help.", file=sys.stderr)
        sys.exit(1)

    print(f"Loading: {in_path}")
    prs = Presentation(str(in_path))
    pairing = _resolve_pairing(args.pairing)
    palette = _resolve_palette(args.palette)
    bg_hex = args.bg_hex or palette["bg"]

    # Dark theme detection: auto-detect when --all or --dark-theme is used
    is_dark_theme = False
    if args.dark_theme or args.all:
        try:
            from pptx.enum.dml import MSO_FILL as _MSO_FILL
            for slide in prs.slides:
                try:
                    bg_fill = slide.background.fill
                    if bg_fill.type == _MSO_FILL.SOLID:
                        rgb = bg_fill.fore_color.rgb
                        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
                        lum = (0.2126 * (r/255) + 0.7152 * (g/255) + 0.0722 * (b/255))
                        if lum < 0.2:
                            is_dark_theme = True
                            break
                except Exception:
                    pass
        except Exception:
            pass

    if is_dark_theme:
        print("WARNING: Dark theme detected -- skipping --colors and --theme-bg to preserve dark palette")
        args.colors = False
        args.theme_bg = False

    print(f"Font pairing: {pairing['label']}")
    print(f"Palette:      {palette['label']}")
    print(f"Background:   #{bg_hex}")
    # Backward compatibility: --margin-in overrides --margin-pct if provided
    margin_pct = args.margin_pct
    if args.margin_in is not None:
        # Convert absolute inches to percentage based on slide width
        sw_in = prs.slide_width / 914400
        margin_pct = args.margin_in / sw_in
        print(f"Margin:       {args.margin_in} in (= {margin_pct:.1%} of slide width)")
    else:
        print(f"Margin:       {margin_pct:.0%} of slide dimension")
    if args.type_scale:
        # Determine effective ratio: explicit --ratio overrides --scenario
        effective_ratio = args.ratio  # may be None
        scenario_name = args.scenario
        if effective_ratio is None and scenario_name is None:
            effective_ratio = 1.25  # default
        print(f"Type scale:   "
              f"ratio={effective_ratio or '(from scenario)'}"
              f"{f' scenario={scenario_name}' if scenario_name else ''}")
    if args.spacing_grid:
        print(f"Spacing grid: {args.grid_in} in")
    print()

    if args.fonts:
        s = fix_fonts(prs, pairing, dry_run=args.dry_run)
        print(f"[fonts]   runs={s['runs_seen']} changed={s['runs_changed']} "
              f"slides_touched={s['slides_touched']}")
    if args.colors:
        s = fix_colors(prs, palette, dry_run=args.dry_run)
        print(f"[colors]  fills={s['fills_changed']} lines={s['lines_changed']} "
              f"text_colors={s['text_colors_changed']}")
    if args.margins:
        s = fix_margins(prs, margin_pct=margin_pct, dry_run=args.dry_run)
        print(f"[margins] shapes_moved={s['shapes_moved']}")
    if args.images:
        s = fix_images(prs, dry_run=args.dry_run)
        print(f"[images]  fixed={s['images_fixed']} skipped={s['images_skipped']}")
    if args.image_borders:
        s = fix_image_borders(prs, dry_run=args.dry_run)
        print(f"[img-borders] bordered={s['images_bordered']} skipped={s['images_skipped']}")
    if args.bullets:
        s = fix_bullets(prs, dry_run=args.dry_run)
        print(f"[bullets] paragraphs_split={s['paragraphs_split']}")
    if args.theme_bg:
        s = fix_theme_background(prs, bg_hex, dry_run=args.dry_run)
        print(f"[theme-bg] slides_updated={s['slides_updated']}")
    if args.type_scale:
        s = fix_type_scale(prs, ratio=args.ratio or 1.25, scenario=args.scenario,
                           dry_run=args.dry_run)
        scenario_str = f" scenario={s['scenario']}" if s.get("scenario") else ""
        print(f"[type-scale] base={s['base_pt']}pt ratio={s['ratio_used']}"
              f"{scenario_str} "
              f"runs={s['runs_seen']} changed={s['runs_changed']}")
    if args.spacing_grid:
        s = fix_spacing_grid(prs, grid_in=args.grid_in, dry_run=args.dry_run)
        print(f"[spacing-grid] snapped={s['shapes_snapped']} "
              f"skipped={s['shapes_skipped']}")
    if args.hollow:
        s = fix_hollow(prs, dry_run=args.dry_run)
        print(f"[hollow]  examined={s['containers_examined']} "
              f"fixed={s['containers_fixed']} "
              f"skipped={s['containers_skipped']}")

    if args.dry_run:
        print("\n(dry-run; no file written)")
        return

    print(f"\nSaving: {out_path}")
    prs.save(out_path)
    print("Done. Render both versions to compare:")
    print(f"  python scripts/render_slides.py \"{in_path}\"")
    print(f"  python scripts/render_slides.py \"{out_path}\"")


if __name__ == "__main__":
    main()
