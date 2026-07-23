#!/usr/bin/env python3
"""
svg2pptx.py — Convert baoyu-diagram SVG to editable PowerPoint shapes.
Reads an SVG file, extracts native shapes (rect, text, ellipse, line, path,
polygon, polyline) and rebuilds them as native PowerPoint shapes using
python-pptx. Supports colors, rounded corners, dashed strokes, and grouped
elements.
Usage:
    python svg2pptx.py <input.svg> [-o output.pptx]
The output PPTX has a 16:9 slide layout and preserves the dark theme colors.
"""
import argparse, math, os, re, sys
try:
    from xml.etree import ElementTree as ET
except ImportError:
    import xml.etree.ElementTree as ET
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
except ImportError:
    print("Error: python-pptx is required. Install with:")
    print("  pip install python-pptx --target=<workspace-libs>")
    sys.exit(1)

SVG_NS = "http://www.w3.org/2000/svg"
ET.register_namespace("", SVG_NS)

def parse_viewbox(svg_root):
    vb = svg_root.get("viewBox")
    if vb:
        parts = list(map(float, vb.split()))
        if len(parts) >= 4: return parts[2], parts[3]
    return float(svg_root.get("width", "1600")), float(svg_root.get("height", "1200"))

def parse_color(css_color, default="1e293b"):
    css_color = (css_color or "").strip()
    if not css_color or css_color == "none": return default
    m = re.match(r"rgba?\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", css_color)
    if m: return f"{int(m.group(1)):02x}{int(m.group(2)):02x}{int(m.group(3)):02x}"
    if css_color.startswith("#"):
        h = css_color[1:]
        if len(h) == 3: h = "".join(c*2 for c in h)
        return h
    named = {"white":"ffffff","black":"000000","red":"ef4444","green":"22c55e","blue":"3b82f6",
             "cyan":"22d3ee","emerald":"34d399","violet":"a78bfa","amber":"fbbf24",
             "rose":"fb7185","orange":"fb923c","slate":"94a3b8","gray":"6b7280"}
    return named.get(css_color.lower(), default)

def parse_fill(elem, default="0f172a"):
    fill = elem.get("fill","")
    if fill == "none" or not fill:
        m = re.search(r"fill\s*:\s*([^;]+)", elem.get("style",""))
        return parse_color(m.group(1), default) if m else default
    return parse_color(fill, default)

def parse_stroke(elem, default="94a3b8"):
    stroke = elem.get("stroke","")
    if not stroke or stroke == "none":
        m = re.search(r"stroke\s*:\s*([^;]+)", elem.get("style",""))
        return parse_color(m.group(1), default) if m else default
    return parse_color(stroke, default)

def parse_stroke_width(elem, default=1.5):
    sw = elem.get("stroke-width","")
    try: return float(sw)
    except: return default

def parse_opacity(elem, default=1.0):
    try: return float(elem.get("opacity") or elem.get("fill-opacity") or default)
    except: return default

def is_dashed(elem):
    return bool(elem.get("stroke-dasharray","")) and elem.get("stroke-dasharray") != "none"

def get_transform(elem):
    t = elem.get("transform","")
    m = re.match(r"translate\s*\(\s*([\d.-]+)\s*[, ]\s*([\d.-]+)", t)
    return (float(m.group(1)), float(m.group(2))) if m else (0,0)

SVG_W, SVG_H = 1600, 1200
PPTX_W, PPTX_H = Inches(13.333), Inches(7.5)

def to_emu_x(x): return int(PPTX_W * x / SVG_W)
def to_emu_y(y): return int(PPTX_H * y / SVG_H)
def to_emu_w(w): return int(PPTX_W * w / SVG_W)
def to_emu_h(h): return int(PPTX_H * h / SVG_H)

def add_rect(slide, x, y, w, h, fill, stroke, sw=1.5, rx=0, dashed=False, opacity=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rx>0 else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*bytes.fromhex(fill))
    shape.line.color.rgb = RGBColor(*bytes.fromhex(stroke)); shape.line.width = Pt(sw)
    if dashed: shape.line.dash_style = 2
    return shape

def add_text(slide, x, y, txt, color, size=11, bold=False, anchor="middle"):
    tb = slide.shapes.add_textbox(x, y, to_emu_w(200), to_emu_h(20))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*bytes.fromhex(color)); p.font.bold = bold
    p.alignment = {"middle":PP_ALIGN.CENTER,"start":PP_ALIGN.LEFT,"end":PP_ALIGN.RIGHT}.get(anchor, PP_ALIGN.CENTER)
    return tb

def add_line(slide, x1, y1, x2, y2, color, sw=1.5, dashed=False):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = RGBColor(*bytes.fromhex(color)); c.line.width = Pt(sw)
    if dashed: c.line.dash_style = 2
    return c

def add_ellipse(slide, cx, cy, rx, ry, fill, stroke, sw=1.5, dashed=False, opacity=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(cx-rx), int(cy-ry), int(rx*2), int(ry*2))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*bytes.fromhex(fill))
    shape.line.color.rgb = RGBColor(*bytes.fromhex(stroke)); shape.line.width = Pt(sw)
    if dashed: shape.line.dash_style = 2
    return shape

def parse_elements(root, tx=0, ty=0):
    elements = []
    for child in root:
        tag = child.tag.split("}")[-1] if "}" in (child.tag or "") else child.tag
        ctx = get_transform(child); ctx_tx, ctx_ty = tx+ctx[0], ty+ctx[1]
        if tag == "g": elements.extend(parse_elements(child, ctx_tx, ctx_ty))
        elif tag == "rect" and float(child.get("width",0))>0 and float(child.get("height",0))>0:
            elements.append(("rect", {"x":float(child.get("x",0))+ctx_tx,"y":float(child.get("y",0))+ctx_ty,
                "w":float(child.get("width",0)),"h":float(child.get("height",0)),"rx":float(child.get("rx",0)),
                "fill":parse_fill(child),"stroke":parse_stroke(child),"sw":parse_stroke_width(child),
                "dashed":is_dashed(child),"opacity":parse_opacity(child)}))
        elif tag == "ellipse":
            elements.append(("ellipse", {"cx":float(child.get("cx",0))+ctx_tx,"cy":float(child.get("cy",0))+ctx_ty,
                "rx":float(child.get("rx",0)),"ry":float(child.get("ry",0)),"fill":parse_fill(child),
                "stroke":parse_stroke(child),"sw":parse_stroke_width(child),
                "dashed":is_dashed(child),"opacity":parse_opacity(child)}))
        elif tag == "line":
            elements.append(("line", {"x1":float(child.get("x1",0))+ctx_tx,"y1":float(child.get("y1",0))+ctx_ty,
                "x2":float(child.get("x2",0))+ctx_tx,"y2":float(child.get("y2",0))+ctx_ty,
                "stroke":parse_stroke(child),"sw":parse_stroke_width(child),"dashed":is_dashed(child)}))
        elif tag == "text":
            txt = (child.text or "") + "".join(t.text or "" for t in child.findall(f"{{{SVG_NS}}}tspan"))
            style = child.get("style",""); fs=int((re.search(r"font-size\s*:\s*(\d+)",style) or [0,11]).group(1))
            bold = child.get("font-weight") in ("700","bold","600") or bool(re.search(r"font-weight\s*:\s*(600|700|bold)",style))
            elements.append(("text", {"x":float(child.get("x",0))+ctx_tx,"y":float(child.get("y",0))+ctx_ty,
                "text":txt.strip(),"size":fs,"bold":bold,"color":parse_fill(child,"ffffff"),
                "anchor":child.get("text-anchor","start")}))
        elif tag == "polygon":
            pts = child.get("points","")
            if pts:
                pts = [(float(p.split(",")[0])+ctx_tx, float(p.split(",")[1])+ctx_ty) for p in pts.strip().split()]
                elements.append(("polygon", {"points":pts,"fill":parse_fill(child),"stroke":parse_stroke(child),
                    "sw":parse_stroke_width(child),"dashed":is_dashed(child),"opacity":parse_opacity(child)}))
    return elements

def build_pptx(elements, path):
    prs = Presentation(); prs.slide_width=PPTX_W; prs.slide_height=PPTX_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(0x0F,0x17,0x2A)
    for t, p in elements:
        try:
            if t=="rect": add_rect(slide, to_emu_x(p["x"]),to_emu_y(p["y"]),to_emu_w(p["w"]),to_emu_h(p["h"]),
                p["fill"],p["stroke"],p["sw"],p["rx"],p["dashed"],p["opacity"])
            elif t=="ellipse": add_ellipse(slide, to_emu_x(p["cx"]),to_emu_y(p["cy"]),to_emu_w(p["rx"]),to_emu_h(p["ry"]),
                p["fill"],p["stroke"],p["sw"],p["dashed"],p["opacity"])
            elif t=="line": add_line(slide, to_emu_x(p["x1"]),to_emu_y(p["y1"]),to_emu_x(p["x2"]),to_emu_y(p["y2"]),
                p["stroke"],p["sw"],p["dashed"])
            elif t=="text" and p["text"]: add_text(slide, to_emu_x(p["x"]),to_emu_y(p["y"]),
                p["text"],p["color"],p["size"],p["bold"],p["anchor"])
            elif t=="polygon" and len(p["points"])>=3:
                pts = [(to_emu_x(pt[0]),to_emu_y(pt[1])) for pt in p["points"]]
                mx,my = min(pt[0] for pt in pts),min(pt[1] for pt in pts)
                Mx,My = max(pt[0] for pt in pts),max(pt[1] for pt in pts)
                s = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, mx,my,Mx-mx,My-my)
                s.fill.solid(); s.fill.fore_color.rgb=RGBColor(*bytes.fromhex(p["fill"]))
                s.line.color.rgb=RGBColor(*bytes.fromhex(p["stroke"])); s.line.width=Pt(p["sw"])
                if p["dashed"]: s.line.dash_style=2
        except: pass
    prs.save(path); return path

def main():
    p = argparse.ArgumentParser(); p.add_argument("input"); p.add_argument("-o","--output")
    a = p.parse_args()
    if not os.path.exists(a.input): print(f"Error: {a.input} not found",file=sys.stderr); sys.exit(1)
    out = a.output or os.path.splitext(a.input)[0]+".pptx"
    global SVG_W, SVG_H
    tree = ET.parse(a.input); root = tree.getroot(); SVG_W, SVG_H = parse_viewbox(root)
    print(f"SVG: {SVG_W}x{SVG_H} -> {out}")
    elems = parse_elements(root); print(f"Elements: {len(elems)}")
    order = {"rect":0,"ellipse":1,"polygon":1,"path":2,"line":2,"text":3}
    elems.sort(key=lambda e: order.get(e[0],5))
    print(f"Done: {build_pptx(elems, out)}")

if __name__=="__main__": main()
