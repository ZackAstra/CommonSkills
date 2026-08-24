# -*- coding: utf-8 -*-
"""
SVG -> fully-editable PPTX converter (native PowerPoint shapes). Generic version.

Usage: python _svg2pptx_editable.py <input.svg> [output.pptx]

Element mapping:
  rect      -> rounded/plain rectangle (fill, stroke, dash, rx, gradient via url(#id))
  text      -> textbox (font size/weight/color, anchor, CSS class support)
  polygon   -> freeform (opacity pre-blended over background)
  path      -> freeform with curves sampled to polylines (icons)
  circle    -> oval
  line      -> straight connector
  image     -> picture (href resolved relative to the SVG file)
  g         -> translate/scale transform composition + attribute inheritance
"""
import re, math, os, sys
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE
from pptx.oxml.ns import qn

SVG_PATH = os.path.abspath(sys.argv[1])
OUT_PATH = os.path.abspath(sys.argv[2]) if len(sys.argv) > 2 else \
    SVG_PATH.replace(".svg", "-editable.pptx")
SVG_DIR = os.path.dirname(SVG_PATH)

svg_text = open(SVG_PATH, encoding="utf-8").read()
root = ET.fromstring(svg_text)
NS = "{http://www.w3.org/2000/svg}"
XLINK = "{http://www.w3.org/1999/xlink}href"

vb = root.get("viewBox")
if vb:
    SVG_W, SVG_H = [float(v) for v in vb.split()][2:4]
else:
    SVG_W, SVG_H = float(root.get("width")), float(root.get("height"))

SLIDE_W_IN = 13.333
SLIDE_W = int(SLIDE_W_IN * 914400)
SLIDE_H = int(SLIDE_W * SVG_H / SVG_W)
SX = SLIDE_W / SVG_W
PT_PER_PX = (SLIDE_W_IN * 72) / SVG_W

def emu(v): return Emu(int(round(v * SX)))

ident = lambda x, y: (x, y)

# ---------------- style block ----------------
def parse_style(text):
    css = {}
    m = re.search(r"<style>(.*?)</style>", text, re.S)
    if not m: return css
    for sel, body in re.findall(r"([^{}]+)\{([^}]+)\}", m[0]):
        props = {}
        for kv in body.split(";"):
            if ":" in kv:
                k, v = kv.split(":", 1)
                props[k.strip()] = v.strip()
        for s in sel.split(","):
            s = s.strip()
            if s.startswith("."):
                css[s[1:]] = props
    return css

css = parse_style(svg_text)

# ---------------- gradient defs ----------------
def parse_gradients(root_el):
    grads = {}
    for lg in root_el.iter(NS + "linearGradient"):
        gid = lg.get("id")
        stops = [(float(st.get("offset")), st.get("stop-color"))
                 for st in lg.iter(NS + "stop")]
        x1, y1 = float(lg.get("x1", 0)), float(lg.get("y1", 0))
        x2, y2 = float(lg.get("x2", 0)), float(lg.get("y2", 1))
        angle = math.degrees(math.atan2(y2 - y1, x2 - x1))  # 90 = top->bottom
        grads[gid] = (stops, angle)
    return grads

GRADS = parse_gradients(root)

def hex2rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def blend(fg, bg, alpha):
    fg, bg = fg.lstrip("#"), bg.lstrip("#")
    c = [round(int(fg[i:i+2],16)*alpha + int(bg[i:i+2],16)*(1-alpha)) for i in (0,2,4)]
    return RGBColor(*c)

# ---------------- transforms ----------------
def parse_transform(t):
    ops = re.findall(r"(translate|scale)\(([^)]*)\)", t or "")
    def f(x, y):
        # SVG transform list: p' = op1(op2(...(p))) — later ops apply to the point first
        for name, args in reversed(ops):
            a = [float(v) for v in re.split(r"[,\s]+", args.strip()) if v]
            if name == "translate":
                x += a[0]; y += a[1] if len(a) > 1 else 0.0
            else:
                x *= a[0]; y *= a[1] if len(a) > 1 else a[0]
        return x, y
    return f

def tf_scale(tf):
    x0, y0 = tf(0, 0); x1, y1 = tf(1, 0)
    return math.hypot(x1 - x0, y1 - y0)

# ---------------- path parser (M L H V C Q A Z, abs+rel) ----------------
TOKEN = re.compile(r"([MmLlHhVvCcQqAaZz])|(-?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?)")

def arc_points(x1, y1, rx, ry, phi, laf, sf, x2, y2, n=14):
    if rx == 0 or ry == 0:
        return [(x2, y2)]
    phi = math.radians(phi)
    cp, sp = math.cos(phi), math.sin(phi)
    dx, dy = (x1-x2)/2.0, (y1-y2)/2.0
    x1p = cp*dx + sp*dy
    y1p = -sp*dx + cp*dy
    rx, ry = abs(rx), abs(ry)
    lam = x1p**2/rx**2 + y1p**2/ry**2
    if lam > 1:
        s = math.sqrt(lam); rx *= s; ry *= s
    num = rx**2*ry**2 - rx**2*y1p**2 - ry**2*x1p**2
    den = rx**2*y1p**2 + ry**2*x1p**2
    co = math.sqrt(max(0.0, num/den)) if den else 0.0
    if laf == sf: co = -co
    cxp = co*rx*y1p/ry
    cyp = -co*ry*x1p/rx
    cx = cp*cxp - sp*cyp + (x1+x2)/2.0
    cy = sp*cxp + cp*cyp + (y1+y2)/2.0
    def ang(ux, uy, vx, vy):
        d = math.hypot(ux, uy)*math.hypot(vx, vy)
        c = max(-1, min(1, (ux*vx+uy*vy)/d))
        a = math.acos(c)
        if ux*vy - uy*vx < 0: a = -a
        return a
    th1 = ang(1, 0, (x1p-cxp)/rx, (y1p-cyp)/ry)
    dth = ang((x1p-cxp)/rx, (y1p-cyp)/ry, (-x1p-cxp)/rx, (-y1p-cyp)/ry)
    if not sf and dth > 0: dth -= 2*math.pi
    if sf and dth < 0: dth += 2*math.pi
    pts = []
    for i in range(1, n+1):
        t = th1 + dth*i/n
        xp, yp = rx*math.cos(t), ry*math.sin(t)
        pts.append((cp*xp - sp*yp + cx, sp*xp + cp*yp + cy))
    return pts

def cubic_points(p0, p1, p2, p3, n=12):
    pts = []
    for i in range(1, n+1):
        t = i/n; mt = 1-t
        x = mt**3*p0[0] + 3*mt*mt*t*p1[0] + 3*mt*t*t*p2[0] + t**3*p3[0]
        y = mt**3*p0[1] + 3*mt*mt*t*p1[1] + 3*mt*t*t*p2[1] + t**3*p3[1]
        pts.append((x, y))
    return pts

def parse_path(d):
    toks = TOKEN.findall(d)
    items = [cmd if cmd else float(num) for cmd, num in toks]
    subs, pts = [], []
    i, x, y, sx0, sy0 = 0, 0.0, 0.0, 0.0, 0.0
    cmd = None
    def is_cmd(v): return isinstance(v, str)
    while i < len(items):
        if is_cmd(items[i]):
            cmd = items[i]; i += 1
            if cmd in "Zz":
                subs.append((pts, True)); pts = []
                x, y = sx0, sy0
                cmd = None
            continue
        rel = cmd.islower(); c = cmd.upper()
        try:
            if c == "M":
                x1, y1 = items[i], items[i+1]; i += 2
                if rel: x1 += x; y1 += y
                if pts: subs.append((pts, False))
                pts = [(x1, y1)]; x, y = x1, y1; sx0, sy0 = x1, y1
                cmd = "l" if rel else "L"
            elif c == "L":
                x1, y1 = items[i], items[i+1]; i += 2
                if rel: x1 += x; y1 += y
                pts.append((x1, y1)); x, y = x1, y1
            elif c == "H":
                x1 = items[i]; i += 1
                if rel: x1 += x
                pts.append((x1, y)); x = x1
            elif c == "V":
                y1 = items[i]; i += 1
                if rel: y1 += y
                pts.append((x, y1)); y = y1
            elif c == "C":
                x1, y1, x2, y2, x3, y3 = items[i:i+6]; i += 6
                if rel:
                    x1 += x; y1 += y; x2 += x; y2 += y; x3 += x; y3 += y
                pts += cubic_points((x, y), (x1, y1), (x2, y2), (x3, y3))
                x, y = x3, y3
            elif c == "Q":
                x1, y1, x2, y2 = items[i:i+4]; i += 4
                if rel:
                    x1 += x; y1 += y; x2 += x; y2 += y
                c1 = (x + 2/3*(x1-x), y + 2/3*(y1-y))
                c2 = (x2 + 2/3*(x1-x2), y2 + 2/3*(y1-y2))
                pts += cubic_points((x, y), c1, c2, (x2, y2))
                x, y = x2, y2
            elif c == "A":
                rx, ry, phi, laf, sf, x2, y2 = items[i:i+7]; i += 7
                if rel: x2 += x; y2 += y
                pts += arc_points(x, y, rx, ry, phi, int(laf), int(sf), x2, y2)
                x, y = x2, y2
            else:
                i += 1
        except (IndexError, TypeError):
            break
    if pts: subs.append((pts, False))
    return subs

# ---------------- pptx setup ----------------
prs = Presentation()
prs.slide_width = Emu(SLIDE_W)
prs.slide_height = Emu(SLIDE_H)
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes

def set_fill(shape, fill, opacity=None, bg="#ffffff"):
    if not fill or fill == "none":
        shape.fill.background()
        return
    m = re.match(r"url\(#(.+)\)", fill)
    if m and m.group(1) in GRADS:
        stops, angle = GRADS[m.group(1)]
        shape.fill.gradient()
        gs = shape.fill.gradient_stops
        for i, (pos, color) in enumerate(stops[:len(gs)]):
            gs[i].position = pos
            gs[i].color.rgb = hex2rgb(color)
        try:
            shape.fill.gradient_angle = angle
        except Exception:
            pass
        return
    shape.fill.solid()
    if opacity is not None and float(opacity) < 1:
        shape.fill.fore_color.rgb = blend(fill, bg, float(opacity))
    else:
        shape.fill.fore_color.rgb = hex2rgb(fill)

def set_line(shape, stroke, sw=None, dash=None):
    if not stroke or stroke == "none":
        shape.line.fill.background()
        return
    shape.line.color.rgb = hex2rgb(stroke)
    if sw: shape.line.width = Emu(int(float(sw) * SX))
    if dash:
        shape.line.dash_style = MSO_LINE.DASH

def fnum(v):
    return float(re.sub(r"[^\d.]", "", str(v)))

def add_rect(el, tf=ident, inh=None):
    inh = inh or {}
    x = float(el.get("x", 0)); y = float(el.get("y", 0))
    w = float(el.get("width", 0)); h = float(el.get("height", 0))
    rx = float(el.get("rx", 0) or 0)
    if tf is not ident:
        s = tf_scale(tf)
        x, y = tf(x, y); w *= s; h *= s; rx *= s
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
    sp = shapes.add_shape(st, emu(x), emu(y), emu(w), emu(h))
    if rx > 0:
        try: sp.adjustments[0] = min(0.5, rx / min(w, h))
        except Exception: pass
    rfill = el.get("fill")
    if rfill is None: rfill = inh.get("fill")
    set_fill(sp, rfill)
    set_line(sp, el.get("stroke") or inh.get("stroke"),
             el.get("stroke-width") or inh.get("stroke-width"),
             el.get("stroke-dasharray"))
    sp.shadow.inherit = False
    return sp

def add_circle(el, tf=ident, inh=None):
    inh = inh or {}
    cx = float(el.get("cx", 0)); cy = float(el.get("cy", 0)); r = float(el.get("r", 0))
    if tf is not ident:
        s = tf_scale(tf)
        cx, cy = tf(cx, cy); r *= s
    sp = shapes.add_shape(MSO_SHAPE.OVAL, emu(cx - r), emu(cy - r), emu(2*r), emu(2*r))
    cfill = el.get("fill")
    if cfill is None: cfill = inh.get("fill")
    set_fill(sp, cfill)
    set_line(sp, el.get("stroke") or inh.get("stroke"),
             el.get("stroke-width") or inh.get("stroke-width"), None)
    sp.shadow.inherit = False
    return sp

def add_line(el, tf=ident, inh=None):
    inh = inh or {}
    x1 = float(el.get("x1", 0)); y1 = float(el.get("y1", 0))
    x2 = float(el.get("x2", 0)); y2 = float(el.get("y2", 0))
    if tf is not ident:
        x1, y1 = tf(x1, y1); x2, y2 = tf(x2, y2)
    conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    set_line(conn, el.get("stroke") or inh.get("stroke"),
             el.get("stroke-width") or inh.get("stroke-width"), None)
    conn.shadow.inherit = False
    return conn

def add_image(el, tf=ident):
    x = float(el.get("x", 0)); y = float(el.get("y", 0))
    w = float(el.get("width", 0)); h = float(el.get("height", 0))
    if tf is not ident:
        s = tf_scale(tf)
        x, y = tf(x, y); w *= s; h *= s
    href = el.get("href") or el.get(XLINK)
    path = href if os.path.isabs(href) else os.path.join(SVG_DIR, href)
    pic = shapes.add_picture(path, emu(x), emu(y), emu(w), emu(h))
    pic.shadow.inherit = False
    return pic

def est_text_width(text, fs):
    w = 0.0
    for ch in text:
        w += fs * (1.0 if ord(ch) > 0x2E7F else 0.56)
    return w

def add_text(el, tf=ident, inh=None):
    content = "".join(el.itertext())
    if not content.strip(): return
    inh = inh or {}
    cls = el.get("class", "")
    props = css.get(cls, {}) if cls else {}
    fill = el.get("fill") or props.get("fill") or inh.get("fill", "#000000")
    fs = fnum(el.get("font-size") or props.get("font-size") or inh.get("font-size", 16))
    bold = (el.get("font-weight") or props.get("font-weight", "")) in ("700", "bold")
    anchor = el.get("text-anchor", "start")
    ls = el.get("letter-spacing")
    x = float(el.get("x", 0)); y = float(el.get("y", 0))
    tw = est_text_width(content, fs) * (1 + (float(ls)/fs if ls else 0))
    th = fs * 1.35
    if anchor == "middle":
        bx = x - tw/2
    elif anchor == "end":
        bx = x - tw
    else:
        bx = x
    by = y - fs * 1.02
    tb = shapes.add_textbox(emu(bx - fs*0.35), emu(by), emu(tw + fs*0.7), emu(th))
    tfm = tb.text_frame
    tfm.word_wrap = False
    tfm.margin_left = tfm.margin_right = tfm.margin_top = tfm.margin_bottom = 0
    tfm.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tfm.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if anchor == "middle" else (PP_ALIGN.RIGHT if anchor == "end" else PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = content
    f = r.font
    f.size = Pt(fs * PT_PER_PX)
    f.bold = bold
    f.name = "微软雅黑"
    f.color.rgb = hex2rgb(fill)
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {"typeface": "微软雅黑"})
    rPr.append(ea)

def add_freeform(points, fill, stroke, sw, opacity, bg, tf, close=True):
    pts = [tf(px, py) for px, py in points]
    fb = shapes.build_freeform(Emu(int(pts[0][0]*SX)), Emu(int(pts[0][1]*SX)), scale=1)
    fb.add_line_segments([(Emu(int(px*SX)), Emu(int(py*SX))) for px, py in pts[1:]], close=close)
    sp = fb.convert_to_shape()
    set_fill(sp, fill, opacity, bg)
    set_line(sp, stroke, sw, None)
    sp.shadow.inherit = False
    return sp

# ---------------- walk svg ----------------
def tag(el): return el.tag.replace(NS, "")

def walk(el, tf, bg, inh=None):
    inh = dict(inh or {})
    for k in ("fill", "font-size", "font-weight", "stroke", "stroke-width"):
        if el.get(k): inh[k] = el.get(k)
    for child in el:
        t = tag(child)
        if t in ("style", "defs"):
            continue
        if t == "g":
            t2 = child.get("transform")
            if t2:
                n = parse_transform(t2)
                ntf = n if tf is ident else (lambda x, y, b=tf, n=n: b(*n(x, y)))
            else:
                ntf = tf
            walk(child, ntf, bg, inh)
        elif t == "rect":
            add_rect(child, tf, inh)
        elif t == "circle":
            add_circle(child, tf, inh)
        elif t == "line":
            add_line(child, tf, inh)
        elif t == "image":
            add_image(child, tf)
        elif t == "text":
            add_text(child, tf, inh)
        elif t == "polygon":
            pts = []
            for pair in child.get("points", "").split():
                px, py = pair.split(",")
                pts.append((float(px), float(py)))
            add_freeform(pts, child.get("fill"), child.get("stroke"),
                         child.get("stroke-width"), child.get("opacity"), bg, ident)
        elif t == "path":
            pfill = child.get("fill")
            if pfill is None: pfill = inh.get("fill")
            pstroke = child.get("stroke") or inh.get("stroke")
            psw = child.get("stroke-width") or inh.get("stroke-width")
            for pts, closed in parse_path(child.get("d", "")):
                if len(pts) >= 2:
                    add_freeform(pts, pfill, pstroke, psw, None, bg, tf, close=closed)

# background color = fill of the first full-canvas rect, fallback white
BG = "#ffffff"
for child in root:
    if tag(child) == "rect" and not child.get("x"):
        if child.get("fill") and not child.get("fill").startswith("url"):
            BG = child.get("fill")
        break
walk(root, ident, BG, {})

prs.save(OUT_PATH)
print("saved:", OUT_PATH, os.path.getsize(OUT_PATH), "bytes,", len(shapes._spTree), "elements")
